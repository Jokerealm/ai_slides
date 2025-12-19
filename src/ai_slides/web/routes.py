"""
Web interface routes for ai_slides
"""

from fastapi import APIRouter, Request, Form, UploadFile, File, HTTPException, Depends
from fastapi.responses import HTMLResponse, StreamingResponse, FileResponse, JSONResponse
from fastapi.templating import Jinja2Templates
from pydantic import BaseModel
import json
import uuid
import asyncio
import time
import os
import zipfile
import tempfile
import shutil
from pathlib import Path
from datetime import datetime
import urllib.parse
import subprocess
import logging
import time
from typing import Optional, Dict, Any, List

from ..api.models import PPTGenerationRequest, PPTProject, TodoBoard, FileOutlineGenerationRequest
from ..services.enhanced_ppt_service import EnhancedPPTService
from ..core.config import ai_config
from ..ai import get_ai_provider, get_role_provider, AIMessage, MessageRole
from ..database.database import get_db
from sqlalchemy.orm import Session
from ..utils.thread_pool import run_blocking_io, to_thread
import re
from bs4 import BeautifulSoup

# Configure logger for this module
logger = logging.getLogger(__name__)

router = APIRouter()
templates = Jinja2Templates(directory="src/ai_slides/web/templates")

# Add custom filters
def timestamp_to_datetime(timestamp):
    """Convert timestamp to readable datetime string"""
    try:
        if isinstance(timestamp, (int, float)):
            return datetime.fromtimestamp(timestamp).strftime("%Y-%m-%d %H:%M:%S")
        return str(timestamp)
    except (ValueError, OSError):
        return "æ— æ•ˆæ—¶é—´"

def strftime_filter(timestamp, format_string="%Y-%m-%d %H:%M"):
    """Jinja2 strftime filter"""
    try:
        if isinstance(timestamp, (int, float)):
            dt = datetime.fromtimestamp(timestamp)
            return dt.strftime(format_string)
        return str(timestamp)
    except (ValueError, OSError):
        return "æ— æ•ˆæ—¶é—´"

# Register custom filters
templates.env.filters["timestamp_to_datetime"] = timestamp_to_datetime
templates.env.filters["strftime"] = strftime_filter

# Import shared service instances to ensure data consistency
from ..services.service_instances import ppt_service
from ..services.pyppeteer_pdf_converter import get_pdf_converter
from ..services.pdf_to_pptx_converter import get_pdf_to_pptx_converter

# Favicon route - return 203 Non-Authoritative Information
@router.get("/favicon.ico")
async def favicon():
    """Return 203 for favicon requests"""
    return JSONResponse(content={"message": "No favicon"}, status_code=203)

# AIç¼–è¾‘è¯·æ±‚æ•°æ®æ¨¡å‹
class AISlideEditRequest(BaseModel):
    slideIndex: int
    slideTitle: str
    slideContent: str
    userRequest: str
    projectInfo: Dict[str, Any]
    slideOutline: Optional[Dict[str, Any]] = None
    chatHistory: Optional[List[Dict[str, str]]] = None
    images: Optional[List[Dict[str, str]]] = None  # æ–°å¢:å›¾ç‰‡ä¿¡æ¯åˆ—è¡¨
    visionEnabled: Optional[bool] = False  # æ–°å¢:è§†è§‰æ¨¡å¼å¯ç”¨çŠ¶æ€
    slideScreenshot: Optional[str] = None  # æ–°å¢:å¹»ç¯ç‰‡æˆªå›¾æ•°æ®(base64æ ¼å¼)

# AIè¦ç‚¹å¢å¼ºè¯·æ±‚æ•°æ®æ¨¡å‹
class AIBulletPointEnhanceRequest(BaseModel):
    slideIndex: int
    slideTitle: str
    slideContent: str
    userRequest: str
    projectInfo: Dict[str, Any]
    slideOutline: Optional[Dict[str, Any]] = None
    contextInfo: Optional[Dict[str, Any]] = None  # åŒ…å«åŸå§‹è¦ç‚¹,å…¶ä»–è¦ç‚¹ç­‰ä¸Šä¸‹æ–‡ä¿¡æ¯

# å›¾åƒé‡æ–°ç”Ÿæˆè¯·æ±‚æ•°æ®æ¨¡å‹
class AIImageRegenerateRequest(BaseModel):
    slide_index: int
    image_info: Dict[str, Any]
    slide_content: Dict[str, Any]
    project_topic: str
    project_scenario: str
    regeneration_reason: Optional[str] = None

# ä¸€é”®é…å›¾è¯·æ±‚æ•°æ®æ¨¡å‹
class AIAutoImageGenerateRequest(BaseModel):
    slide_index: int
    slide_content: Dict[str, Any]
    project_topic: str
    project_scenario: str


class AutoLayoutRepairRequest(BaseModel):
    html_content: str
    slide_data: Dict[str, Any]


class SpeechScriptGenerationRequest(BaseModel):
    generation_type: str  # "single", "multi", "full"
    slide_indices: Optional[List[int]] = None  # For single and multi generation
    customization: Dict[str, Any] = {}  # Customization options

class SpeechScriptExportRequest(BaseModel):
    export_format: str  # "docx", "markdown"
    scripts_data: List[Dict[str, Any]]
    include_metadata: bool = True

# å›¾ç‰‡å¯¼å‡ºPPTXè¯·æ±‚æ•°æ®æ¨¡å‹
class ImagePPTXExportRequest(BaseModel):
    slides: Optional[List[Dict[str, Any]]] = None  # åŒ…å«index, html_content, title
    images: Optional[List[Dict[str, Any]]] = None  # åŒ…å«index, data(base64), width, height (å‘åå…¼å®¹)

# Helper function to extract slides from HTML content
async def _extract_slides_from_html(slides_html: str, existing_slides_data: list) -> list:
    """
    Extract individual slides from combined HTML content and update slides_data
    """
    try:
        # Parse HTML content
        soup = BeautifulSoup(slides_html, 'html.parser')

        # Find all slide containers - look for common slide patterns
        slide_containers = []

        # Try different patterns to find slides
        patterns = [
            {'class': re.compile(r'slide')},
            {'class': re.compile(r'page')},
            {'style': re.compile(r'width:\s*1280px.*height:\s*720px', re.IGNORECASE)},
            {'style': re.compile(r'aspect-ratio:\s*16\s*/\s*9', re.IGNORECASE)}
        ]

        for pattern in patterns:
            containers = soup.find_all('div', pattern)
            if containers:
                slide_containers = containers
                break

        # If no specific slide containers found, try to split by common separators
        if not slide_containers:
            # Look for sections or divs that might represent slides
            all_divs = soup.find_all('div')
            # Filter divs that might be slides (have substantial content)
            slide_containers = [div for div in all_divs
                             if div.get_text(strip=True) and len(div.get_text(strip=True)) > 50]

        updated_slides_data = []

        # If we found slide containers, extract them
        if slide_containers:
            for i, container in enumerate(slide_containers):
                # Try to extract title from the slide
                title = f"ç¬¬{i+1}é¡µ"
                title_elements = container.find_all(['h1', 'h2', 'h3', 'h4', 'h5', 'h6'])
                if title_elements:
                    title = title_elements[0].get_text(strip=True) or title

                # Get the HTML content of this slide
                slide_html = str(container)

                # Create slide data
                slide_data = {
                    "page_number": i + 1,
                    "title": title,
                    "html_content": slide_html,
                    "is_user_edited": True  # Mark as user edited since it came from editor
                }

                # If we have existing slide data, preserve some fields
                if i < len(existing_slides_data):
                    existing_slide = existing_slides_data[i]
                    # Preserve any additional fields from existing data
                    for key, value in existing_slide.items():
                        if key not in slide_data:
                            slide_data[key] = value

                updated_slides_data.append(slide_data)

        # If we couldn't extract individual slides, treat the entire content as slides
        if not updated_slides_data and existing_slides_data:
            # Fall back to using existing slides structure but mark as edited
            for i, existing_slide in enumerate(existing_slides_data):
                slide_data = existing_slide.copy()
                slide_data["is_user_edited"] = True
                updated_slides_data.append(slide_data)

        # If we still have no slides but have HTML content, create a single slide
        if not updated_slides_data and slides_html.strip():
            slide_data = {
                "page_number": 1,
                "title": "ç¼–è¾‘åçš„PPT",
                "html_content": slides_html,
                "is_user_edited": True
            }
            updated_slides_data.append(slide_data)

        logger.info(f"Extracted {len(updated_slides_data)} slides from HTML content")
        return updated_slides_data

    except Exception as e:
        logger.error(f"Error extracting slides from HTML: {e}")
        # Fall back to marking existing slides as edited
        if existing_slides_data:
            updated_slides_data = []
            for slide in existing_slides_data:
                slide_copy = slide.copy()
                slide_copy["is_user_edited"] = True
                updated_slides_data.append(slide_copy)
            return updated_slides_data
        else:
            return []


def _clean_html_for_pdf(original_html: str, slide_number: int, total_slides: int) -> str:
    """Clean complete HTML document for PDF generation by removing navigation elements"""
    import re
    
    # Remove navigation elements that might interfere with PDF generation
    cleaned_html = original_html
    
    # Remove navigation divs and buttons
    cleaned_html = re.sub(r'<div[^>]*class="[^"]*navigation[^"]*"[^>]*>.*?</div>', '', cleaned_html, flags=re.DOTALL | re.IGNORECASE)
    cleaned_html = re.sub(r'<button[^>]*class="[^"]*nav[^"]*"[^>]*>.*?</button>', '', cleaned_html, flags=re.DOTALL | re.IGNORECASE)
    cleaned_html = re.sub(r'<a[^>]*class="[^"]*nav[^"]*"[^>]*>.*?</a>', '', cleaned_html, flags=re.DOTALL | re.IGNORECASE)
    
    # Remove fullscreen buttons
    cleaned_html = re.sub(r'<button[^>]*fullscreen[^>]*>.*?</button>', '', cleaned_html, flags=re.DOTALL | re.IGNORECASE)
    
    # Add PDF-specific styles
    pdf_styles = """<style>
/* PDF optimization styles */
* {
    -webkit-print-color-adjust: exact !important;
    print-color-adjust: exact !important;
}
html, body {
    width: 100% !important;
    height: 100vh !important;
    margin: 0 !important;
    padding: 0 !important;
    overflow: hidden !important;
}
/* Hide any remaining navigation elements */
.navigation, .nav-btn, .fullscreen-btn, .slide-navigation {
    display: none !important;
}
</style>"""
    
    # Insert PDF styles before closing head tag
    head_pattern = r'</head>'
    cleaned_html = re.sub(head_pattern, pdf_styles + '\n</head>', cleaned_html, flags=re.IGNORECASE)
    
    return cleaned_html


async def _generate_pdf_slide_html(slide, slide_number: int, total_slides: int, topic: str) -> str:
    """Generate PDF-optimized HTML for individual slide without navigation elements"""
    slide_html = slide.get('html_content', '')
    slide_title = slide.get('title', f'ç¬¬{slide_number}é¡µ')
    
    # Check if it's already a complete HTML document
    import re
    if slide_html.strip().lower().startswith('<!doctype') or slide_html.strip().lower().startswith('<html'):
        # It's a complete HTML document, clean it for PDF
        return _clean_html_for_pdf(slide_html, slide_number, total_slides)
    else:
        # It's just content, wrap it in a PDF-optimized structure
        slide_content = slide_html
        
        return f"""<!DOCTYPE html>
<html lang="zh-CN">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>{topic} - {slide_title}</title>
    <style>
        * {{
            margin: 0;
            padding: 0;
            box-sizing: border-box;
        }}
        html, body {{
            width: 100%;
            height: 100vh;
            margin: 0;
            padding: 0;
            font-family: 'Microsoft YaHei', 'PingFang SC', sans-serif;
            overflow: hidden;
        }}
        .slide-container {{
            width: 100vw;
            height: 100vh;
            display: flex;
            align-items: center;
            justify-content: center;
            position: relative;
        }}
        .slide-content {{
            width: 100%;
            height: 100%;
            display: flex;
            align-items: center;
            justify-content: center;
            position: relative;
        }}
        /* Ensure all backgrounds and colors are preserved for PDF */
        * {{
            -webkit-print-color-adjust: exact !important;
            print-color-adjust: exact !important;
        }}
    </style>
</head>
<body>
    <div class="slide-container">
        <div class="slide-content">
            {slide_content}
        </div>
    </div>
</body>
</html>"""


async def _generate_pdf_with_pyppeteer(project, output_path: str, individual: bool = False) -> bool:
    """Generate PDF using Pyppeteer (Python)"""
    try:
        pdf_converter = get_pdf_converter()
        
        with tempfile.TemporaryDirectory() as temp_dir:
            temp_path = Path(temp_dir)
            
            # Always generate individual HTML files for each slide for better page separation
            # This ensures each slide becomes a separate PDF page
            html_files = []
            for i, slide in enumerate(project.slides_data):
                # Use a specialized PDF-optimized HTML generator without navigation
                slide_html = await _generate_pdf_slide_html(slide, i+1, len(project.slides_data), project.topic)
                html_file = temp_path / f"slide_{i+1}.html"
                
                # Write HTML file in thread pool to avoid blocking
                def write_html_file(content, path):
                    with open(path, 'w', encoding='utf-8') as f:
                        f.write(content)
                
                await run_blocking_io(write_html_file, slide_html, str(html_file))
                html_files.append(str(html_file))
            
            # Use Pyppeteer to convert multiple files and merge them
            pdf_dir = temp_path / "pdfs"
            await run_blocking_io(pdf_dir.mkdir)
            
            logging.info(f"Starting PDF generation for {len(html_files)} files")
            
            # Convert HTML files to PDFs and merge them
            pdf_files = await pdf_converter.convert_multiple_html_to_pdf(
                html_files, 
                str(pdf_dir), 
                output_path
            )
            
            if pdf_files and os.path.exists(output_path):
                logging.info("Pyppeteer PDF generation successful")
                return True
            else:
                logging.error("Pyppeteer PDF generation failed: No output file created")
                return False
                
    except Exception as e:
        logging.error(f"Pyppeteer PDF generation failed: {e}")
        return False


@router.get("/home", response_class=HTMLResponse)
async def web_home(
    request: Request
):
    """Main web interface home page - redirect to dashboard for existing users"""
    # Check if user has projects, if so redirect to dashboard
    try:
        projects_response = await ppt_service.project_manager.list_projects(page=1, page_size=1)
        if projects_response.total > 0:
            # User has projects, redirect to dashboard
            from fastapi.responses import RedirectResponse
            return RedirectResponse(url="/dashboard", status_code=302)
    except:
        pass  # If error, show index page

    # New user or error, show index page
    return templates.TemplateResponse("index.html", {
        "request": request,
        "ai_provider": ai_config.default_ai_provider,
        "available_providers": ai_config.get_available_providers()
    })

@router.get("/ai-config", response_class=HTMLResponse)
async def web_ai_config(
    request: Request
):
    """AI configuration page"""
    from ..services.config_service import get_config_service

    config_service = get_config_service()
    current_config = config_service.get_all_config()

    return templates.TemplateResponse("ai_config.html", {
        "request": request,
        "current_provider": ai_config.default_ai_provider,
        "available_providers": ai_config.get_available_providers(),
        "provider_status": {
            provider: ai_config.is_provider_available(provider)
            for provider in ai_config.get_available_providers()
        },
        "current_config": current_config,
        "user": user.to_dict()
    })


@router.get("/image-generation-test", response_class=HTMLResponse)
async def web_image_generation_test(
    request: Request
):
    """AIå›¾ç‰‡ç”Ÿæˆæµ‹è¯•é¡µé¢"""
    return templates.TemplateResponse("image_generation_test.html", {
        "request": request,
        "user": user.to_dict()
    })


@router.post("/api/ai/providers/openai/models")
async def get_openai_models(
    request: Request
):
    """Proxy endpoint to get OpenAI models list, avoiding CORS issues - uses frontend provided config"""
    try:
        import aiohttp
        import json
        
        # Get configuration from frontend request
        data = await request.json()
        base_url = data.get('base_url', 'https://api.openai.com/v1')
        api_key = data.get('api_key', '')
        
        logger.info(f"Frontend requested models from: {base_url}")
        
        if not api_key:
            return {"success": False, "error": "API Key is required"}
        
        # Ensure base URL ends with /v1
        if not base_url.endswith('/v1'):
            base_url = base_url.rstrip('/') + '/v1'
        
        models_url = f"{base_url}/models"
        logger.info(f"Fetching models from: {models_url}")
        
        # Make request to OpenAI API using frontend provided credentials
        async with aiohttp.ClientSession() as session:
            headers = {
                'Authorization': f'Bearer {api_key}',
                'Content-Type': 'application/json'
            }
            
            async with session.get(models_url, headers=headers, timeout=30) as response:
                if response.status == 200:
                    data = await response.json()
                    
                    # Filter and sort models
                    models = []
                    if 'data' in data and isinstance(data['data'], list):
                        for model in data['data']:
                            if model.get('id'):
                                models.append({
                                    'id': model['id'],
                                    'created': model.get('created', 0),
                                    'owned_by': model.get('owned_by', 'unknown')
                                })
                        
                        # Sort models with GPT-4 first, then GPT-3.5, then others
                        def get_priority(model_id):
                            if 'gpt-4' in model_id:
                                return 0
                            elif 'gpt-3.5' in model_id:
                                return 1
                            else:
                                return 2
                        
                        models.sort(key=lambda x: (get_priority(x['id']), x['id']))
                    logger.info(f"Successfully fetched {len(models)} models from {base_url}")
                    return {"success": True, "models": models}
                else:
                    error_text = await response.text()
                    logger.error(f"Failed to fetch models from {base_url}: {response.status} - {error_text}")
                    return {"success": False, "error": f"API returned status {response.status}: {error_text}"}
                    
    except Exception as e:
        logger.error(f"Error fetching OpenAI models from frontend config: {e}")
        return {"success": False, "error": str(e)}

@router.post("/api/ai/providers/openai/test")
async def test_openai_provider_proxy(
    request: Request
):
    """Proxy endpoint to test OpenAI provider, avoiding CORS issues - uses frontend provided config"""
    try:
        import aiohttp
        
        # Get configuration from frontend request
        data = await request.json()
        base_url = data.get('base_url', 'https://api.openai.com/v1')
        api_key = data.get('api_key', '')
        model = data.get('model', 'gpt-4o')
        
        logger.info(f"Frontend requested test with: base_url={base_url}, model={model}")
        
        if not api_key:
            return {"success": False, "error": "API Key is required"}
        
        # Ensure base URL ends with /v1
        if not base_url.endswith('/v1'):
            base_url = base_url.rstrip('/') + '/v1'
        
        chat_url = f"{base_url}/chat/completions"
        logger.info(f"Testing OpenAI provider at: {chat_url}")
        
        # Make test request to OpenAI API using frontend provided credentials
        async with aiohttp.ClientSession() as session:
            headers = {
                'Authorization': f'Bearer {api_key}',
                'Content-Type': 'application/json'
            }
            
            payload = {
                "model": model,
                "messages": [
                    {
                        "role": "user",
                        "content": "Say 'Hello, I am working!' in exactly 5 words."
                    }
                ],
                "temperature": 0
            }
            
            async with session.post(chat_url, headers=headers, json=payload, timeout=30) as response:
                if response.status == 200:
                    data = await response.json()
                    
                    logger.info(f"Test successful for {base_url} with model {model}")
                    
                    # Return with consistent format that frontend expects
                    return {
                        "success": True,
                        "status": "success",  # Add status field for compatibility
                        "provider": "openai",
                        "model": model,
                        "response_preview": data['choices'][0]['message']['content'],
                        "usage": data.get('usage', {
                            "prompt_tokens": 0,
                            "completion_tokens": 0,
                            "total_tokens": 0
                        })
                    }
                else:
                    error_text = await response.text()
                    try:
                        error_data = json.loads(error_text)
                        error_message = error_data.get('error', {}).get('message', f"API returned status {response.status}")
                    except:
                        error_message = f"API returned status {response.status}: {error_text}"
                    
                    logger.error(f"Test failed for {base_url}: {error_message}")
                    
                    return {
                        "success": False,
                        "status": "error",  # Add status field for compatibility
                        "error": error_message
                    }
                    
    except Exception as e:
        logger.error(f"Error testing OpenAI provider with frontend config: {e}")
        return {
            "success": False,
            "status": "error",  # Add status field for compatibility
            "error": str(e)
        }

@router.get("/scenarios", response_class=HTMLResponse)
async def web_scenarios(
    request: Request
):
    """Scenarios selection page"""
    scenarios = [
        {"id": "general", "name": "é€šç”¨", "description": "é€‚ç”¨äºå„ç§é€šç”¨åœºæ™¯çš„PPTæ¨¡æ¿", "icon": "ğŸ“‹"},
        {"id": "tourism", "name": "æ—…æ¸¸è§‚å…‰", "description": "æ—…æ¸¸çº¿è·¯,æ™¯ç‚¹ä»‹ç»ç­‰æ—…æ¸¸ç›¸å…³PPT", "icon": "ğŸŒ"},
        {"id": "education", "name": "å„¿ç«¥ç§‘æ™®", "description": "é€‚åˆå„¿ç«¥çš„ç§‘æ™®æ•™è‚²PPT", "icon": "ğŸ“"},
        {"id": "analysis", "name": "æ·±å…¥åˆ†æ", "description": "æ•°æ®åˆ†æ,ç ”ç©¶æŠ¥å‘Šç­‰æ·±åº¦åˆ†æPPT", "icon": "ğŸ“Š"},
        {"id": "history", "name": "å†å²æ–‡åŒ–", "description": "å†å²äº‹ä»¶,æ–‡åŒ–ä»‹ç»ç­‰äººæ–‡ç±»PPT", "icon": "ğŸ›ï¸"},
        {"id": "technology", "name": "ç§‘æŠ€æŠ€æœ¯", "description": "æŠ€æœ¯ä»‹ç»,äº§å“å‘å¸ƒç­‰ç§‘æŠ€ç±»PPT", "icon": "ğŸ’»"},
        {"id": "business", "name": "æ–¹æ¡ˆæ±‡æŠ¥", "description": "å•†ä¸šè®¡åˆ’,é¡¹ç›®æ±‡æŠ¥ç­‰å•†åŠ¡PPT", "icon": "ğŸ’¼"}
    ]
    return templates.TemplateResponse("scenarios.html", {"request": request, "scenarios": scenarios})

# Legacy route removed - now using /projects/create for new project workflow

# Legacy task status route removed - now using project detail pages

# Legacy preview route removed - now using project-based preview at /projects/{project_id}/fullscreen

# Legacy tasks list route removed - now using /projects for project management

@router.post("/upload", response_class=HTMLResponse)
async def web_upload_file(
    request: Request,
    file: UploadFile = File(...)
):
    """Upload file via web interface"""
    try:
        # Validate file type
        allowed_types = [".docx", ".pdf", ".txt", ".md"]
        file_extension = "." + file.filename.split(".")[-1].lower()

        if file_extension not in allowed_types:
            return templates.TemplateResponse("upload_result.html", {
                "request": request,
                "success": False,
                "error": f"Unsupported file type. Allowed types: {', '.join(allowed_types)}"
            })

        # Read file content in thread pool to avoid blocking
        content = await file.read()

        # Process file in thread pool
        processed_content = await ppt_service.process_uploaded_file(
            filename=file.filename,
            content=content,
            file_type=file_extension
        )

        return templates.TemplateResponse("upload_result.html", {
            "request": request,
            "success": True,
            "filename": file.filename,
            "size": len(content),
            "type": file_extension,
            "processed_content": processed_content[:500] + "..." if len(processed_content) > 500 else processed_content
        })

    except Exception as e:
        return templates.TemplateResponse("upload_result.html", {
            "request": request,
            "success": False,
            "error": str(e)
        })

@router.get("/demo", response_class=HTMLResponse)
async def web_demo(
    request: Request
):
    """Demo page with sample PPT"""
    # Create a demo PPT
    demo_request = PPTGenerationRequest(
        scenario="technology",
        topic="äººå·¥æ™ºèƒ½æŠ€æœ¯å‘å±•è¶‹åŠ¿",
        requirements="é¢å‘æŠ€æœ¯äººå‘˜çš„æ·±åº¦åˆ†æ",
        network_mode=False,
        language="zh"
    )

    task_id = "demo-" + str(uuid.uuid4())[:8]
    result = await ppt_service.generate_ppt(task_id, demo_request)

    return templates.TemplateResponse("demo.html", {
        "request": request,
        "task_id": task_id,
        "outline": result.get("outline"),
        "slides_html": result.get("slides_html"),
        "demo_topic": demo_request.topic
    })

# Diagnostic endpoint for Playwright
@router.get("/api/diagnostics/playwright")
async def diagnose_playwright():
    """Diagnose Playwright availability"""
    import sys
    
    result = {
        "python_executable": sys.executable,
        "python_version": sys.version,
        "playwright_module_available": False,
        "playwright_async_api_available": False,
        "pdf_converter_available": False,
        "error": None
    }
    
    # Test 1: Import playwright module
    try:
        import playwright
        result["playwright_module_available"] = True
        result["playwright_location"] = playwright.__file__
    except ImportError as e:
        result["error"] = f"Cannot import playwright: {e}"
        return JSONResponse(result)
    
    # Test 2: Import async_api
    try:
        from playwright.async_api import async_playwright
        result["playwright_async_api_available"] = True
    except ImportError as e:
        result["error"] = f"Cannot import playwright.async_api: {e}"
        return JSONResponse(result)
    
    # Test 3: Check PDF converter
    try:
        pdf_converter = get_pdf_converter()
        result["pdf_converter_available"] = pdf_converter.is_available()
    except Exception as e:
        result["error"] = f"PDF converter check failed: {e}"
    
    return JSONResponse(result)

# New Project Management Routes

@router.get("/dashboard", response_class=HTMLResponse)
async def web_dashboard(
    request: Request
):
    """Project dashboard with overview"""
    try:
        # Get project statistics
        projects_response = await ppt_service.project_manager.list_projects(page=1, page_size=100)
        projects = projects_response.projects

        total_projects = len(projects)
        completed_projects = len([p for p in projects if p.status == "completed"])
        in_progress_projects = len([p for p in projects if p.status == "in_progress"])
        draft_projects = len([p for p in projects if p.status == "draft"])

        # Get recent projects (last 5)
        recent_projects = sorted(projects, key=lambda x: x.updated_at, reverse=True)[:5]

        # Get active TODO boards
        active_todo_boards = []
        for project in projects:
            if project.status == "in_progress" and project.todo_board:
                todo_board = await ppt_service.get_project_todo_board(project.project_id)
                if todo_board:
                    active_todo_boards.append(todo_board)

        return templates.TemplateResponse("project_dashboard.html", {
            "request": request,
            "total_projects": total_projects,
            "completed_projects": completed_projects,
            "in_progress_projects": in_progress_projects,
            "draft_projects": draft_projects,
            "recent_projects": recent_projects,
            "active_todo_boards": active_todo_boards[:3]  # Show max 3 boards
        })

    except Exception as e:
        return templates.TemplateResponse("error.html", {
            "request": request,
            "error": str(e)
        })

@router.get("/projects", response_class=HTMLResponse)
async def web_projects_list(
    request: Request,
    page: int = 1,
    status: str = None
):
    """List all projects"""
    try:
        projects_response = await ppt_service.project_manager.list_projects(
            page=page, page_size=10, status=status
        )

        return templates.TemplateResponse("projects_list.html", {
            "request": request,
            "projects": projects_response.projects,
            "total": projects_response.total,
            "page": projects_response.page,
            "page_size": projects_response.page_size,
            "status_filter": status
        })

    except Exception as e:
        return templates.TemplateResponse("error.html", {
            "request": request,
            "error": str(e)
        })

@router.get("/projects/{project_id}", response_class=HTMLResponse)
async def web_project_detail(
    request: Request,
    project_id: str
):
    """Project detail page"""
    try:
        project = await ppt_service.project_manager.get_project(project_id)
        if not project:
            return templates.TemplateResponse("error.html", {
                "request": request,
                "error": "Project not found"
            })

        todo_board = await ppt_service.get_project_todo_board(project_id)
        versions = await ppt_service.project_manager.get_project_versions(project_id)

        return templates.TemplateResponse("project_detail.html", {
            "request": request,
            "project": project,
            "todo_board": todo_board,
            "versions": versions
        })

    except Exception as e:
        return templates.TemplateResponse("error.html", {
            "request": request,
            "error": str(e)
        })

@router.get("/projects/{project_id}/todo", response_class=HTMLResponse)
async def web_project_todo_board(
    request: Request,
    project_id: str
):
    """TODO board page for a project with integrated editor"""
    try:
        # Validate project_id format (should be UUID-like)
        if project_id in ["template-selection", "todo", "edit", "preview", "fullscreen"]:
            error_msg = f"æ— æ•ˆçš„é¡¹ç›®ID: {project_id}.\n\n"
            error_msg += "å¯èƒ½çš„åŸå› :\n"
            error_msg += "1. URLæ ¼å¼é”™è¯¯,æ­£ç¡®æ ¼å¼åº”ä¸º: /projects/[é¡¹ç›®ID]/todo\n"
            error_msg += "2. æ‚¨å¯èƒ½è®¿é—®äº†é”™è¯¯çš„é“¾æ¥\n\n"
            error_msg += "å»ºè®®è§£å†³æ–¹æ¡ˆ:\n"
            error_msg += "â€¢ è¿”å›é¡¹ç›®åˆ—è¡¨é¡µé¢é€‰æ‹©æ­£ç¡®çš„é¡¹ç›®\n"
            error_msg += "â€¢ æ£€æŸ¥æµè§ˆå™¨åœ°å€æ ä¸­çš„URLæ˜¯å¦å®Œæ•´"

            return templates.TemplateResponse("error.html", {
                "request": request,
                "error": error_msg
            })

        # Check if project exists first
        project = await ppt_service.project_manager.get_project(project_id)
        if not project:
            return templates.TemplateResponse("error.html", {
                "request": request,
                "error": f"é¡¹ç›®ä¸å­˜åœ¨ (ID: {project_id}).è¯·æ£€æŸ¥é¡¹ç›®IDæ˜¯å¦æ­£ç¡®."
            })

        todo_board = await ppt_service.get_project_todo_board(project_id)
        if not todo_board:
            return templates.TemplateResponse("error.html", {
                "request": request,
                "error": f"é¡¹ç›® '{project.topic}' çš„TODOçœ‹æ¿ä¸å­˜åœ¨.è¯·è”ç³»æŠ€æœ¯æ”¯æŒ."
            })

        # Check if we should use the integrated editor version
        project = await ppt_service.project_manager.get_project(project_id)
        use_integrated_editor = (
            project and
            project.confirmed_requirements and
            len(todo_board.stages) > 2 and
            (todo_board.stages[1].status in ['running', 'completed'] or
             todo_board.stages[2].status in ['running', 'completed'])
        )

        # Also use integrated editor if PPT creation stage is about to start or running
        if (project and project.confirmed_requirements and len(todo_board.stages) > 2 and
            todo_board.stages[1].status == 'completed'):
            use_integrated_editor = True

        template_name = "todo_board_with_editor.html" if use_integrated_editor else "todo_board.html"

        # Ensure project is not None for template
        template_context = {
            "request": request,
            "todo_board": todo_board
        }

        # Only add project if it exists
        if project:
            template_context["project"] = project

        return templates.TemplateResponse(template_name, template_context)

    except Exception as e:
        return templates.TemplateResponse("error.html", {
            "request": request,
            "error": str(e)
        })



@router.get("/projects/{project_id}/fullscreen", response_class=HTMLResponse)
async def web_project_fullscreen(
    request: Request,
    project_id: str
):
    """Fullscreen preview of project PPT with modern presentation interface"""
    try:
        # ç›´æ¥ä»æ•°æ®åº“è·å–æœ€æ–°çš„é¡¹ç›®æ•°æ®,ç¡®ä¿æ•°æ®å®æ—¶æ€§
        from ..services.db_project_manager import DatabaseProjectManager
        db_manager = DatabaseProjectManager()
        project = await db_manager.get_project(project_id)

        if not project:
            return templates.TemplateResponse("error.html", {
                "request": request,
                "error": "é¡¹ç›®æœªæ‰¾åˆ°"
            })

        # æ£€æŸ¥æ˜¯å¦æœ‰å¹»ç¯ç‰‡æ•°æ®
        if not project.slides_data or len(project.slides_data) == 0:
            return templates.TemplateResponse("error.html", {
                "request": request,
                "error": "PPTå°šæœªç”Ÿæˆæˆ–æ— å¹»ç¯ç‰‡å†…å®¹"
            })

        # ä½¿ç”¨æ–°çš„åˆ†äº«æ¼”ç¤ºæ¨¡æ¿
        return templates.TemplateResponse("project_fullscreen_presentation.html", {
            "request": request,
            "project": project,
            "slides_count": len(project.slides_data)
        })

    except Exception as e:
        logger.error(f"Error in fullscreen presentation: {e}")
        return templates.TemplateResponse("error.html", {
            "request": request,
            "error": f"åŠ è½½æ¼”ç¤ºæ—¶å‡ºé”™: {str(e)}"
        })

# Share link routes removed - åˆ†äº«é“¾æ¥åŠŸèƒ½å·²åˆ é™¤

@router.get("/api/projects/{project_id}/slides-data")
async def get_project_slides_data(
    project_id: str
):
    """è·å–é¡¹ç›®æœ€æ–°çš„å¹»ç¯ç‰‡æ•°æ® - ç”¨äºåˆ†äº«æ¼”ç¤ºå®æ—¶æ›´æ–°"""
    try:
        # ç›´æ¥ä»æ•°æ®åº“è·å–æœ€æ–°æ•°æ®
        from ..services.db_project_manager import DatabaseProjectManager
        db_manager = DatabaseProjectManager()
        project = await db_manager.get_project(project_id)

        if not project:
            raise HTTPException(status_code=404, detail="é¡¹ç›®æœªæ‰¾åˆ°")

        if not project.slides_data or len(project.slides_data) == 0:
            return {
                "status": "no_slides",
                "message": "PPTå°šæœªç”Ÿæˆ",
                "slides_data": [],
                "total_slides": 0
            }

        return {
            "status": "success",
            "slides_data": project.slides_data,
            "total_slides": len(project.slides_data),
            "project_title": project.title,
            "updated_at": project.updated_at
        }

    except Exception as e:
        logger.error(f"Error getting slides data: {e}")
        raise HTTPException(status_code=500, detail=f"è·å–å¹»ç¯ç‰‡æ•°æ®å¤±è´¥: {str(e)}")


# Share link generation, disable, and info routes removed

@router.get("/test/slides-navigation", response_class=HTMLResponse)
async def test_slides_navigation(
    request: Request
):
    """æµ‹è¯•å¹»ç¯ç‰‡å¯¼èˆªåŠŸèƒ½"""
    with open("test_slides_navigation.html", "r", encoding="utf-8") as f:
        content = f.read()
    return HTMLResponse(content=content)

@router.get("/temp/{file_path:path}")
async def serve_temp_file(
    file_path: str
):
    """Serve temporary slide files"""
    try:
        # Construct the full path to the temp file using system temp directory
        import tempfile
        temp_dir = Path(tempfile.gettempdir()) / "ai_slides"
        full_path = temp_dir / file_path

        # Security check: ensure the file is within the temp directory
        if not str(full_path.resolve()).startswith(str(temp_dir.resolve())):
            raise HTTPException(status_code=403, detail="Access denied")

        # Check if file exists
        if not full_path.exists():
            raise HTTPException(status_code=404, detail="File not found")

        # Return the file
        return FileResponse(
            path=str(full_path),
            media_type="text/html; charset=utf-8",
            headers={"Cache-Control": "no-cache"}
        )

    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@router.post("/projects/create", response_class=HTMLResponse)
async def web_create_project(
    request: Request,
    scenario: str = Form(...),
    topic: str = Form(...),
    requirements: str = Form(None),
    language: str = Form("zh"),
    network_mode: bool = Form(False)
):
    """Create new project via web interface"""
    try:
        # Create project request
        project_request = PPTGenerationRequest(
            scenario=scenario,
            topic=topic,
            requirements=requirements,
            network_mode=network_mode,
            language=language
        )

        # Create project with TODO board (without starting workflow yet)
        project = await ppt_service.project_manager.create_project(project_request)

        # Update project status to in_progress
        await ppt_service.project_manager.update_project_status(project.project_id, "in_progress")

        # Redirect directly to TODO page without showing redirect page
        from fastapi.responses import RedirectResponse
        return RedirectResponse(
            url=f"/projects/{project.project_id}/todo",
            status_code=302
        )

    except Exception as e:
        return templates.TemplateResponse("error.html", {
            "request": request,
            "error": str(e)
        })

@router.post("/projects/{project_id}/start-workflow")
async def start_project_workflow(
    project_id: str
):
    """Start the AI workflow for a project (only if requirements are confirmed)"""
    try:
        # Get project
        project = await ppt_service.project_manager.get_project(project_id)
        if not project:
            raise HTTPException(status_code=404, detail="Project not found")

        # Check if requirements are confirmed
        if not project.confirmed_requirements:
            return {"status": "waiting", "message": "Waiting for requirements confirmation"}

        # Extract network_mode from project metadata
        network_mode = False
        if project.project_metadata and isinstance(project.project_metadata, dict):
            network_mode = project.project_metadata.get("network_mode", False)

        # Create project request from project data
        confirmed_requirements = project.confirmed_requirements or {}
        project_request = PPTGenerationRequest(
            scenario=project.scenario,
            topic=project.topic,
            requirements=project.requirements,
            language="zh",  # Default language
            network_mode=network_mode,
            target_audience=confirmed_requirements.get('target_audience', 'æ™®é€šå¤§ä¼—'),
            ppt_style=confirmed_requirements.get('ppt_style', 'general'),
            custom_style_prompt=confirmed_requirements.get('custom_style_prompt'),
            description=confirmed_requirements.get('description')
        )

        # Start the workflow in background
        asyncio.create_task(ppt_service._execute_project_workflow(project_id, project_request))

        return {"status": "success", "message": "Workflow started"}

    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@router.get("/projects/{project_id}/requirements", response_class=HTMLResponse)
async def project_requirements_page(
    request: Request,
    project_id: str
):
    """Show project requirements confirmation page"""
    try:
        # Get project
        project = await ppt_service.project_manager.get_project(project_id)
        if not project:
            raise HTTPException(status_code=404, detail="Project not found")

        # æä¾›é»˜è®¤çš„å±•ç¤ºç±»å‹é€‰é¡¹,ä¸å†è°ƒç”¨AIç”Ÿæˆå»ºè®®
        default_type_options = [
            "æŠ€æœ¯åˆ†äº«",
            "äº§å“ä»‹ç»",
            "å­¦æœ¯æŠ¥å‘Š",
            "å•†ä¸šæ±‡æŠ¥",
            "æ•™å­¦è¯¾ä»¶",
            "é¡¹ç›®å±•ç¤º",
            "æ•°æ®åˆ†æ",
            "ç»¼åˆä»‹ç»"
        ]

        return templates.TemplateResponse("project_requirements.html", {
            "request": request,
            "project": project,
            "ai_suggestions": {
                "type_options": default_type_options
            }
        })

    except Exception as e:
        return templates.TemplateResponse("error.html", {
            "request": request,
            "error": str(e)
        })

# ç§»é™¤AIç”Ÿæˆéœ€æ±‚å»ºè®®çš„APIç«¯ç‚¹,æ”¹ä¸ºä½¿ç”¨é»˜è®¤é€‰é¡¹

@router.get("/projects/{project_id}/outline-stream")
async def stream_outline_generation(
    project_id: str
):
    """Stream outline generation for a project"""
    try:
        project = await ppt_service.project_manager.get_project(project_id)
        if not project:
            raise HTTPException(status_code=404, detail="Project not found")

        async def generate():
            try:
                async for chunk in ppt_service.generate_outline_streaming(project_id):
                    yield chunk
            except Exception as e:
                import json
                error_response = {'error': str(e)}
                yield f"data: {json.dumps(error_response)}\n\n"

        return StreamingResponse(generate(), media_type="text/plain")

    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@router.post("/projects/{project_id}/generate-outline")
async def generate_outline(
    project_id: str
):
    """Generate outline for a project (non-streaming)"""
    try:
        project = await ppt_service.project_manager.get_project(project_id)
        if not project:
            raise HTTPException(status_code=404, detail="Project not found")

        # Check if project has confirmed requirements
        if not project.confirmed_requirements:
            return {
                "status": "error",
                "error": "é¡¹ç›®éœ€æ±‚å°šæœªç¡®è®¤,è¯·å…ˆå®Œæˆéœ€æ±‚ç¡®è®¤æ­¥éª¤"
            }

        # Create PPTGenerationRequest from project data
        confirmed_requirements = project.confirmed_requirements

        # Extract network_mode from project metadata
        network_mode = False
        if project.project_metadata and isinstance(project.project_metadata, dict):
            network_mode = project.project_metadata.get("network_mode", False)

        project_request = PPTGenerationRequest(
            scenario=project.scenario,
            topic=confirmed_requirements.get('topic', project.topic),
            requirements=project.requirements,
            language="zh",  # Default language
            network_mode=network_mode,
            target_audience=confirmed_requirements.get('target_audience', 'æ™®é€šå¤§ä¼—'),
            ppt_style=confirmed_requirements.get('ppt_style', 'general'),
            custom_style_prompt=confirmed_requirements.get('custom_style_prompt'),
            description=confirmed_requirements.get('description')
        )

        # Extract page count settings from confirmed requirements
        page_count_settings = confirmed_requirements.get('page_count_settings', {})

        # Generate outline using AI with page count settings
        outline = await ppt_service.generate_outline(project_request, page_count_settings)

        # Convert outline to dict format
        outline_dict = {
            "title": outline.title,
            "slides": outline.slides,
            "metadata": outline.metadata
        }

        # Format as JSON
        import json
        formatted_json = json.dumps(outline_dict, ensure_ascii=False, indent=2)

        # Update outline generation stage
        await ppt_service._update_outline_generation_stage(project_id, outline_dict)

        return {
            "status": "success",
            "outline_content": formatted_json,
            "message": "Outline generated successfully"
        }

    except Exception as e:
        logger.error(f"Error generating outline: {e}")
        return {
            "status": "error",
            "error": str(e)
        }

@router.post("/projects/{project_id}/regenerate-outline")
async def regenerate_outline(
    project_id: str,
    request: Request
):
    """Regenerate outline for a project (overwrites existing outline) with optional custom requirements"""
    try:
        # Get request body to extract custom requirements if provided
        request_data = {}
        try:
            request_data = await request.json()
        except:
            pass  # If no body or invalid JSON, use empty dict
        
        custom_requirements = request_data.get('custom_requirements', '')
        
        project = await ppt_service.project_manager.get_project(project_id)
        if not project:
            raise HTTPException(status_code=404, detail="Project not found")

        # Check if project has confirmed requirements
        if not project.confirmed_requirements:
            return {
                "status": "error",
                "error": "é¡¹ç›®éœ€æ±‚å°šæœªç¡®è®¤,è¯·å…ˆå®Œæˆéœ€æ±‚ç¡®è®¤æ­¥éª¤"
            }

        # Create project request from confirmed requirements
        confirmed_requirements = project.confirmed_requirements
        
        # å¦‚æœæä¾›äº†è‡ªå®šä¹‰éœ€æ±‚,å°†å…¶è¿½åŠ æˆ–è¦†ç›–åŸæœ‰éœ€æ±‚
        final_requirements = confirmed_requirements.get('requirements', project.requirements)
        if custom_requirements:
            # å°†è‡ªå®šä¹‰éœ€æ±‚è¿½åŠ åˆ°åŸæœ‰éœ€æ±‚
            if final_requirements:
                final_requirements = f"{final_requirements}\n\nã€æœ¬æ¬¡é‡æ–°ç”Ÿæˆçš„é¢å¤–è¦æ±‚ã€‘\n{custom_requirements}"
            else:
                final_requirements = custom_requirements
        
        project_request = PPTGenerationRequest(
            scenario=confirmed_requirements.get('scenario', 'general'),
            topic=confirmed_requirements.get('topic', project.topic),
            requirements=final_requirements,
            language="zh",  # Default language
            network_mode=confirmed_requirements.get('network_mode', False),
            target_audience=confirmed_requirements.get('target_audience', 'æ™®é€šå¤§ä¼—'),
            ppt_style=confirmed_requirements.get('ppt_style', 'general'),
            custom_style_prompt=confirmed_requirements.get('custom_style_prompt'),
            description=confirmed_requirements.get('description')
        )

        # Extract page count settings from confirmed requirements
        page_count_settings = confirmed_requirements.get('page_count_settings', {})

        # Check if this is a file-based project
        is_file_project = confirmed_requirements.get('content_source') == 'file'

        if is_file_project:
            # Check if file path exists
            file_path = confirmed_requirements.get('file_path')
            if not file_path:
                return {
                    "status": "error",
                    "error": "æ–‡ä»¶è·¯å¾„ä¿¡æ¯ä¸¢å¤±,è¯·é‡æ–°ä¸Šä¼ æ–‡ä»¶å¹¶ç¡®è®¤éœ€æ±‚"
                }

            # Use file-based outline generation
            file_request = FileOutlineGenerationRequest(
                file_path=file_path,
                filename=confirmed_requirements.get('filename', 'uploaded_file'),
                topic=project_request.topic,
                scenario=project_request.scenario,
                requirements=confirmed_requirements.get('requirements', ''),
                target_audience=confirmed_requirements.get('target_audience', 'æ™®é€šå¤§ä¼—'),
                page_count_mode=page_count_settings.get('mode', 'ai_decide'),
                min_pages=page_count_settings.get('min_pages', 5),
                max_pages=page_count_settings.get('max_pages', 20),
                fixed_pages=page_count_settings.get('fixed_pages', 10),
                ppt_style=confirmed_requirements.get('ppt_style', 'general'),
                custom_style_prompt=confirmed_requirements.get('custom_style_prompt'),
                file_processing_mode=confirmed_requirements.get('file_processing_mode', 'markitdown'),
                content_analysis_depth=confirmed_requirements.get('content_analysis_depth', 'standard')
            )

            result = await ppt_service.generate_outline_from_file(file_request)

            if not result.success:
                return {
                    "status": "error",
                    "error": result.error or "æ–‡ä»¶å¤§çº²ç”Ÿæˆå¤±è´¥"
                }

            # Update outline generation stage
            await ppt_service._update_outline_generation_stage(project_id, result.outline)

            # Format outline as JSON string
            import json
            outline_content = json.dumps(result.outline, ensure_ascii=False, indent=2)

            return {
                "status": "success",
                "outline_content": outline_content,
                "message": "File-based outline regenerated successfully"
            }
        else:
            # Use standard outline generation
            outline = await ppt_service.generate_outline(project_request, page_count_settings)

            # Convert outline to dict format
            outline_dict = {
                "title": outline.title,
                "slides": outline.slides,
                "metadata": outline.metadata
            }

            # Format as JSON
            import json
            formatted_json = json.dumps(outline_dict, ensure_ascii=False, indent=2)

            # Update outline generation stage
            await ppt_service._update_outline_generation_stage(project_id, outline_dict)

            return {
                "status": "success",
                "outline_content": formatted_json,
                "message": "Outline regenerated successfully"
            }

    except Exception as e:
        logger.error(f"Error regenerating outline: {e}")
        return {
            "status": "error",
            "error": str(e)
        }

@router.post("/projects/{project_id}/generate-file-outline")
async def generate_file_outline(
    project_id: str
):
    """Generate outline from uploaded file (non-streaming)"""
    try:
        project = await ppt_service.project_manager.get_project(project_id)
        if not project:
            raise HTTPException(status_code=404, detail="Project not found")

        # Check if project has file-generated outline
        file_generated_outline = None

        # é¦–å…ˆæ£€æŸ¥é¡¹ç›®çš„outlineå­—æ®µ
        if project.outline and project.outline.get('slides'):
            # æ£€æŸ¥æ˜¯å¦æ˜¯ä»æ–‡ä»¶ç”Ÿæˆçš„å¤§çº²
            metadata = project.outline.get('metadata', {})
            if metadata.get('generated_with_summeryfile') or metadata.get('generated_with_file'):
                file_generated_outline = project.outline
                logger.info(f"Project {project_id} has file-generated outline in project.outline, using it")

        # å¦‚æœé¡¹ç›®outlineä¸­æ²¡æœ‰,å†æ£€æŸ¥confirmed_requirements
        if not file_generated_outline and project.confirmed_requirements and project.confirmed_requirements.get('file_generated_outline'):
            file_generated_outline = project.confirmed_requirements['file_generated_outline']
            logger.info(f"Project {project_id} has file-generated outline in confirmed_requirements, using it")

        # If no existing outline but file upload is configured, wait a bit and check again
        if not file_generated_outline and project.confirmed_requirements and project.confirmed_requirements.get('content_source') == 'file':
            logger.info(f"Project {project_id} has file upload but no outline yet, waiting for file processing...")

            # Wait for file processing to complete (it should be done during requirements confirmation)
            import asyncio
            max_wait_time = 10  # Maximum wait time in seconds
            wait_interval = 1   # Check every 1 second

            for i in range(max_wait_time):
                await asyncio.sleep(wait_interval)

                # Refresh project data
                project = await ppt_service.project_manager.get_project(project_id)
                if project.confirmed_requirements and project.confirmed_requirements.get('file_generated_outline'):
                    file_generated_outline = project.confirmed_requirements['file_generated_outline']
                    logger.info(f"Project {project_id} file outline found after waiting {i+1} seconds")
                    break

            if not file_generated_outline:
                logger.warning(f"Project {project_id} file outline not found after waiting {max_wait_time} seconds")

        if file_generated_outline:
            # Return the existing file-generated outline
            import json
            existing_outline = {
                "title": file_generated_outline.get('title', project.topic),
                "slides": file_generated_outline.get('slides', []),
                "metadata": file_generated_outline.get('metadata', {})
            }

            # Ensure metadata includes correct identification
            if 'metadata' not in existing_outline:
                existing_outline['metadata'] = {}
            existing_outline['metadata']['generated_with_summeryfile'] = True
            existing_outline['metadata']['generated_at'] = time.time()

            formatted_json = json.dumps(existing_outline, ensure_ascii=False, indent=2)

            # Update outline generation stage
            await ppt_service._update_outline_generation_stage(project_id, existing_outline)

            return {
                "status": "success",
                "outline_content": formatted_json,
                "message": "File outline generated successfully"
            }
        else:
            # Check if there's an uploaded file that needs processing
            if (project.confirmed_requirements and
                (project.confirmed_requirements.get('uploaded_files') or
                 project.confirmed_requirements.get('content_source') == 'file')):
                logger.info(f"Project {project_id} has uploaded files, starting file outline generation")

                # Start file outline generation using summeryfile
                try:
                    # Create a request object for file outline generation
                    from ..api.models import FileOutlineGenerationRequest

                    # Get file information from confirmed requirements
                    uploaded_files = project.confirmed_requirements.get('uploaded_files', [])
                    if uploaded_files:
                        file_info = uploaded_files[0]  # Use first file
                        # ä½¿ç”¨ç¡®è®¤çš„è¦æ±‚æˆ–é¡¹ç›®åˆ›å»ºæ—¶çš„è¦æ±‚ä½œä¸ºfallback
                        confirmed_reqs = project.confirmed_requirements.get('requirements', '')
                        project_reqs = project.requirements or ''
                        final_reqs = confirmed_reqs or project_reqs

                        file_request = FileOutlineGenerationRequest(
                            filename=file_info.get('filename', 'uploaded_file'),
                            file_path=file_info.get('file_path', ''),
                            topic=project.topic,
                            scenario='general',
                            requirements=final_reqs,
                            target_audience=project.confirmed_requirements.get('target_audience', 'æ™®é€šå¤§ä¼—'),
                            page_count_mode=project.confirmed_requirements.get('page_count_settings', {}).get('mode', 'ai_decide'),
                            min_pages=project.confirmed_requirements.get('page_count_settings', {}).get('min_pages', 8),
                            max_pages=project.confirmed_requirements.get('page_count_settings', {}).get('max_pages', 15),
                            fixed_pages=project.confirmed_requirements.get('page_count_settings', {}).get('fixed_pages', 10),
                            ppt_style=project.confirmed_requirements.get('ppt_style', 'general'),
                            custom_style_prompt=project.confirmed_requirements.get('custom_style_prompt'),
                            file_processing_mode=project.confirmed_requirements.get('file_processing_mode', 'markitdown'),
                            content_analysis_depth=project.confirmed_requirements.get('content_analysis_depth', 'standard')
                        )

                        # Generate outline from file using summeryfile
                        outline_response = await ppt_service.generate_outline_from_file(file_request)

                        if outline_response.success and outline_response.outline:
                            # Format the generated outline
                            import json
                            formatted_outline = outline_response.outline

                            # Ensure metadata includes correct identification
                            if 'metadata' not in formatted_outline:
                                formatted_outline['metadata'] = {}
                            formatted_outline['metadata']['generated_with_summeryfile'] = True
                            formatted_outline['metadata']['generated_at'] = time.time()

                            formatted_json = json.dumps(formatted_outline, ensure_ascii=False, indent=2)

                            # Update outline generation stage
                            await ppt_service._update_outline_generation_stage(project_id, formatted_outline)

                            return {
                                "status": "success",
                                "outline_content": formatted_json,
                                "message": "File outline generated successfully"
                            }
                        else:
                            error_msg = outline_response.error if hasattr(outline_response, 'error') else "Unknown error"
                            return {
                                "status": "error",
                                "error": f"Failed to generate outline from uploaded file: {error_msg}"
                            }
                    else:
                        return {
                            "status": "error",
                            "error": "No uploaded file information found in project requirements."
                        }

                except Exception as gen_error:
                    logger.error(f"Error generating outline from file: {gen_error}")
                    return {
                        "status": "error",
                        "error": f"Failed to generate outline from file: {str(gen_error)}"
                    }
            else:
                # No file outline found and no uploaded files
                return {
                    "status": "error",
                    "error": "No file outline found. Please ensure you uploaded a file during requirements confirmation."
                }

    except Exception as e:
        logger.error(f"Error generating file outline: {e}")
        return {
            "status": "error",
            "error": str(e)
        }

@router.post("/projects/{project_id}/update-outline")
async def update_project_outline(
    project_id: str,
    request: Request
):
    """Update project outline content"""
    try:
        data = await request.json()
        outline_content = data.get('outline_content', '')

        success = await ppt_service.update_project_outline(project_id, outline_content)
        if success:
            return {"status": "success", "message": "Outline updated"}
        else:
            raise HTTPException(status_code=500, detail="Failed to update outline")

    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@router.post("/projects/{project_id}/confirm-outline")
async def confirm_project_outline(
    project_id: str
):
    """Confirm project outline and enable PPT generation"""
    try:
        success = await ppt_service.confirm_project_outline(project_id)
        if success:
            return {"status": "success", "message": "Outline confirmed"}
        else:
            raise HTTPException(status_code=500, detail="Failed to confirm outline")

    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@router.get("/projects/{project_id}/todo-editor")
async def web_project_todo_editor(
    request: Request,
    project_id: str,
    auto_start: bool = False
):
    """Project TODO board with editor"""
    try:
        project = await ppt_service.project_manager.get_project(project_id)
        if not project:
            return templates.TemplateResponse("error.html", {
                "request": request,
                "error": "Project not found"
            })

        return templates.TemplateResponse("todo_board_with_editor.html", {
            "request": request,
            "todo_board": project.todo_board,
            "project": project,
            "auto_start": auto_start
        })

    except Exception as e:
        return templates.TemplateResponse("error.html", {
            "request": request,
            "error": str(e)
        })

@router.post("/projects/{project_id}/confirm-requirements")
async def confirm_project_requirements(
    request: Request,
    project_id: str,
    topic: str = Form(...),
    audience_type: str = Form(...),
    custom_audience: str = Form(None),
    page_count_mode: str = Form("ai_decide"),
    min_pages: int = Form(8),
    max_pages: int = Form(15),
    fixed_pages: int = Form(10),
    ppt_style: str = Form("general"),
    custom_style_prompt: str = Form(None),
    description: str = Form(None),
    content_source: str = Form("manual"),
    file_upload: List[UploadFile] = File(None),
    file_processing_mode: str = Form("markitdown"),
    content_analysis_depth: str = Form("standard")
):
    """Confirm project requirements and generate TODO list - æ”¯æŒå¤šæ–‡ä»¶ä¸Šä¼ å’Œè”ç½‘æœç´¢é›†æˆ"""
    try:
        # Get project to access original requirements
        project = await ppt_service.project_manager.get_project(project_id)
        if not project:
            raise HTTPException(status_code=404, detail="Project not found")

        # Extract network_mode from project metadata (set during project creation)
        network_mode = False
        if project.project_metadata and isinstance(project.project_metadata, dict):
            network_mode = project.project_metadata.get("network_mode", False)

        # Process audience information
        target_audience = audience_type
        if audience_type == "è‡ªå®šä¹‰" and custom_audience:
            target_audience = custom_audience

        # Handle file upload if content source is file
        file_outline = None
        if content_source == "file" and file_upload:
            # Process uploaded files (support multiple files) and generate outline
            # ä½¿ç”¨é¡¹ç›®åˆ›å»ºæ—¶çš„ network_mode å‚æ•°
            file_outline = await _process_uploaded_files_for_outline(
                file_upload, topic, target_audience, page_count_mode, min_pages, max_pages,
                fixed_pages, ppt_style, custom_style_prompt,
                file_processing_mode, content_analysis_depth, project.requirements,
                enable_web_search=network_mode,  # ä½¿ç”¨é¡¹ç›®çš„ network_mode
                scenario=project.scenario,  # ä¼ é€’åœºæ™¯å‚æ•°
                language="zh"  # ä¼ é€’è¯­è¨€å‚æ•°
            )

            # Update topic if it was extracted from file
            if file_outline and file_outline.get('title') and not topic.strip():
                topic = file_outline['title']

        # Process page count settings
        page_count_settings = {
            "mode": page_count_mode,
            "min_pages": min_pages if page_count_mode == "custom_range" else None,
            "max_pages": max_pages if page_count_mode == "custom_range" else None,
            "fixed_pages": fixed_pages if page_count_mode == "fixed" else None
        }

        # Update project with confirmed requirements
        confirmed_requirements = {
            "topic": topic,
            "requirements": project.requirements,  # ä½¿ç”¨é¡¹ç›®åˆ›å»ºæ—¶çš„å…·ä½“è¦æ±‚
            "target_audience": target_audience,
            "audience_type": audience_type,
            "custom_audience": custom_audience if audience_type == "è‡ªå®šä¹‰" else None,
            "page_count_settings": page_count_settings,
            "ppt_style": ppt_style,
            "custom_style_prompt": custom_style_prompt if ppt_style == "custom" else None,
            "description": description,
            "content_source": content_source,
            "file_processing_mode": file_processing_mode if content_source == "file" else None,
            "content_analysis_depth": content_analysis_depth if content_source == "file" else None,
            "file_generated_outline": file_outline
        }

        # å¦‚æœæ˜¯æ–‡ä»¶é¡¹ç›®,ä¿å­˜æ–‡ä»¶ä¿¡æ¯
        if content_source == "file" and file_outline and 'file_info' in file_outline:
            file_info = file_outline['file_info']
            file_path = file_info.get('file_path') or file_info.get('merged_file_path')
            filename = file_info.get('filename') or file_info.get('merged_filename')
            uploaded_files = file_info.get('uploaded_files')

            file_metadata = {}
            if file_path:
                file_metadata["file_path"] = file_path
            if filename:
                file_metadata["filename"] = filename
            if uploaded_files:
                file_metadata["uploaded_files"] = uploaded_files

            if file_metadata:
                confirmed_requirements.update(file_metadata)

        # Store confirmed requirements in project
        # ç›´æ¥ç¡®è®¤éœ€æ±‚å¹¶æ›´æ–°TODOæ¿,æ— éœ€AIç”Ÿæˆå¾…åŠæ¸…å•
        success = await ppt_service.confirm_requirements_and_update_workflow(project_id, confirmed_requirements)

        if not success:
            raise Exception("éœ€æ±‚ç¡®è®¤å¤±è´¥")

        # Return JSON success response for AJAX request
        from fastapi.responses import JSONResponse
        return JSONResponse({
            "status": "success",
            "message": "éœ€æ±‚ç¡®è®¤å®Œæˆ",
            "redirect_url": f"/projects/{project_id}/todo"
        })

    except Exception as e:
        from fastapi.responses import JSONResponse
        return JSONResponse({
            "status": "error",
            "message": str(e)
        }, status_code=500)

@router.get("/projects/{project_id}/stage-stream/{stage_id}")
async def stream_stage_response(
    project_id: str,
    stage_id: str
):
    """Stream AI response for a complete stage"""

    async def generate_stage_stream():
        try:
            # Get project and stage info
            project = await ppt_service.project_manager.get_project(project_id)
            if not project:
                yield f"data: {json.dumps({'error': 'Project not found'})}\n\n"
                return

            if not project.confirmed_requirements:
                yield f"data: {json.dumps({'error': 'Project requirements not confirmed'})}\n\n"
                return

            todo_board = await ppt_service.get_project_todo_board(project_id)
            if not todo_board:
                yield f"data: {json.dumps({'error': 'TODO board not found'})}\n\n"
                return

            # Find the stage
            stage = None
            for s in todo_board.stages:
                if s.id == stage_id:
                    stage = s
                    break

            if not stage:
                yield f"data: {json.dumps({'error': 'Stage not found'})}\n\n"
                return

            # Extract confirmed requirements from project
            confirmed_requirements = project.confirmed_requirements

            # Check if stage is already running or completed
            if stage.status == "running":
                yield f"data: {json.dumps({'error': 'Stage is already running'})}\n\n"
                return
            elif stage.status == "completed":
                yield f"data: {json.dumps({'error': 'Stage is already completed'})}\n\n"
                return

            # Update stage status to running
            await ppt_service.project_manager.update_stage_status(
                project_id, stage_id, "running", 0.0
            )

            # Execute the complete stage using the enhanced service
            try:
                if stage_id == "outline_generation":
                    response_content = await ppt_service._execute_outline_generation(
                        project_id, confirmed_requirements, ppt_service._load_prompts_md_system_prompt()
                    )
                elif stage_id == "ppt_creation":
                    response_content = await ppt_service._execute_ppt_creation(
                        project_id, confirmed_requirements, ppt_service._load_prompts_md_system_prompt()
                    )
                else:
                    # Fallback for other stages
                    response_content = await ppt_service._execute_general_stage(
                        project_id, stage_id, confirmed_requirements
                    )

                # Stream the response word by word for better UX
                if isinstance(response_content, dict):
                    content_text = response_content.get('message', str(response_content))
                else:
                    content_text = str(response_content)

                words = content_text.split()
                for i, word in enumerate(words):
                    yield f"data: {json.dumps({'content': word + ' ', 'done': False})}\n\n"
                    await asyncio.sleep(0.05)  # Small delay for streaming effect

            except Exception as e:
                # Fallback to basic stage execution
                prompt = f"""
ä½œä¸ºPPTç”ŸæˆåŠ©æ‰‹,è¯·å®Œæˆä»¥ä¸‹é˜¶æ®µä»»åŠ¡:

é¡¹ç›®ä¸»é¢˜:{project.topic}
é¡¹ç›®åœºæ™¯:{project.scenario}
é¡¹ç›®è¦æ±‚:{project.requirements or 'æ— ç‰¹æ®Šè¦æ±‚'}

å½“å‰é˜¶æ®µ:{stage.name}
é˜¶æ®µæè¿°:{stage.description}

è¯·æ ¹æ®ä»¥ä¸Šä¿¡æ¯å®Œæˆå½“å‰é˜¶æ®µçš„å®Œæ•´ä»»åŠ¡,å¹¶æä¾›è¯¦ç»†çš„æ‰§è¡Œç»“æœ.
"""

                # Stream AI response using real streaming
                async for chunk in ppt_service.ai_provider.stream_text_completion(
                    prompt=prompt,
                    max_tokens=2000,
                    temperature=0.7
                ):
                    if chunk:
                        yield f"data: {json.dumps({'content': chunk, 'done': False})}\n\n"

            # Update stage status to completed
            await ppt_service.project_manager.update_stage_status(
                project_id, stage_id, "completed", 100.0
            )

            # Send completion signal
            yield f"data: {json.dumps({'content': '', 'done': True})}\n\n"

        except Exception as e:
            yield f"data: {json.dumps({'error': str(e)})}\n\n"

    return StreamingResponse(
        generate_stage_stream(),
        media_type="text/plain",
        headers={"Cache-Control": "no-cache", "Connection": "keep-alive"}
    )



@router.get("/projects/{project_id}/edit", response_class=HTMLResponse)
async def edit_project_ppt(
    request: Request,
    project_id: str
):
    """Edit PPT slides with advanced editor"""
    try:
        project = await ppt_service.project_manager.get_project(project_id)
        if not project:
            raise HTTPException(status_code=404, detail="Project not found")

        # å…è®¸ç¼–è¾‘å™¨åœ¨PPTç”Ÿæˆè¿‡ç¨‹ä¸­æ˜¾ç¤º,æä¾›æ›´å¥½çš„ç”¨æˆ·ä½“éªŒ
        # å¦‚æœæ²¡æœ‰slides_data,åˆ›å»ºä¸€ä¸ªç©ºçš„ç»“æ„ä¾›ç¼–è¾‘å™¨ä½¿ç”¨
        if not project.slides_data:
            project.slides_data = []

        return templates.TemplateResponse("project_slides_editor.html", {
            "request": request,
            "project": project,
            "enable_auto_layout_repair": ai_config.enable_auto_layout_repair
        })

    except Exception as e:
        return templates.TemplateResponse("error.html", {
            "request": request,
            "error": str(e)
        })

@router.post("/api/projects/{project_id}/update-html")
async def update_project_html(
    project_id: str,
    request: Request
):
    """Update project HTML content and mark all slides as user-edited"""
    try:
        data = await request.json()
        slides_html = data.get('slides_html', '')

        project = await ppt_service.project_manager.get_project(project_id)
        if not project:
            raise HTTPException(status_code=404, detail="Project not found")

        # Update project HTML
        project.slides_html = slides_html
        project.updated_at = time.time()

        # è§£æHTMLå†…å®¹,æå–å„ä¸ªé¡µé¢å¹¶æ ‡è®°ä¸ºç”¨æˆ·ç¼–è¾‘
        if project.slides_data and slides_html:
            try:
                # è§£æHTMLå†…å®¹,æå–å„ä¸ªé¡µé¢
                updated_slides_data = await _extract_slides_from_html(slides_html, project.slides_data)

                # æ ‡è®°æ‰€æœ‰é¡µé¢ä¸ºç”¨æˆ·ç¼–è¾‘çŠ¶æ€
                for slide_data in updated_slides_data:
                    slide_data["is_user_edited"] = True

                # æ›´æ–°é¡¹ç›®çš„slides_data
                project.slides_data = updated_slides_data

                logger.info(f"Marked {len(updated_slides_data)} slides as user-edited for project {project_id}")

            except Exception as parse_error:
                logger.warning(f"Failed to parse HTML content for slide extraction: {parse_error}")
                # å¦‚æœè§£æå¤±è´¥,è‡³å°‘æ ‡è®°ç°æœ‰çš„slides_dataä¸ºç”¨æˆ·ç¼–è¾‘
                if project.slides_data:
                    for slide_data in project.slides_data:
                        slide_data["is_user_edited"] = True

        # ä¿å­˜æ›´æ–°çš„HTMLå’Œslides_dataåˆ°æ•°æ®åº“
        try:
            from ..services.db_project_manager import DatabaseProjectManager
            db_manager = DatabaseProjectManager()

            # ä¿å­˜å¹»ç¯ç‰‡HTMLå’Œæ•°æ®åˆ°æ•°æ®åº“
            save_success = await db_manager.save_project_slides(
                project_id,
                project.slides_html,
                project.slides_data or []
            )

            if save_success:
                logger.info(f"Successfully saved updated HTML and slides data to database for project {project_id}")
            else:
                logger.error(f"Failed to save updated HTML and slides data to database for project {project_id}")

        except Exception as save_error:
            logger.error(f"Exception while saving updated HTML and slides data to database: {save_error}")
            # ç»§ç»­è¿”å›æˆåŠŸ,å› ä¸ºå†…å­˜ä¸­çš„æ•°æ®å·²ç»æ›´æ–°

        return {"status": "success", "message": "HTML updated successfully and slides marked as user-edited"}

    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@router.get("/api/projects/{project_id}")
async def get_project_data(
    project_id: str
):
    """Get project data for real-time updates"""
    try:
        project = await ppt_service.project_manager.get_project(project_id)
        if not project:
            raise HTTPException(status_code=404, detail="Project not found")

        return {
            "project_id": project.project_id,
            "title": project.title,
            "status": project.status,
            "slides_data": project.slides_data or [],
            "slides_count": len(project.slides_data) if project.slides_data else 0,
            "updated_at": project.updated_at
        }

    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@router.put("/api/projects/{project_id}/slides")
async def update_project_slides(
    project_id: str,
    request: Request
):
    """Update project slides data"""
    try:
        logger.info(f"ğŸ”„ å¼€å§‹æ›´æ–°é¡¹ç›® {project_id} çš„å¹»ç¯ç‰‡æ•°æ®")

        data = await request.json()
        slides_data = data.get('slides_data', [])

        logger.info(f"ğŸ“Š æ¥æ”¶åˆ° {len(slides_data)} é¡µå¹»ç¯ç‰‡æ•°æ®")

        project = await ppt_service.project_manager.get_project(project_id)
        if not project:
            logger.error(f"âŒ é¡¹ç›® {project_id} ä¸å­˜åœ¨")
            raise HTTPException(status_code=404, detail="Project not found")

        logger.info(f"ğŸ“ æ›´æ–°é¡¹ç›®å¹»ç¯ç‰‡æ•°æ®...")

        # Update project slides data
        project.slides_data = slides_data
        project.updated_at = time.time()

        # Regenerate combined HTML
        if slides_data:
            # å®‰å…¨åœ°è·å–å¤§çº²æ ‡é¢˜
            outline_title = project.title
            if project.outline:
                if isinstance(project.outline, dict):
                    outline_title = project.outline.get('title', project.title)
                elif hasattr(project.outline, 'title'):
                    outline_title = project.outline.title

            project.slides_html = ppt_service._combine_slides_to_full_html(
                slides_data, outline_title
            )

        # æ ‡è®°æ‰€æœ‰å¹»ç¯ç‰‡ä¸ºç”¨æˆ·ç¼–è¾‘çŠ¶æ€
        for i, slide_data in enumerate(project.slides_data):
            slide_data["is_user_edited"] = True

        # ä¿å­˜æ›´æ–°çš„å¹»ç¯ç‰‡æ•°æ®åˆ°æ•°æ®åº“
        save_success = False
        save_error_message = None

        try:
            from ..services.db_project_manager import DatabaseProjectManager
            db_manager = DatabaseProjectManager()

            # ä¿å­˜å¹»ç¯ç‰‡æ•°æ®åˆ°æ•°æ®åº“
            save_success = await db_manager.save_project_slides(
                project_id,
                project.slides_html or "",
                project.slides_data
            )

            if save_success:
                logger.info(f"Successfully saved updated slides data to database for project {project_id}")
            else:
                logger.error(f"Failed to save updated slides data to database for project {project_id}")
                save_error_message = "Failed to save slides data to database"

        except Exception as save_error:
            logger.error(f"âŒ ä¿å­˜å¹»ç¯ç‰‡æ•°æ®åˆ°æ•°æ®åº“æ—¶å‘ç”Ÿå¼‚å¸¸: {save_error}")
            import traceback
            traceback.print_exc()
            save_success = False
            save_error_message = str(save_error)

        # æ ¹æ®ä¿å­˜ç»“æœè¿”å›ç›¸åº”çš„å“åº”
        if save_success:
            return {
                "status": "success",
                "success": True,
                "message": "Slides updated and saved to database successfully"
            }
        else:
            # å³ä½¿æ•°æ®åº“ä¿å­˜å¤±è´¥,å†…å­˜ä¸­çš„æ•°æ®å·²ç»æ›´æ–°,æ‰€ä»¥ä»ç„¶è¿”å›æˆåŠŸ,ä½†åŒ…å«è­¦å‘Šä¿¡æ¯
            return {
                "status": "success",
                "success": True,
                "message": "Slides updated in memory successfully",
                "warning": f"Database save failed: {save_error_message}",
                "database_saved": False
            }

    except Exception as e:
        logger.error(f"âŒ æ›´æ–°é¡¹ç›®å¹»ç¯ç‰‡æ•°æ®æ—¶å‘ç”Ÿé”™è¯¯: {e}")
        import traceback
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=str(e))

@router.post("/api/projects/{project_id}/regenerate-html")
async def regenerate_project_html(project_id: str):
    """Regenerate project HTML with fixed encoding"""
    try:
        project = await ppt_service.project_manager.get_project(project_id)
        if not project:
            raise HTTPException(status_code=404, detail="Project not found")

        if not project.slides_data:
            raise HTTPException(status_code=400, detail="No slides data found")

        # Regenerate combined HTML using the fixed method
        # å®‰å…¨åœ°è·å–å¤§çº²æ ‡é¢˜
        outline_title = project.title
        if project.outline:
            if isinstance(project.outline, dict):
                outline_title = project.outline.get('title', project.title)
            elif hasattr(project.outline, 'title'):
                outline_title = project.outline.title

        project.slides_html = ppt_service._combine_slides_to_full_html(
            project.slides_data, outline_title
        )

        project.updated_at = time.time()

        # ä¿å­˜é‡æ–°ç”Ÿæˆçš„HTMLåˆ°æ•°æ®åº“
        try:
            from ..services.db_project_manager import DatabaseProjectManager
            db_manager = DatabaseProjectManager()

            # ä¿å­˜å¹»ç¯ç‰‡æ•°æ®åˆ°æ•°æ®åº“
            save_success = await db_manager.save_project_slides(
                project_id,
                project.slides_html,
                project.slides_data
            )

            if save_success:
                logger.info(f"Successfully saved regenerated HTML to database for project {project_id}")
            else:
                logger.error(f"Failed to save regenerated HTML to database for project {project_id}")

        except Exception as save_error:
            logger.error(f"Exception while saving regenerated HTML to database: {save_error}")
            # ç»§ç»­è¿”å›æˆåŠŸ,å› ä¸ºå†…å­˜ä¸­çš„æ•°æ®å·²ç»æ›´æ–°

        return {
            "success": True,
            "message": "Project HTML regenerated successfully"
        }

    except Exception as e:
        return {"success": False, "error": str(e)}

@router.post("/api/projects/{project_id}/slides/{slide_number}/regenerate")
async def regenerate_slide(project_id: str, slide_number: int):
    """Regenerate a specific slide"""
    try:
        project = await ppt_service.project_manager.get_project(project_id)
        if not project:
            raise HTTPException(status_code=404, detail="Project not found")

        if not project.outline:
            raise HTTPException(status_code=400, detail="Project outline not found")

        if not project.confirmed_requirements:
            raise HTTPException(status_code=400, detail="Project requirements not confirmed")

        # Handle different outline structures
        if isinstance(project.outline, dict):
            slides = project.outline.get('slides', [])
        else:
            # If outline is a PPTOutline object
            slides = project.outline.slides if hasattr(project.outline, 'slides') else []

        if slide_number < 1 or slide_number > len(slides):
            raise HTTPException(status_code=400, detail="Invalid slide number")

        slide_data = slides[slide_number - 1]

        # Load system prompt
        system_prompt = ppt_service._load_prompts_md_system_prompt()

        # Ensure project has a global template selected (use default if none selected)
        selected_template = await ppt_service._ensure_global_master_template_selected(project_id)

        # Regenerate the slide using template-based generation if template is available
        if selected_template:
            logger.info(f"Regenerating slide {slide_number} using template: {selected_template['template_name']}")
            new_html_content = await ppt_service._generate_slide_with_template(
                slide_data, selected_template, slide_number, len(slides), project.confirmed_requirements
            )
        else:
            # Fallback to original generation method if no template available
            logger.warning(f"No template available for project {project_id}, using fallback generation")
            new_html_content = await ppt_service._generate_single_slide_html_with_prompts(
                slide_data, project.confirmed_requirements, system_prompt, slide_number, len(slides),
                slides, project.slides_data, project_id=project_id
            )

        # Update the slide in project data
        if not project.slides_data:
            project.slides_data = []

        # Ensure slides_data has enough entries
        while len(project.slides_data) < slide_number:
            new_page_number = len(project.slides_data) + 1
            project.slides_data.append({
                "page_number": new_page_number,
                "title": f"ç¬¬{new_page_number}é¡µ",
                "html_content": "<div>å¾…ç”Ÿæˆ</div>",
                "slide_type": "content",
                "content_points": [],
                "is_user_edited": False
            })

        # Update the specific slide - ä¿æŒä¸ç°æœ‰æ•°æ®ç»“æ„ä¸€è‡´
        existing_slide = project.slides_data[slide_number - 1] if slide_number <= len(project.slides_data) else {}

        # æ›´æ–°å¹»ç¯ç‰‡æ•°æ®,ä¿ç•™ç°æœ‰å­—æ®µå¹¶ç¡®ä¿å¿…è¦å­—æ®µå­˜åœ¨
        updated_slide = {
            "page_number": slide_number,
            "title": slide_data.get('title', f'ç¬¬{slide_number}é¡µ'),
            "html_content": new_html_content,
            "slide_type": slide_data.get('slide_type', existing_slide.get('slide_type', 'content')),
            "content_points": slide_data.get('content_points', existing_slide.get('content_points', [])),
            "is_user_edited": existing_slide.get('is_user_edited', False),
            # ä¿ç•™å…¶ä»–å¯èƒ½å­˜åœ¨çš„å­—æ®µ
            **{k: v for k, v in existing_slide.items() if k not in ['page_number', 'title', 'html_content', 'slide_type', 'content_points', 'is_user_edited']}
        }

        project.slides_data[slide_number - 1] = updated_slide

        # Regenerate combined HTML
        outline_title = project.title
        if isinstance(project.outline, dict):
            outline_title = project.outline.get('title', project.title)
        elif hasattr(project.outline, 'title'):
            outline_title = project.outline.title

        project.slides_html = ppt_service._combine_slides_to_full_html(
            project.slides_data, outline_title
        )

        project.updated_at = time.time()

        # ä¿å­˜æ›´æ–°åçš„å¹»ç¯ç‰‡æ•°æ®åˆ°æ•°æ®åº“
        try:
            from ..services.db_project_manager import DatabaseProjectManager
            db_manager = DatabaseProjectManager()

            # åªä¿å­˜å•ä¸ªé‡æ–°ç”Ÿæˆçš„å¹»ç¯ç‰‡,è€Œä¸æ˜¯æ•´ä¸ªé¡¹ç›®çš„å¹»ç¯ç‰‡æ•°æ®
            # è¿™æ ·å¯ä»¥é¿å…åˆ é™¤æ‰€æœ‰å¹»ç¯ç‰‡å†é‡æ–°åˆ›å»ºçš„é—®é¢˜
            save_success = await db_manager.save_single_slide(
                project_id,
                slide_number - 1,  # è½¬æ¢ä¸º0åŸºç´¢å¼•
                updated_slide
            )

            if save_success:
                logger.info(f"Successfully saved regenerated slide {slide_number} to database for project {project_id}")

                # åŒæ—¶æ›´æ–°é¡¹ç›®çš„slides_htmlå­—æ®µ
                await db_manager.update_project_data(project_id, {
                    "slides_html": project.slides_html,
                    "updated_at": project.updated_at
                })
            else:
                logger.error(f"Failed to save regenerated slide {slide_number} to database for project {project_id}")

        except Exception as save_error:
            logger.error(f"Exception while saving regenerated slide to database: {save_error}")
            # ç»§ç»­è¿”å›æˆåŠŸ,å› ä¸ºå†…å­˜ä¸­çš„æ•°æ®å·²ç»æ›´æ–°

        return {
            "success": True,
            "message": f"Slide {slide_number} regenerated successfully",
            "slide_data": project.slides_data[slide_number - 1]
        }

    except Exception as e:
        return {"success": False, "error": str(e)}

@router.post("/api/projects/{project_id}/slides/{slide_index}/auto-repair-layout")
async def auto_repair_layout(
    project_id: str,
    slide_index: int,
    request: AutoLayoutRepairRequest
):
    """Run multimodal layout inspection and repair workflow for a single slide."""
    try:
        if slide_index < 1:
            raise HTTPException(status_code=400, detail="Slide index must be >= 1")

        html_content = (request.html_content or "").strip()
        if not html_content:
            raise HTTPException(status_code=400, detail="HTML content is required")

        project = await ppt_service.project_manager.get_project(project_id)
        if not project:
            raise HTTPException(status_code=404, detail="Project not found")

        slides_data = project.slides_data or []
        total_pages = len(slides_data)
        if total_pages == 0:
            total_pages = request.slide_data.get("total_pages") or request.slide_data.get("totalSlides") or slide_index

        slide_payload = dict(request.slide_data or {})
        slide_payload.setdefault("page_number", slide_index)
        slide_payload.setdefault("title", slide_payload.get("title", f"ç¬¬{slide_index}é¡µ"))

        repaired_html = await ppt_service._apply_auto_layout_repair(
            html_content,
            slide_payload,
            slide_index,
            total_pages or slide_index
        )

        changed = repaired_html.strip() != html_content

        if project.slides_data is None:
            project.slides_data = []

        while len(project.slides_data) < slide_index:
            page_number = len(project.slides_data) + 1
            project.slides_data.append({
                "page_number": page_number,
                "title": f"ç¬¬{page_number}é¡µ",
                "html_content": "",
                "slide_type": "content",
                "content_points": [],
                "is_user_edited": False
            })

        existing_slide = project.slides_data[slide_index - 1]
        updated_slide = {
            **existing_slide,
            "page_number": slide_index,
            "title": slide_payload.get("title", existing_slide.get("title", f"ç¬¬{slide_index}é¡µ")),
            "html_content": repaired_html,
        }

        project.slides_data[slide_index - 1] = updated_slide

        if changed:
            outline_title = project.title
            if isinstance(project.outline, dict):
                outline_title = project.outline.get('title', project.title)
            elif hasattr(project.outline, 'title'):
                outline_title = project.outline.title

            project.slides_html = ppt_service._combine_slides_to_full_html(
                project.slides_data,
                outline_title
            )
            project.updated_at = time.time()

        try:
            from ..services.db_project_manager import DatabaseProjectManager
            db_manager = DatabaseProjectManager()
            await db_manager.save_single_slide(project_id, slide_index - 1, updated_slide)

            if changed:
                await db_manager.update_project_data(project_id, {
                    "slides_html": project.slides_html,
                    "updated_at": project.updated_at
                })

        except Exception as save_error:
            logger.error(f"Failed to persist auto layout repair result: {save_error}")

        return {
            "success": True,
            "repaired_html": repaired_html,
            "changed": changed
        }

    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Auto layout repair failed: {e}", exc_info=True)
        raise HTTPException(status_code=500, detail=str(e))

@router.post("/api/ai/slide-edit")
async def ai_slide_edit(
    request: AISlideEditRequest
):
    """AIç¼–è¾‘å¹»ç¯ç‰‡æ¥å£"""
    try:
        # è·å–AIæä¾›è€…
        provider, settings = get_role_provider("editor")

        # æ„å»ºAIç¼–è¾‘ä¸Šä¸‹æ–‡
        outline_info = ""
        if request.slideOutline:
            outline_info = f"""
å½“å‰å¹»ç¯ç‰‡å¤§çº²ä¿¡æ¯:
- å¹»ç¯ç‰‡ç±»å‹:{request.slideOutline.get('slide_type', 'æœªçŸ¥')}
- æè¿°:{request.slideOutline.get('description', 'æ— ')}
- è¦ç‚¹:{', '.join(request.slideOutline.get('content_points', [])) if request.slideOutline.get('content_points') else 'æ— '}
"""

        context = f"""
ä½ æ˜¯ä¸€ä½ä¸“ä¸šçš„PPTè®¾è®¡å¸ˆå’Œç¼–è¾‘åŠ©æ‰‹.ç”¨æˆ·æƒ³è¦å¯¹å½“å‰å¹»ç¯ç‰‡è¿›è¡Œç¼–è¾‘ä¿®æ”¹.

å½“å‰å¹»ç¯ç‰‡ä¿¡æ¯:
- é¡µç :ç¬¬{request.slideIndex}é¡µ
- æ ‡é¢˜:{request.slideTitle}
- é¡¹ç›®ä¸»é¢˜:{request.projectInfo.get('title', 'æœªçŸ¥')}
- é¡¹ç›®åœºæ™¯:{request.projectInfo.get('scenario', 'æœªçŸ¥')}
{outline_info}
ç”¨æˆ·çš„ç¼–è¾‘è¦æ±‚:
{request.userRequest}

å½“å‰å¹»ç¯ç‰‡çš„HTMLå†…å®¹:
{request.slideContent}

è¯·æ ¹æ®ç”¨æˆ·çš„è¦æ±‚å’Œå¹»ç¯ç‰‡å¤§çº²ä¿¡æ¯,æä¾›ä»¥ä¸‹å†…å®¹:
1. å¯¹ç”¨æˆ·è¦æ±‚çš„ç†è§£å’Œåˆ†æ
2. å…·ä½“çš„ä¿®æ”¹å»ºè®®
3. å¦‚æœéœ€è¦,æä¾›ä¿®æ”¹åçš„å®Œæ•´HTMLä»£ç 

æ³¨æ„äº‹é¡¹:
- ç¡®ä¿ä¿®æ”¹åçš„å†…å®¹ç¬¦åˆPPTæ¼”ç¤ºçš„ä¸“ä¸šæ ‡å‡†å’Œå¤§çº²è¦æ±‚
- ç”Ÿæˆçš„HTMLåº”è¯¥æ˜¯å®Œæ•´çš„,åŒ…å«å¿…è¦çš„CSSæ ·å¼
- ä¿æŒ1280x720çš„PPTæ ‡å‡†å°ºå¯¸
- å‚è€ƒå¤§çº²ä¿¡æ¯ä¸­çš„è¦ç‚¹å’Œæè¿°æ¥ä¼˜åŒ–å†…å®¹
"""

        # æ„å»ºAIæ¶ˆæ¯,åŒ…å«å¯¹è¯å†å²
        messages = [
            AIMessage(role=MessageRole.SYSTEM, content="ä½ æ˜¯ä¸€ä½ä¸“ä¸šçš„PPTè®¾è®¡å¸ˆå’Œç¼–è¾‘åŠ©æ‰‹,æ“…é•¿æ ¹æ®ç”¨æˆ·éœ€æ±‚ä¿®æ”¹å’Œä¼˜åŒ–PPTå†…å®¹.")
        ]

        # æ·»åŠ å¯¹è¯å†å²
        if request.chatHistory:
            logger.debug(f"AIç¼–è¾‘æ¥æ”¶åˆ°å¯¹è¯å†å²,å…± {len(request.chatHistory)} æ¡æ¶ˆæ¯")
            for i, chat_msg in enumerate(request.chatHistory):
                role = MessageRole.USER if chat_msg.get('role') == 'user' else MessageRole.ASSISTANT
                content = chat_msg.get('content', '')
                logger.debug(f"å¯¹è¯å†å² {i+1}: {role.value} - {content[:100]}...")
                messages.append(AIMessage(role=role, content=content))
        else:
            logger.debug("AIç¼–è¾‘æœªæ¥æ”¶åˆ°å¯¹è¯å†å²")

        # æ·»åŠ å½“å‰ç”¨æˆ·è¯·æ±‚
        messages.append(AIMessage(role=MessageRole.USER, content=context))

        # è°ƒç”¨AIç”Ÿæˆå›å¤
        response = await provider.chat_completion(
            messages=messages,
            max_tokens=ai_config.max_tokens,
            temperature=0.7,
            model=settings.get('model')
        )

        ai_response = response.content

        # æ£€æŸ¥æ˜¯å¦åŒ…å«HTMLä»£ç 
        new_html_content = None
        if "```html" in ai_response:
            import re
            html_match = re.search(r'```html\s*(.*?)\s*```', ai_response, re.DOTALL)
            if html_match:
                new_html_content = html_match.group(1).strip()

        return {
            "success": True,
            "response": ai_response,
            "newHtmlContent": new_html_content
        }

    except Exception as e:
        logger.error(f"AIç¼–è¾‘è¯·æ±‚å¤±è´¥: {e}")
        return {
            "success": False,
            "error": str(e),
            "response": "æŠ±æ­‰,AIç¼–è¾‘æœåŠ¡æš‚æ—¶ä¸å¯ç”¨.è¯·ç¨åé‡è¯•."
        }

@router.post("/api/ai/slide-edit/stream")
async def ai_slide_edit_stream(
    request: AISlideEditRequest
):
    """AIç¼–è¾‘å¹»ç¯ç‰‡æµå¼æ¥å£"""
    try:
        # è·å–AIæä¾›è€…
        provider, settings = get_role_provider("editor")

        # æ„å»ºAIç¼–è¾‘ä¸Šä¸‹æ–‡
        outline_info = ""
        if request.slideOutline:
            outline_info = f"""
å½“å‰å¹»ç¯ç‰‡å¤§çº²ä¿¡æ¯:
{request.slideOutline}
"""

        # æ„å»ºå›¾ç‰‡ä¿¡æ¯
        images_info = ""
        if request.images and len(request.images) > 0:
            images_info = f"""

ç”¨æˆ·ä¸Šä¼ çš„å›¾ç‰‡ä¿¡æ¯:
"""
            for i, image in enumerate(request.images, 1):
                images_info += f"""
- å›¾ç‰‡{i}:{image.get('name', 'æœªçŸ¥')}
  - URL:{image.get('url', '')}
  - å¤§å°:{image.get('size', 'æœªçŸ¥')}
  - è¯´æ˜:è¯·åˆ†æè¿™å¼ å›¾ç‰‡çš„å†…å®¹,ç†è§£ç”¨æˆ·çš„æ„å›¾,å¹¶æ ¹æ®ç¼–è¾‘è¦æ±‚è¿›è¡Œç›¸åº”çš„å¤„ç†
"""

        # æ„å»ºè§†è§‰ä¸Šä¸‹æ–‡ä¿¡æ¯
        vision_context = ""
        if request.visionEnabled and request.slideScreenshot:
            vision_context = f"""

ğŸ” è§†è§‰ä¸Šä¸‹æ–‡:
- å½“å‰å¹»ç¯ç‰‡çš„è§†è§‰æˆªå›¾å·²æä¾›
- è¯·ç»“åˆæˆªå›¾ä¸­çš„è§†è§‰å†…å®¹æ¥ç†è§£ç”¨æˆ·çš„ç¼–è¾‘éœ€æ±‚
- æ³¨æ„æˆªå›¾ä¸­çš„å¸ƒå±€,é¢œè‰²,å­—ä½“,å›¾ç‰‡ä½ç½®ç­‰è§†è§‰å…ƒç´ 
- åœ¨æä¾›ç¼–è¾‘å»ºè®®æ—¶,è¯·è€ƒè™‘å½“å‰çš„è§†è§‰å‘ˆç°æ•ˆæœ
"""

        context = f"""
ä½ æ˜¯ä¸€ä½ä¸“ä¸šçš„PPTè®¾è®¡å¸ˆå’Œç¼–è¾‘åŠ©æ‰‹.ç”¨æˆ·æƒ³è¦å¯¹å½“å‰å¹»ç¯ç‰‡è¿›è¡Œç¼–è¾‘ä¿®æ”¹.

å½“å‰å¹»ç¯ç‰‡ä¿¡æ¯:
- é¡µç :ç¬¬{request.slideIndex}é¡µ
- æ ‡é¢˜:{request.slideTitle}
- é¡¹ç›®ä¸»é¢˜:{request.projectInfo.get('title', 'æœªçŸ¥')}
- é¡¹ç›®åœºæ™¯:{request.projectInfo.get('scenario', 'æœªçŸ¥')}
{outline_info}{images_info}{vision_context}
ç”¨æˆ·çš„ç¼–è¾‘è¦æ±‚:
{request.userRequest}

å½“å‰å¹»ç¯ç‰‡çš„HTMLå†…å®¹:
{request.slideContent}

è¯·æ ¹æ®ç”¨æˆ·çš„è¦æ±‚å’Œå¹»ç¯ç‰‡å¤§çº²ä¿¡æ¯,æä¾›ä»¥ä¸‹å†…å®¹:
1. å¯¹ç”¨æˆ·è¦æ±‚çš„ç†è§£å’Œåˆ†æ
2. å…·ä½“çš„ä¿®æ”¹å»ºè®®
3. é»˜è®¤æä¾›ä¿®æ”¹åçš„å®Œæ•´HTMLä»£ç 

æ³¨æ„äº‹é¡¹:
- ä¿æŒåŸæœ‰çš„è®¾è®¡é£æ ¼å’Œå¸ƒå±€ç»“æ„
- ç¡®ä¿ä¿®æ”¹åçš„å†…å®¹ç¬¦åˆPPTæ¼”ç¤ºçš„ä¸“ä¸šæ ‡å‡†å’Œå¤§çº²è¦æ±‚
- å¦‚æœç”¨æˆ·è¦æ±‚ä¸æ˜ç¡®,è¯·æä¾›å¤šä¸ªå¯é€‰æ–¹æ¡ˆ
- ç”Ÿæˆçš„HTMLåº”è¯¥æ˜¯å®Œæ•´çš„,åŒ…å«å¿…è¦çš„CSSæ ·å¼
- ä¿æŒ1280x720çš„PPTæ ‡å‡†å°ºå¯¸
- å‚è€ƒå¤§çº²ä¿¡æ¯ä¸­çš„è¦ç‚¹å’Œæè¿°æ¥ä¼˜åŒ–å†…å®¹
"""

        # æ„å»ºAIæ¶ˆæ¯,åŒ…å«å¯¹è¯å†å²
        messages = [
            AIMessage(role=MessageRole.SYSTEM, content="ä½ æ˜¯ä¸€ä½ä¸“ä¸šçš„PPTè®¾è®¡å¸ˆå’Œç¼–è¾‘åŠ©æ‰‹,æ“…é•¿æ ¹æ®ç”¨æˆ·éœ€æ±‚ä¿®æ”¹å’Œä¼˜åŒ–PPTå†…å®¹.")
        ]

        # æ·»åŠ å¯¹è¯å†å²
        if request.chatHistory:
            logger.info(f"AIæµå¼ç¼–è¾‘æ¥æ”¶åˆ°å¯¹è¯å†å²,å…± {len(request.chatHistory)} æ¡æ¶ˆæ¯")
            for i, chat_msg in enumerate(request.chatHistory):
                role = MessageRole.USER if chat_msg.get('role') == 'user' else MessageRole.ASSISTANT
                content = chat_msg.get('content', '')
                logger.info(f"å¯¹è¯å†å² {i+1}: {role.value} - {content[:100]}...")
                messages.append(AIMessage(role=role, content=content))
        else:
            logger.info("AIæµå¼ç¼–è¾‘æœªæ¥æ”¶åˆ°å¯¹è¯å†å²")

        # æ·»åŠ å½“å‰ç”¨æˆ·è¯·æ±‚(æ”¯æŒå¤šæ¨¡æ€å†…å®¹)
        if request.visionEnabled and request.slideScreenshot:
            # åˆ›å»ºå¤šæ¨¡æ€æ¶ˆæ¯,åŒ…å«æ–‡æœ¬å’Œå›¾ç‰‡
            from ..ai.base import TextContent, ImageContent
            user_content = [
                TextContent(text=context),
                ImageContent(image_url={"url": request.slideScreenshot})
            ]
            messages.append(AIMessage(role=MessageRole.USER, content=user_content))
        else:
            # æ™®é€šæ–‡æœ¬æ¶ˆæ¯
            messages.append(AIMessage(role=MessageRole.USER, content=context))

        async def generate_ai_stream():
            try:
                # å‘é€å¼€å§‹ä¿¡å·
                yield f"data: {json.dumps({'type': 'start', 'content': ''})}\n\n"

                # æµå¼ç”ŸæˆAIå›å¤
                full_response = ""
                if hasattr(provider, 'stream_chat_completion'):
                    async for chunk in provider.stream_chat_completion(
                        messages=messages,
                        max_tokens=ai_config.max_tokens,
                        temperature=0.7,
                        model=settings.get('model')
                    ):
                        if chunk:
                            full_response += chunk
                            yield f"data: {json.dumps({'type': 'content', 'content': chunk})}\n\n"
                else:
                    response = await provider.chat_completion(
                        messages=messages,
                        max_tokens=ai_config.max_tokens,
                        temperature=0.7,
                        model=settings.get('model')
                    )
                    if response.content:
                        full_response = response.content
                        yield f"data: {json.dumps({'type': 'content', 'content': response.content})}\n\n"

                # æ£€æŸ¥æ˜¯å¦åŒ…å«HTMLä»£ç  - æ”¹è¿›ç‰ˆæœ¬,æ”¯æŒå¤šç§æ ¼å¼
                new_html_content = None
                import re
                
                # å°è¯•å¤šç§HTMLä»£ç å—æ ¼å¼
                html_patterns = [
                    r'```html\s*(.*?)\s*```',  # æ ‡å‡†æ ¼å¼
                    r'```HTML\s*(.*?)\s*```',  # å¤§å†™
                    r'```\s*html\s*(.*?)\s*```',  # å¸¦ç©ºæ ¼
                    r'<html[^>]*>.*?</html>',  # å®Œæ•´HTMLæ–‡æ¡£
                    r'<div[^>]*style[^>]*>.*?</div>',  # PPTå¹»ç¯ç‰‡div
                ]
                
                for pattern in html_patterns:
                    html_match = re.search(pattern, full_response, re.DOTALL | re.IGNORECASE)
                    if html_match:
                        new_html_content = html_match.group(1).strip() if html_match.groups() else html_match.group(0).strip()
                        logger.info(f"HTMLå†…å®¹æå–æˆåŠŸ,ä½¿ç”¨æ¨¡å¼: {pattern},å†…å®¹é•¿åº¦: {len(new_html_content)}")
                        break
                
                if not new_html_content:
                    logger.warning(f"æœªèƒ½ä»AIå“åº”ä¸­æå–HTMLå†…å®¹.å“åº”é•¿åº¦: {len(full_response)}")
                    logger.debug(f"AIå®Œæ•´å“åº”: {full_response[:500]}...")

                # å‘é€å®Œæˆä¿¡å·
                yield f"data: {json.dumps({'type': 'complete', 'content': '', 'newHtmlContent': new_html_content, 'fullResponse': full_response})}\n\n"

            except Exception as e:
                logger.error(f"AIæµå¼ç¼–è¾‘è¯·æ±‚å¤±è´¥: {e}")
                yield f"data: {json.dumps({'type': 'error', 'content': '', 'error': str(e)})}\n\n"

        return StreamingResponse(
            generate_ai_stream(),
            media_type="text/event-stream",
            headers={
                "Cache-Control": "no-cache",
                "Connection": "keep-alive",
                "Access-Control-Allow-Origin": "*",
                "Access-Control-Allow-Headers": "Cache-Control"
            }
        )

    except Exception as e:
        logger.error(f"AIæµå¼ç¼–è¾‘è¯·æ±‚å¤±è´¥: {e}")
        return {
            "success": False,
            "error": str(e),
            "response": "æŠ±æ­‰,AIç¼–è¾‘æœåŠ¡æš‚æ—¶ä¸å¯ç”¨.è¯·ç¨åé‡è¯•."
        }

# å¤§çº²AIä¼˜åŒ–è¯·æ±‚æ•°æ®æ¨¡å‹
class OutlineAIOptimizeRequest(BaseModel):
    outline_content: str  # JSONæ ¼å¼çš„å¤§çº²å†…å®¹
    user_request: str  # ç”¨æˆ·çš„ä¼˜åŒ–éœ€æ±‚
    project_info: Dict[str, Any]  # é¡¹ç›®ä¿¡æ¯
    optimization_type: str = "full"  # full=å…¨å¤§çº²ä¼˜åŒ–, single=å•é¡µä¼˜åŒ–
    slide_index: Optional[int] = None  # å½“optimization_type=singleæ—¶ä½¿ç”¨

@router.post("/api/ai/optimize-outline")
async def ai_optimize_outline(
    request: OutlineAIOptimizeRequest
):
    """AIä¼˜åŒ–å¤§çº²æ¥å£ - æ”¯æŒå…¨å¤§çº²ä¼˜åŒ–å’Œå•é¡µä¼˜åŒ–"""
    try:
        # è·å–AIæä¾›è€…
        provider, settings = get_role_provider("editor")
        
        # è§£æå¤§çº²JSON
        try:
            outline_data = json.loads(request.outline_content)
        except json.JSONDecodeError as e:
            return {
                "success": False,
                "error": f"å¤§çº²JSONæ ¼å¼é”™è¯¯: {str(e)}"
            }
        
        # æ ¹æ®ä¼˜åŒ–ç±»å‹æ„å»ºä¸åŒçš„æç¤ºè¯
        if request.optimization_type == "single" and request.slide_index is not None:
            # å•é¡µä¼˜åŒ–
            if request.slide_index < 0 or request.slide_index >= len(outline_data.get('slides', [])):
                return {
                    "success": False,
                    "error": "æ— æ•ˆçš„å¹»ç¯ç‰‡ç´¢å¼•"
                }
            
            slide = outline_data['slides'][request.slide_index]
            
            context = f"""
ä½ æ˜¯ä¸€ä½ä¸“ä¸šçš„PPTå¤§çº²è®¾è®¡ä¸“å®¶.ç”¨æˆ·æƒ³è¦ä¼˜åŒ–PPTå¤§çº²ä¸­çš„ç¬¬{request.slide_index + 1}é¡µå†…å®¹.

é¡¹ç›®ä¿¡æ¯:
- ä¸»é¢˜:{request.project_info.get('topic', 'æœªçŸ¥')}
- åœºæ™¯:{request.project_info.get('scenario', 'é€šç”¨')}
- ç›®æ ‡å—ä¼—:{request.project_info.get('target_audience', 'æ™®é€šå¤§ä¼—')}

å½“å‰é¡µé¢ä¿¡æ¯:
- é¡µç :ç¬¬{slide.get('page_number', request.slide_index + 1)}é¡µ
- æ ‡é¢˜:{slide.get('title', 'æœªå‘½å')}
- ç±»å‹:{slide.get('slide_type', 'content')}
- å†…å®¹è¦ç‚¹:{json.dumps(slide.get('content_points', []), ensure_ascii=False, indent=2)}

ç”¨æˆ·çš„ä¼˜åŒ–éœ€æ±‚:
{request.user_request}

è¯·æ ¹æ®ç”¨æˆ·éœ€æ±‚ä¼˜åŒ–è¿™ä¸€é¡µçš„å†…å®¹.

ã€é‡è¦ã€‘ç›´æ¥è¿”å›ä¼˜åŒ–åçš„JSONæ•°æ®,ä¸è¦åŒ…å«ä»»ä½•è§£é‡Šæ€§æ–‡å­—æˆ–markdownæ ‡è®°(å¦‚```json).

è¿”å›æ ¼å¼ç¤ºä¾‹:
{{
  "page_number": {slide.get('page_number', request.slide_index + 1)},
  "title": "ä¼˜åŒ–åçš„æ ‡é¢˜",
  "subtitle": "å‰¯æ ‡é¢˜(å¯é€‰)",
  "content_points": ["è¦ç‚¹1", "è¦ç‚¹2", "è¦ç‚¹3"],
  "slide_type": "content",
  "description": "é¡µé¢æè¿°(å¯é€‰)"
}}

ä¼˜åŒ–è¦æ±‚:
1. ä¿æŒä¸æ•´ä½“å¤§çº²çš„è¿è´¯æ€§å’Œé€»è¾‘æ€§
2. ç¡®ä¿å†…å®¹è¦ç‚¹æ¸…æ™°,å…·ä½“,æœ‰ä»·å€¼
3. æ ‡é¢˜è¦ç®€æ´æœ‰åŠ›,èƒ½å¤Ÿå‡†ç¡®æ¦‚æ‹¬é¡µé¢å†…å®¹
4. content_pointsæ•°ç»„ä¸­çš„å­—ç¬¦ä¸²å¯ä»¥åŒ…å«ä»£ç ç¤ºä¾‹(ç”¨```æ ‡è®°),è¿™æ˜¯åˆæ³•çš„JSONå­—ç¬¦ä¸²å†…å®¹
5. ã€å…³é”®ã€‘åªè¿”å›çº¯JSONå¯¹è±¡,ä¸è¦ç”¨```jsonåŒ…è£¹æ•´ä¸ªJSON,ä¸è¦æ·»åŠ ä»»ä½•å…¶ä»–è§£é‡Šæ–‡å­—
"""
        else:
            # å…¨å¤§çº²ä¼˜åŒ–
            context = f"""
ä½ æ˜¯ä¸€ä½ä¸“ä¸šçš„PPTå¤§çº²è®¾è®¡ä¸“å®¶.ç”¨æˆ·æƒ³è¦ä¼˜åŒ–æ•´ä¸ªPPTå¤§çº².

é¡¹ç›®ä¿¡æ¯:
- ä¸»é¢˜:{request.project_info.get('topic', 'æœªçŸ¥')}
- åœºæ™¯:{request.project_info.get('scenario', 'é€šç”¨')}
- ç›®æ ‡å—ä¼—:{request.project_info.get('target_audience', 'æ™®é€šå¤§ä¼—')}
- å½“å‰é¡µæ•°:{len(outline_data.get('slides', []))}é¡µ

å½“å‰å¤§çº²:
{json.dumps(outline_data, ensure_ascii=False, indent=2)}

ç”¨æˆ·çš„ä¼˜åŒ–éœ€æ±‚:
{request.user_request}

è¯·æ ¹æ®ç”¨æˆ·éœ€æ±‚ä¼˜åŒ–æ•´ä¸ªå¤§çº².

ã€é‡è¦ã€‘ç›´æ¥è¿”å›å®Œæ•´çš„ä¼˜åŒ–åçš„JSONæ•°æ®,ä¸è¦åŒ…å«ä»»ä½•è§£é‡Šæ€§æ–‡å­—,markdownæ ‡è®°æˆ–æ³¨é‡Š.

è¿”å›æ ¼å¼ç¤ºä¾‹:
{{
  "title": "ä¼˜åŒ–åçš„PPTæ ‡é¢˜",
  "slides": [
    {{
      "page_number": 1,
      "title": "é¡µé¢æ ‡é¢˜",
      "subtitle": "å‰¯æ ‡é¢˜(å¯é€‰)",
      "content_points": ["è¦ç‚¹1", "è¦ç‚¹2"],
      "slide_type": "title",
      "description": "é¡µé¢æè¿°(å¯é€‰)"
    }}
  ],
  "metadata": {{
    "scenario": "{request.project_info.get('scenario', 'é€šç”¨')}",
    "language": "zh",
    "target_audience": "{request.project_info.get('target_audience', 'æ™®é€šå¤§ä¼—')}",
    "optimized": true
  }}
}}

ä¼˜åŒ–è¦æ±‚:
1. ä¿æŒå¤§çº²çš„æ•´ä½“é€»è¾‘æ€§å’Œè¿è´¯æ€§
2. ç¡®ä¿æ¯é¡µå†…å®¹è¦ç‚¹æ¸…æ™°,å…·ä½“,æœ‰ä»·å€¼
3. å¯ä»¥è°ƒæ•´é¡µé¢é¡ºåº,åˆå¹¶æˆ–æ‹†åˆ†é¡µé¢,ä½†è¦ä¿æŒæ€»ä½“ç»“æ„åˆç†
4. æ ‡é¢˜è¦ç®€æ´æœ‰åŠ›
5. ã€å…³é”®ã€‘åªè¿”å›çº¯JSONæ ¼å¼,ä¸è¦æ·»åŠ ä»»ä½•è§£é‡Š,æ³¨é‡Šæˆ–markdownæ ‡è®°
"""
        
        # æ„å»ºAIæ¶ˆæ¯
        messages = [
            AIMessage(role=MessageRole.SYSTEM, content="ä½ æ˜¯ä¸€ä½ä¸“ä¸šçš„PPTå¤§çº²è®¾è®¡ä¸“å®¶,æ“…é•¿ä¼˜åŒ–å’Œæ”¹è¿›PPTå¤§çº²ç»“æ„å’Œå†…å®¹.ä½ çš„å›å¤å¿…é¡»æ˜¯çº¯JSONæ ¼å¼,ä¸è¦åŒ…å«ä»»ä½•è§£é‡Šæ€§æ–‡å­—,markdownæ ‡è®°æˆ–æ³¨é‡Š."),
            AIMessage(role=MessageRole.USER, content=context)
        ]
        
        # è°ƒç”¨AIç”Ÿæˆå›å¤
        response = await provider.chat_completion(
            messages=messages,
            temperature=0.7,
            model=settings.get('model')
        )
        
        ai_response = response.content
        
        # æ™ºèƒ½æå–JSONå†…å®¹
        import re
        
        def extract_json_from_response(text: str) -> str:
            """ä»AIå“åº”ä¸­æå–JSONå†…å®¹,æ”¯æŒå¤šç§æ ¼å¼"""
            
            # ä¼˜å…ˆæ–¹æ³•: æŸ¥æ‰¾ç¬¬ä¸€ä¸ª{åˆ°æœ€åä¸€ä¸ª}ä¹‹é—´çš„å†…å®¹
            # è¿™æ ·å¯ä»¥é¿å…é”™è¯¯æå–content_pointså­—æ®µå†…çš„ä»£ç å—
            first_brace = text.find('{')
            last_brace = text.rfind('}')
            if first_brace != -1 and last_brace != -1 and first_brace < last_brace:
                potential_json = text[first_brace:last_brace + 1]
                # å°è¯•è§£æ,å¦‚æœæˆåŠŸåˆ™è¿”å›
                try:
                    json.loads(potential_json)
                    return potential_json.strip()
                except json.JSONDecodeError:
                    # å¦‚æœè§£æå¤±è´¥,å°è¯•æ¸…ç†æ³¨é‡Šåå†è¯•
                    cleaned_json = re.sub(r'//[^\n]*', '', potential_json)  # å•è¡Œæ³¨é‡Š
                    cleaned_json = re.sub(r'/\*.*?\*/', '', cleaned_json, flags=re.DOTALL)  # å¤šè¡Œæ³¨é‡Š
                    try:
                        json.loads(cleaned_json)
                        return cleaned_json.strip()
                    except json.JSONDecodeError:
                        pass  # ç»§ç»­å°è¯•å…¶ä»–æ–¹æ³•
            
            # å¤‡ç”¨æ–¹æ³•: æå–markdownä»£ç å—ä¸­çš„JSON(ä»…å½“æ ‡è®°ä¸ºjsonæ—¶)
            # ä½¿ç”¨æ›´ä¸¥æ ¼çš„åŒ¹é…,ç¡®ä¿æ˜¯JSONä»£ç å—è€Œä¸æ˜¯å…¶ä»–ä»£ç å—
            json_match = re.search(r'```json\s*(.*?)\s*```', text, re.DOTALL | re.IGNORECASE)
            if json_match:
                extracted = json_match.group(1).strip()
                # éªŒè¯æå–çš„å†…å®¹æ˜¯å¦æ˜¯æœ‰æ•ˆJSON
                try:
                    json.loads(extracted)
                    return extracted
                except json.JSONDecodeError:
                    pass  # ç»§ç»­å°è¯•å…¶ä»–æ–¹æ³•
            
            # æœ€åå°è¯•: ç›´æ¥è¿”å›æ¸…ç†åçš„æ–‡æœ¬
            cleaned = text.strip()
            if cleaned.startswith('{'):
                return cleaned
            
            # å°è¯•æ‰¾åˆ°JSONå¼€å§‹çš„ä½ç½®
            for line in cleaned.split('\n'):
                line = line.strip()
                if line.startswith('{'):
                    start_idx = cleaned.find(line)
                    return cleaned[start_idx:].strip()
            
            return cleaned
        
        optimized_json = extract_json_from_response(ai_response)
        
        # éªŒè¯JSONæ ¼å¼
        try:
            optimized_data = json.loads(optimized_json)
        except json.JSONDecodeError as e:
            # æä¾›æ›´è¯¦ç»†çš„é”™è¯¯ä¿¡æ¯,å¸®åŠ©è°ƒè¯•
            return {
                "success": False,
                "error": f"AIè¿”å›çš„å†…å®¹ä¸æ˜¯æœ‰æ•ˆçš„JSONæ ¼å¼: {str(e)}",
                "raw_response": ai_response,
                "extracted_json": optimized_json[:500] if len(optimized_json) > 500 else optimized_json
            }
        
        return {
            "success": True,
            "optimized_content": json.dumps(optimized_data, ensure_ascii=False, indent=2),
            "optimization_type": request.optimization_type,
            "raw_response": ai_response
        }
        
    except Exception as e:
        logger.error(f"AIä¼˜åŒ–å¤§çº²è¯·æ±‚å¤±è´¥: {e}")
        return {
            "success": False,
            "error": str(e)
        }

@router.post("/api/ai/regenerate-image")
async def ai_regenerate_image(
    request: AIImageRegenerateRequest
):
    """AIé‡æ–°ç”Ÿæˆå›¾åƒæ¥å£ - å®Œå…¨éµå¾ªenhanced_ppt_service.pyçš„æ ‡å‡†æµç¨‹"""
    try:
        # è·å–å›¾åƒæœåŠ¡å’ŒAIæä¾›è€…
        from ..services.image.image_service import get_image_service

        image_service = get_image_service()
        if not image_service:
            return {
                "success": False,
                "message": "å›¾åƒæœåŠ¡ä¸å¯ç”¨"
            }

        provider, settings = get_role_provider("editor")
        # Ensure we have a general AI provider instance as well (some image processors expect ai_provider)
        ai_provider = get_ai_provider()
        if not ai_provider:
            return {
                "success": False,
                "message": "AIæä¾›è€…ä¸å¯ç”¨"
            }

        # è·å–å›¾åƒé…ç½®
        from ..services.config_service import config_service
        image_config = config_service.get_config_by_category('image_service')

        # æ£€æŸ¥æ˜¯å¦å¯ç”¨å›¾ç‰‡ç”ŸæˆæœåŠ¡
        enable_image_service = image_config.get('enable_image_service', False)
        if not enable_image_service:
            return {
                "success": False,
                "message": "å›¾ç‰‡ç”ŸæˆæœåŠ¡æœªå¯ç”¨,è¯·åœ¨é…ç½®ä¸­å¯ç”¨"
            }

        # ç¬¬ä¸€æ­¥:æ£€æŸ¥å¯ç”¨çš„å›¾ç‰‡æ¥æº(å®Œå…¨éµå¾ªPPTImageProcessorçš„é€»è¾‘)
        from ..services.models.slide_image_info import ImageSource

        enabled_sources = []
        if image_config.get('enable_local_images', True):
            enabled_sources.append(ImageSource.LOCAL)
        if image_config.get('enable_network_search', False):
            enabled_sources.append(ImageSource.NETWORK)
        if image_config.get('enable_ai_generation', False):
            enabled_sources.append(ImageSource.AI_GENERATED)

        if not enabled_sources:
            return {
                "success": False,
                "message": "æ²¡æœ‰å¯ç”¨ä»»ä½•å›¾ç‰‡æ¥æº,è¯·åœ¨é…ç½®ä¸­å¯ç”¨è‡³å°‘ä¸€ç§å›¾ç‰‡æ¥æº"
            }

        # åˆå§‹åŒ–PPTå›¾åƒå¤„ç†å™¨
        from ..services.ppt_image_processor import PPTImageProcessor

        image_processor = PPTImageProcessor(
            image_service=image_service,
            ai_provider=ai_provider
        )

        # æå–å›¾åƒä¿¡æ¯å’Œå¹»ç¯ç‰‡å†…å®¹
        image_info = request.image_info
        slide_content = request.slide_content

        # æ„å»ºå¹»ç¯ç‰‡æ•°æ®ç»“æ„(éµå¾ªPPTImageProcessoræœŸæœ›çš„æ ¼å¼)
        slide_data = {
            'title': slide_content.get('title', ''),
            'content_points': [slide_content.get('title', '')],  # ç®€åŒ–çš„å†…å®¹ç‚¹
        }

        # æ„å»ºç¡®è®¤éœ€æ±‚ç»“æ„
        confirmed_requirements = {
            'project_topic': request.project_topic,
            'project_scenario': request.project_scenario
        }

        # ç¬¬äºŒæ­¥:ç›´æ¥åˆ›å»ºå›¾åƒé‡æ–°ç”Ÿæˆéœ€æ±‚(è·³è¿‡AIé…å›¾é€‚ç”¨æ€§åˆ¤æ–­)
        logger.info(f"å¼€å§‹å›¾ç‰‡é‡æ–°ç”Ÿæˆ,å¯ç”¨çš„æ¥æº: {[source.value for source in enabled_sources]}")

        # åˆ†æåŸå›¾åƒçš„ç”¨é€”å’Œä¸Šä¸‹æ–‡
        image_context = await analyze_image_context(
            image_info, slide_content, request.project_topic, request.project_scenario
        )

        # æ ¹æ®å¯ç”¨çš„æ¥æºå’Œé…ç½®,æ™ºèƒ½é€‰æ‹©æœ€ä½³çš„å›¾ç‰‡æ¥æº
        selected_source = select_best_image_source(enabled_sources, image_config, image_context)

        # åˆ›å»ºå›¾åƒéœ€æ±‚å¯¹è±¡(ç›´æ¥ç”Ÿæˆ,ä¸éœ€è¦AIåˆ¤æ–­æ˜¯å¦é€‚åˆé…å›¾)
        from ..services.models.slide_image_info import ImageRequirement, ImagePurpose

        # å°†å­—ç¬¦ä¸²ç”¨é€”è½¬æ¢ä¸ºImagePurposeæšä¸¾
        purpose_str = image_context.get('image_purpose', 'illustration')
        purpose_mapping = {
            'background': ImagePurpose.BACKGROUND,
            'icon': ImagePurpose.ICON,
            'chart_support': ImagePurpose.CHART_SUPPORT,
            'decoration': ImagePurpose.DECORATION,
            'illustration': ImagePurpose.ILLUSTRATION
        }
        purpose = purpose_mapping.get(purpose_str, ImagePurpose.ILLUSTRATION)

        requirement = ImageRequirement(
            source=selected_source,
            count=1,
            purpose=purpose,
            description=f"é‡æ–°ç”Ÿæˆå›¾åƒ: {image_info.get('alt', '')} - {request.project_topic}",
            priority=5  # é«˜ä¼˜å…ˆçº§,å› ä¸ºæ˜¯ç”¨æˆ·æ˜ç¡®è¯·æ±‚çš„é‡æ–°ç”Ÿæˆ
        )

        logger.info(f"é€‰æ‹©å›¾ç‰‡æ¥æº: {selected_source.value}, ç”¨é€”: {purpose.value}")

        # ç¬¬ä¸‰æ­¥:ç›´æ¥å¤„ç†å›¾ç‰‡ç”Ÿæˆ(å•ä¸ªéœ€æ±‚)
        from ..services.models.slide_image_info import SlideImagesCollection

        images_collection = SlideImagesCollection(page_number=request.slide_index + 1, images=[])

        # æ ¹æ®é€‰æ‹©çš„æ¥æºå¤„ç†å›¾ç‰‡ç”Ÿæˆ
        if requirement.source == ImageSource.LOCAL and ImageSource.LOCAL in enabled_sources:
            local_images = await image_processor._process_local_images(
                requirement, request.project_topic, request.project_scenario,
                slide_content.get('title', ''), slide_content.get('title', '')
            )
            images_collection.images.extend(local_images)

        elif requirement.source == ImageSource.NETWORK and ImageSource.NETWORK in enabled_sources:
            network_images = await image_processor._process_network_images(
                requirement, request.project_topic, request.project_scenario,
                slide_content.get('title', ''), slide_content.get('title', ''), image_config
            )
            images_collection.images.extend(network_images)

        elif requirement.source == ImageSource.AI_GENERATED and ImageSource.AI_GENERATED in enabled_sources:
            ai_images = await image_processor._process_ai_generated_images(
                requirement=requirement,
                project_topic=request.project_topic,
                project_scenario=request.project_scenario,
                slide_title=slide_content.get('title', ''),
                slide_content=slide_content.get('title', ''),
                image_config=image_config,
                page_number=request.slide_index + 1,
                total_pages=1,
                template_html=slide_content.get('html_content', '')
            )
            images_collection.images.extend(ai_images)

        # é‡æ–°è®¡ç®—ç»Ÿè®¡ä¿¡æ¯
        images_collection.__post_init__()

        if images_collection.total_count == 0:
            return {
                "success": False,
                "message": "æœªèƒ½ç”Ÿæˆä»»ä½•å›¾ç‰‡,è¯·æ£€æŸ¥é…ç½®å’Œç½‘ç»œè¿æ¥"
            }

        # è·å–ç¬¬ä¸€å¼ ç”Ÿæˆçš„å›¾åƒ(ç”¨äºæ›¿æ¢)
        new_image = images_collection.images[0]
        new_image_url = new_image.absolute_url

        # æ›¿æ¢HTMLä¸­çš„å›¾åƒ
        updated_html = replace_image_in_html(
            slide_content.get('html_content', ''),
            image_info,
            new_image_url
        )

        logger.info(f"å›¾ç‰‡é‡æ–°ç”ŸæˆæˆåŠŸ: {new_image.source.value}æ¥æº, URL: {new_image_url}")

        return {
            "success": True,
            "message": f"å›¾åƒé‡æ–°ç”ŸæˆæˆåŠŸ(æ¥æº:{new_image.source.value})",
            "new_image_url": new_image_url,
            "new_image_id": new_image.image_id,
            "updated_html_content": updated_html,
            "generation_prompt": getattr(new_image, 'generation_prompt', ''),
            "image_source": new_image.source.value,
            "ai_analysis": {
                "total_images_analyzed": 1,
                "reasoning": f"ç”¨æˆ·è¯·æ±‚é‡æ–°ç”Ÿæˆ{image_context.get('image_purpose', 'å›¾åƒ')},é€‰æ‹©{selected_source.value}æ¥æº",
                "enabled_sources": [source.value for source in enabled_sources],
                "selected_source": selected_source.value
            },
            "image_info": {
                "width": new_image.width,
                "height": new_image.height,
                "format": getattr(new_image, 'format', 'unknown'),
                "alt_text": new_image.alt_text,
                "title": new_image.title,
                "source": new_image.source.value,
                "purpose": new_image.purpose.value
            }
        }

    except Exception as e:
        logger.error(f"AIå›¾åƒé‡æ–°ç”Ÿæˆå¤±è´¥: {e}")
        import traceback
        traceback.print_exc()
        return {
            "success": False,
            "message": f"å›¾åƒé‡æ–°ç”Ÿæˆå¤±è´¥: {str(e)}"
        }

@router.post("/api/ai/auto-generate-slide-images")
async def ai_auto_generate_slide_images(
    request: AIAutoImageGenerateRequest
):
    """AIä¸€é”®é…å›¾æ¥å£ - è‡ªåŠ¨åˆ†æå¹»ç¯ç‰‡å†…å®¹å¹¶ç”Ÿæˆç›¸å…³é…å›¾"""
    try:
        # è·å–å›¾åƒæœåŠ¡å’ŒAIæä¾›è€…
        from ..services.image.image_service import get_image_service

        image_service = get_image_service()
        if not image_service:
            return {
                "success": False,
                "message": "å›¾åƒæœåŠ¡ä¸å¯ç”¨"
            }

        ai_provider = get_ai_provider()
        if not ai_provider:
            return {
                "success": False,
                "message": "AIæä¾›è€…ä¸å¯ç”¨"
            }

        # è·å–å›¾åƒå¤„ç†å™¨
        from ..services.ppt_image_processor import PPTImageProcessor
        image_processor = PPTImageProcessor(image_service, ai_provider)

        slide_content = request.slide_content
        slide_title = slide_content.get('title', f'ç¬¬{request.slide_index + 1}é¡µ')
        slide_html = slide_content.get('html_content', '')

        logger.info(f"å¼€å§‹ä¸ºç¬¬{request.slide_index + 1}é¡µè¿›è¡Œä¸€é”®é…å›¾")

        # ç¬¬ä¸€æ­¥:AIåˆ†æå¹»ç¯ç‰‡å†…å®¹,ç¡®å®šæ˜¯å¦éœ€è¦é…å›¾ä»¥åŠé…å›¾éœ€æ±‚
        analysis_prompt = f"""ä½œä¸ºä¸“ä¸šçš„PPTè®¾è®¡å¸ˆ,è¯·åˆ†æä»¥ä¸‹å¹»ç¯ç‰‡å†…å®¹,åˆ¤æ–­æ˜¯å¦éœ€è¦é…å›¾ä»¥åŠé…å›¾éœ€æ±‚.

é¡¹ç›®ä¸»é¢˜:{request.project_topic}
é¡¹ç›®åœºæ™¯:{request.project_scenario}
å¹»ç¯ç‰‡æ ‡é¢˜:{slide_title}
å¹»ç¯ç‰‡HTMLå†…å®¹:{slide_html[:1000]}...

è¯·åˆ†æ:
1. è¿™ä¸ªå¹»ç¯ç‰‡æ˜¯å¦éœ€è¦é…å›¾?
2. å¦‚æœéœ€è¦,åº”è¯¥é…å‡ å¼ å›¾?
3. æ¯å¼ å›¾çš„ç”¨é€”å’Œæè¿°æ˜¯ä»€ä¹ˆ?
4. å›¾ç‰‡åº”è¯¥æ’å…¥åˆ°ä»€ä¹ˆä½ç½®?

è¯·ä»¥JSONæ ¼å¼å›å¤:
{{
    "needs_images": true/false,
    "image_count": æ•°é‡,
    "images": [
        {{
            "purpose": "å›¾ç‰‡ç”¨é€”(å¦‚:ä¸»è¦æ’å›¾,è£…é¥°å›¾,èƒŒæ™¯å›¾ç­‰)",
            "description": "å›¾ç‰‡å†…å®¹æè¿°",
            "keywords": "æœç´¢å…³é”®è¯",
            "position": "æ’å…¥ä½ç½®(å¦‚:æ ‡é¢˜ä¸‹æ–¹,å†…å®¹ä¸­é—´,é¡µé¢å³ä¾§ç­‰)"
        }}
    ],
    "reasoning": "åˆ†æç†ç”±"
}}"""

        analysis_response = await ai_provider.text_completion(
            prompt=analysis_prompt,
            temperature=0.3
        )

        # è§£æAIåˆ†æç»“æœ
        import json
        try:
            analysis_result = json.loads(analysis_response.content.strip())
        except json.JSONDecodeError:
            # å¦‚æœJSONè§£æå¤±è´¥,ä½¿ç”¨é»˜è®¤é…ç½®
            analysis_result = {
                "needs_images": True,
                "image_count": 1,
                "images": [{
                    "purpose": "ä¸»è¦æ’å›¾",
                    "description": f"ä¸{slide_title}ç›¸å…³çš„é…å›¾",
                    "keywords": f"{request.project_topic} {slide_title}",
                    "position": "å†…å®¹ä¸­é—´"
                }],
                "reasoning": "é»˜è®¤ä¸ºå¹»ç¯ç‰‡æ·»åŠ ä¸€å¼ ä¸»è¦é…å›¾"
            }

        if not analysis_result.get("needs_images", False):
            return {
                "success": True,
                "message": "AIåˆ†æè®¤ä¸ºæ­¤å¹»ç¯ç‰‡ä¸éœ€è¦é…å›¾",
                "updated_html_content": slide_html,
                "generated_images_count": 0,
                "ai_analysis": analysis_result
            }

        # ç¬¬äºŒæ­¥:æ ¹æ®åˆ†æç»“æœç”Ÿæˆå›¾ç‰‡éœ€æ±‚
        from ..services.models.slide_image_info import ImageRequirement, ImagePurpose, ImageSource, SlideImagesCollection

        images_collection = SlideImagesCollection(page_number=request.slide_index + 1, images=[])

        # è·å–å›¾åƒé…ç½®(ä½¿ç”¨ä¸é‡æ–°ç”Ÿæˆå›¾ç‰‡ç›¸åŒçš„é…ç½®é”®)
        from ..services.config_service import config_service
        image_config = config_service.get_config_by_category('image_service')

        # æ£€æŸ¥æ˜¯å¦å¯ç”¨å›¾ç‰‡ç”ŸæˆæœåŠ¡
        enable_image_service = image_config.get('enable_image_service', False)
        if not enable_image_service:
            return {
                "success": False,
                "message": "å›¾ç‰‡ç”ŸæˆæœåŠ¡æœªå¯ç”¨,è¯·åœ¨é…ç½®ä¸­å¯ç”¨"
            }

        # è·å–å¯ç”¨çš„å›¾åƒæ¥æº(ä½¿ç”¨ä¸é‡æ–°ç”Ÿæˆå›¾ç‰‡ç›¸åŒçš„é€»è¾‘)
        from ..services.models.slide_image_info import ImageSource

        enabled_sources = []
        if image_config.get('enable_local_images', True):
            enabled_sources.append(ImageSource.LOCAL)
        if image_config.get('enable_network_search', False):
            enabled_sources.append(ImageSource.NETWORK)
        if image_config.get('enable_ai_generation', False):
            enabled_sources.append(ImageSource.AI_GENERATED)

        if not enabled_sources:
            return {
                "success": False,
                "message": "æ²¡æœ‰å¯ç”¨çš„å›¾åƒæ¥æº,è¯·åœ¨è®¾ç½®ä¸­é…ç½®å›¾åƒè·å–æ–¹å¼"
            }

        # ä½¿ç”¨ä¸é‡æ–°ç”Ÿæˆå›¾ç‰‡å®Œå…¨ç›¸åŒçš„å›¾ç‰‡æ¥æºé€‰æ‹©é€»è¾‘
        image_context = {
            'image_purpose': 'illustration',  # ä¸€é”®é…å›¾é»˜è®¤ä¸ºè¯´æ˜æ€§å›¾ç‰‡
            'slide_title': slide_title,
            'slide_content': slide_html
        }

        selected_source = select_best_image_source(enabled_sources, image_config, image_context)

        # ä¸ºæ¯ä¸ªå›¾ç‰‡éœ€æ±‚ç”Ÿæˆå›¾ç‰‡
        for i, image_info in enumerate(analysis_result.get("images", [])[:3]):  # æœ€å¤š3å¼ å›¾
            # åˆ›å»ºå›¾ç‰‡éœ€æ±‚
            requirement = ImageRequirement(
                purpose=ImagePurpose.ILLUSTRATION,
                description=image_info.get("description", "ç›¸å…³é…å›¾"),
                priority=1,
                source=selected_source,
                count=1
            )

            # æ ¹æ®é€‰æ‹©çš„æ¥æºå¤„ç†å›¾ç‰‡ç”Ÿæˆ
            if requirement.source == ImageSource.AI_GENERATED and ImageSource.AI_GENERATED in enabled_sources:
                ai_images = await image_processor._process_ai_generated_images(
                    requirement=requirement,
                    project_topic=request.project_topic,
                    project_scenario=request.project_scenario,
                    slide_title=slide_title,
                    slide_content=slide_title,
                    image_config=image_config,
                    page_number=request.slide_index + 1,
                    total_pages=1,
                    template_html=slide_html
                )
                images_collection.images.extend(ai_images)

            elif requirement.source == ImageSource.NETWORK and ImageSource.NETWORK in enabled_sources:
                network_images = await image_processor._process_network_images(
                    requirement=requirement,
                    project_topic=request.project_topic,
                    project_scenario=request.project_scenario,
                    slide_title=slide_title,
                    slide_content=slide_title,
                    image_config=image_config
                )
                images_collection.images.extend(network_images)

            elif requirement.source == ImageSource.LOCAL and ImageSource.LOCAL in enabled_sources:
                local_images = await image_processor._process_local_images(
                    requirement=requirement,
                    project_topic=request.project_topic,
                    project_scenario=request.project_scenario,
                    slide_title=slide_title,
                    slide_content=slide_title
                )
                images_collection.images.extend(local_images)

        if not images_collection.images:
            return {
                "success": False,
                "message": "æœªèƒ½ç”Ÿæˆä»»ä½•é…å›¾,è¯·æ£€æŸ¥å›¾åƒæœåŠ¡é…ç½®"
            }

        # ç¬¬ä¸‰æ­¥:å°†ç”Ÿæˆçš„å›¾ç‰‡æ’å…¥åˆ°å¹»ç¯ç‰‡ä¸­
        updated_html = await image_processor._insert_images_into_slide(
            slide_html, images_collection, slide_title
        )

        logger.info(f"ä¸€é”®é…å›¾å®Œæˆ: ç”Ÿæˆ{len(images_collection.images)}å¼ å›¾ç‰‡")

        return {
            "success": True,
            "message": f"ä¸€é”®é…å›¾å®Œæˆ,å·²ç”Ÿæˆ{len(images_collection.images)}å¼ å›¾ç‰‡",
            "updated_html_content": updated_html,
            "generated_images_count": len(images_collection.images),
            "generated_images": [
                {
                    "image_id": img.image_id,
                    "url": img.absolute_url,
                    "description": img.content_description,
                    "source": img.source.value
                } for img in images_collection.images
            ],
            "ai_analysis": analysis_result
        }

    except Exception as e:
        logger.error(f"AIä¸€é”®é…å›¾å¤±è´¥: {e}")
        return {
            "success": False,
            "message": f"ä¸€é”®é…å›¾å¤±è´¥: {str(e)}"
        }

@router.post("/api/ai/enhance-bullet-point")
async def ai_enhance_bullet_point(
    request: AIBulletPointEnhanceRequest
):
    """AIå¢å¼ºè¦ç‚¹æ¥å£"""
    try:
        # è·å–AIæä¾›è€…
        provider, settings = get_role_provider("outline")

        # æ„å»ºä¸Šä¸‹æ–‡ä¿¡æ¯
        context_info = ""
        if request.contextInfo:
            original_point = request.contextInfo.get('originalBulletPoint', '')
            other_points = request.contextInfo.get('otherBulletPoints', [])
            point_index = request.contextInfo.get('pointIndex', 0)

            context_info = f"""
å½“å‰è¦ç‚¹ä¸Šä¸‹æ–‡ä¿¡æ¯:
- è¦ç‚¹ä½ç½®:ç¬¬{point_index + 1}ä¸ªè¦ç‚¹
- åŸå§‹è¦ç‚¹å†…å®¹:{original_point}
- åŒé¡µé¢å…¶ä»–è¦ç‚¹:{', '.join(other_points) if other_points else 'æ— '}
"""

        # æ„å»ºå¤§çº²ä¿¡æ¯
        outline_info = ""
        if request.slideOutline:
            outline_info = f"""
å½“å‰å¹»ç¯ç‰‡å¤§çº²ä¿¡æ¯:
- å¹»ç¯ç‰‡ç±»å‹:{request.slideOutline.get('slide_type', 'æœªçŸ¥')}
- æè¿°:{request.slideOutline.get('description', 'æ— ')}
- æ‰€æœ‰è¦ç‚¹:{', '.join(request.slideOutline.get('content_points', [])) if request.slideOutline.get('content_points') else 'æ— '}
"""

        # æ„å»ºAIå¢å¼ºæç¤ºè¯
        context = f"""
ä½ æ˜¯ä¸€ä½ä¸“ä¸šçš„PPTå†…å®¹ç¼–è¾‘ä¸“å®¶.ç”¨æˆ·éœ€è¦ä½ å¢å¼ºå’Œä¼˜åŒ–ä¸€ä¸ªPPTè¦ç‚¹çš„å†…å®¹.

é¡¹ç›®ä¿¡æ¯:
- é¡¹ç›®æ ‡é¢˜:{request.projectInfo.get('title', 'æœªçŸ¥')}
- é¡¹ç›®ä¸»é¢˜:{request.projectInfo.get('topic', 'æœªçŸ¥')}
- åº”ç”¨åœºæ™¯:{request.projectInfo.get('scenario', 'æœªçŸ¥')}

å¹»ç¯ç‰‡ä¿¡æ¯:
- å¹»ç¯ç‰‡æ ‡é¢˜:{request.slideTitle}
- å¹»ç¯ç‰‡ä½ç½®:ç¬¬{request.slideIndex}é¡µ

{outline_info}

{context_info}

ç”¨æˆ·è¯·æ±‚:{request.userRequest}

è¯·æ ¹æ®ä»¥ä¸Šä¿¡æ¯,å¯¹è¦ç‚¹è¿›è¡Œå¢å¼ºå’Œä¼˜åŒ–.è¦æ±‚:

1. **ä¿æŒæ ¸å¿ƒæ„æ€ä¸å˜**:ä¸è¦æ”¹å˜è¦ç‚¹çš„åŸºæœ¬å«ä¹‰å’Œæ–¹å‘
2. **å¢åŠ å…·ä½“ç»†èŠ‚**:æ·»åŠ æ›´å¤šå…·ä½“çš„æè¿°,æ•°æ®,ä¾‹å­æˆ–è¯´æ˜
3. **æå‡è¡¨è¾¾è´¨é‡**:ä½¿ç”¨æ›´ä¸“ä¸š,æ›´æœ‰å¸å¼•åŠ›çš„è¡¨è¾¾æ–¹å¼
4. **ä¿æŒç®€æ´æ€§**:è™½ç„¶è¦å¢å¼ºå†…å®¹,ä½†ä»è¦ä¿æŒè¦ç‚¹çš„ç®€æ´ç‰¹æ€§,ä¸è¦è¿‡äºå†—é•¿
5. **ä¸å…¶ä»–è¦ç‚¹åè°ƒ**:ç¡®ä¿å¢å¼ºåçš„è¦ç‚¹ä¸åŒé¡µé¢å…¶ä»–è¦ç‚¹åœ¨é£æ ¼å’Œå±‚æ¬¡ä¸Šä¿æŒä¸€è‡´
6. **ç¬¦åˆåœºæ™¯éœ€æ±‚**:æ ¹æ®åº”ç”¨åœºæ™¯è°ƒæ•´è¯­è¨€é£æ ¼å’Œä¸“ä¸šç¨‹åº¦

è¯·ç›´æ¥è¿”å›å¢å¼ºåçš„è¦ç‚¹å†…å®¹,ä¸éœ€è¦é¢å¤–çš„è§£é‡Šæˆ–æ ¼å¼åŒ–.
"""

        # è°ƒç”¨AIç”Ÿæˆå¢å¼ºå†…å®¹
        response = await provider.text_completion(
            prompt=context,
            max_tokens=ai_config.max_tokens // 2,  # ä½¿ç”¨è¾ƒå°çš„tokené™åˆ¶
            temperature=0.7,
            model=settings.get('model')
        )

        enhanced_text = response.content.strip()

        # ç®€å•çš„å†…å®¹éªŒè¯
        if not enhanced_text or len(enhanced_text) < 5:
            raise ValueError("AIç”Ÿæˆçš„å¢å¼ºå†…å®¹è¿‡çŸ­æˆ–ä¸ºç©º")

        return {
            "success": True,
            "enhancedText": enhanced_text,
            "originalText": request.contextInfo.get('originalBulletPoint', '') if request.contextInfo else ""
        }

    except Exception as e:
        logger.error(f"AIè¦ç‚¹å¢å¼ºè¯·æ±‚å¤±è´¥: {e}")
        return {
            "success": False,
            "error": str(e),
            "message": "æŠ±æ­‰,AIè¦ç‚¹å¢å¼ºæœåŠ¡æš‚æ—¶ä¸å¯ç”¨.è¯·ç¨åé‡è¯•."
        }

@router.post("/api/ai/enhance-all-bullet-points")
async def ai_enhance_all_bullet_points(
    request: AIBulletPointEnhanceRequest
):
    """AIå¢å¼ºæ‰€æœ‰è¦ç‚¹æ¥å£"""
    try:
        # è·å–AIæä¾›è€…
        provider, settings = get_role_provider("outline")

        # æ„å»ºä¸Šä¸‹æ–‡ä¿¡æ¯
        context_info = ""
        all_points = []
        if request.contextInfo:
            all_points = request.contextInfo.get('allBulletPoints', [])
            total_points = request.contextInfo.get('totalPoints', 0)

            context_info = f"""
å½“å‰è¦ç‚¹ä¸Šä¸‹æ–‡ä¿¡æ¯:
- è¦ç‚¹æ€»æ•°:{total_points}ä¸ª
- æ‰€æœ‰è¦ç‚¹å†…å®¹:
"""
            for i, point in enumerate(all_points, 1):
                context_info += f"  {i}. {point}\n"

        # æ„å»ºå¤§çº²ä¿¡æ¯
        outline_info = ""
        if request.slideOutline:
            outline_info = f"""
å½“å‰å¹»ç¯ç‰‡å¤§çº²ä¿¡æ¯:
- å¹»ç¯ç‰‡ç±»å‹:{request.slideOutline.get('slide_type', 'æœªçŸ¥')}
- æè¿°:{request.slideOutline.get('description', 'æ— ')}
"""

        # æ„å»ºAIå¢å¼ºæç¤ºè¯
        context = f"""
è¯·å¯¹ä»¥ä¸‹PPTè¦ç‚¹è¿›è¡Œå¢å¼ºå’Œä¼˜åŒ–.

é¡¹ç›®èƒŒæ™¯:
- é¡¹ç›®:{request.projectInfo.get('title', 'æœªçŸ¥')}
- ä¸»é¢˜:{request.projectInfo.get('topic', 'æœªçŸ¥')}
- åœºæ™¯:{request.projectInfo.get('scenario', 'æœªçŸ¥')}
- å¹»ç¯ç‰‡:{request.slideTitle}(ç¬¬{request.slideIndex}é¡µ)

{outline_info}

{context_info}

å¢å¼ºè¦æ±‚:
1. ä¿æŒæ¯ä¸ªè¦ç‚¹çš„æ ¸å¿ƒæ„æ€ä¸å˜
2. æ·»åŠ å…·ä½“ç»†èŠ‚,æ•°æ®æˆ–ä¾‹å­
3. ä½¿ç”¨æ›´ä¸“ä¸š,å‡†ç¡®çš„è¡¨è¾¾
4. ä¿æŒç®€æ´,é¿å…å†—é•¿
5. ç¡®ä¿è¦ç‚¹é—´é€»è¾‘è¿è´¯,é£æ ¼ç»Ÿä¸€
6. ç¬¦åˆ{request.projectInfo.get('scenario', 'å•†åŠ¡')}åœºæ™¯çš„ä¸“ä¸šè¦æ±‚

é‡è¦:è¯·ç›´æ¥è¿”å›å¢å¼ºåçš„è¦ç‚¹åˆ—è¡¨,æ¯è¡Œä¸€ä¸ªè¦ç‚¹,ä¸è¦åŒ…å«ä»»ä½•è§£é‡Š,å¼€åœºç™½æˆ–æ ¼å¼è¯´æ˜.ä¸è¦æ·»åŠ ç¼–å·,ç¬¦å·æˆ–å…¶ä»–æ ‡è®°.

ç¤ºä¾‹æ ¼å¼:
ç¬¬ä¸€ä¸ªå¢å¼ºåçš„è¦ç‚¹å†…å®¹
ç¬¬äºŒä¸ªå¢å¼ºåçš„è¦ç‚¹å†…å®¹
ç¬¬ä¸‰ä¸ªå¢å¼ºåçš„è¦ç‚¹å†…å®¹
"""

        # è°ƒç”¨AIç”Ÿæˆå¢å¼ºå†…å®¹
        response = await provider.text_completion(
            prompt=context,
            max_tokens=ai_config.max_tokens,  # ä½¿ç”¨å®Œæ•´çš„tokené™åˆ¶,å› ä¸ºè¦å¤„ç†å¤šä¸ªè¦ç‚¹
            temperature=0.7,
            model=settings.get('model')
        )

        enhanced_content = response.content.strip()

        # è§£æå¢å¼ºåçš„è¦ç‚¹ - æ”¹è¿›çš„è¿‡æ»¤é€»è¾‘
        enhanced_points = []
        if enhanced_content:
            # æŒ‰è¡Œåˆ†å‰²,è¿‡æ»¤ç©ºè¡Œ
            lines = [line.strip() for line in enhanced_content.split('\n') if line.strip()]

            # è¿‡æ»¤æ‰å¸¸è§çš„æ— å…³å†…å®¹
            filtered_lines = []
            skip_patterns = [
                'å¥½çš„', 'ä½œä¸º', 'æˆ‘å°†', 'æˆ‘ä¼š', 'ä»¥ä¸‹æ˜¯', 'æ ¹æ®', 'è¯·æ³¨æ„', 'éœ€è¦è¯´æ˜',
                'å¢å¼ºåçš„è¦ç‚¹', 'ä¼˜åŒ–åçš„', 'æ”¹è¿›åçš„', 'ä»¥ä¸Š', 'æ€»ç»“', 'å¸Œæœ›',
                'å¦‚æœ‰', 'å¦‚æœ', 'å»ºè®®', 'æ¨è', 'æ³¨æ„', 'æé†’', 'è¯´æ˜',
                'è¦ç‚¹1', 'è¦ç‚¹2', 'è¦ç‚¹3', 'è¦ç‚¹4', 'è¦ç‚¹5',
                'ç¬¬ä¸€', 'ç¬¬äºŒ', 'ç¬¬ä¸‰', 'ç¬¬å››', 'ç¬¬äº”', 'ç¬¬å…­', 'ç¬¬ä¸ƒ', 'ç¬¬å…«', 'ç¬¬ä¹', 'ç¬¬å',
                '1.', '2.', '3.', '4.', '5.', '6.', '7.', '8.', '9.', '10.',
                'â€¢', 'Â·', '-', '*', 'â†’', 'â–ª', 'â–«'
            ]

            for line in lines:
                # è·³è¿‡è¿‡çŸ­çš„è¡Œ(å¯èƒ½æ˜¯æ ¼å¼æ ‡è®°)
                if len(line) < 5:
                    continue

                # è·³è¿‡åŒ…å«å¸¸è§å¼€åœºç™½æ¨¡å¼çš„è¡Œ
                should_skip = False
                for pattern in skip_patterns:
                    if line.startswith(pattern) or (pattern in ['å¥½çš„', 'ä½œä¸º', 'æˆ‘å°†', 'æˆ‘ä¼š'] and pattern in line[:10]):
                        should_skip = True
                        break

                # è·³è¿‡çº¯æ•°å­—æˆ–ç¬¦å·å¼€å¤´çš„è¡Œ(å¯èƒ½æ˜¯ç¼–å·)
                if line[0].isdigit() or line[0] in ['â€¢', 'Â·', '-', '*', 'â†’', 'â–ª', 'â–«']:
                    # ä½†ä¿ç•™å»æ‰ç¼–å·åçš„å†…å®¹
                    cleaned_line = line
                    # ç§»é™¤å¼€å¤´çš„ç¼–å·å’Œç¬¦å·
                    import re
                    cleaned_line = re.sub(r'^[\d\s\.\-\*\â€¢\Â·\â†’\â–ª\â–«]+', '', cleaned_line).strip()
                    if len(cleaned_line) >= 5:
                        filtered_lines.append(cleaned_line)
                    continue

                if not should_skip:
                    filtered_lines.append(line)

            enhanced_points = filtered_lines

        # ç®€å•çš„å†…å®¹éªŒè¯
        if not enhanced_points or len(enhanced_points) == 0:
            raise ValueError("AIç”Ÿæˆçš„å¢å¼ºå†…å®¹ä¸ºç©ºæˆ–è¢«è¿‡æ»¤")

        return {
            "success": True,
            "enhancedPoints": enhanced_points,
            "originalPoints": all_points,
            "totalEnhanced": len(enhanced_points)
        }

    except Exception as e:
        logger.error(f"AIå¢å¼ºæ‰€æœ‰è¦ç‚¹è¯·æ±‚å¤±è´¥: {e}")
        return {
            "success": False,
            "error": str(e),
            "message": "æŠ±æ­‰,AIè¦ç‚¹å¢å¼ºæœåŠ¡æš‚æ—¶ä¸å¯ç”¨.è¯·ç¨åé‡è¯•."
        }

@router.post("/api/projects/{project_id}/speech-script/generate")
async def generate_speech_script(
    project_id: str,
    request: SpeechScriptGenerationRequest
):
    """Generate speech scripts for presentation slides"""
    try:
        import uuid
        import asyncio

        # Generate task ID for progress tracking
        task_id = str(uuid.uuid4())

        # Get project
        project = await ppt_service.project_manager.get_project(project_id)
        if not project:
            return {
                "success": False,
                "error": "Project not found"
            }

        # Check if slides data exists
        if not project.slides_data or len(project.slides_data) == 0:
            return {
                "success": False,
                "error": "No slides data available"
            }

        # Import speech script service
        from ..services.speech_script_service import SpeechScriptService, SpeechScriptCustomization
        from ..services.speech_script_service import SpeechTone, TargetAudience, LanguageComplexity

        # Initialize service
        speech_service = SpeechScriptService()

        # Parse customization options
        customization_data = request.customization
        customization = SpeechScriptCustomization(
            tone=SpeechTone(customization_data.get('tone', 'conversational')),
            target_audience=TargetAudience(customization_data.get('target_audience', 'general_public')),
            language_complexity=LanguageComplexity(customization_data.get('language_complexity', 'moderate')),
            custom_style_prompt=customization_data.get('custom_style_prompt'),
            include_transitions=customization_data.get('include_transitions', True),
            include_timing_notes=customization_data.get('include_timing_notes', False),
            speaking_pace=customization_data.get('speaking_pace', 'normal')
        )

        # Validate request parameters
        if request.generation_type == "single":
            if not request.slide_indices or len(request.slide_indices) != 1:
                return {
                    "success": False,
                    "error": "Single generation requires exactly one slide index"
                }
        elif request.generation_type == "multi":
            if not request.slide_indices:
                return {
                    "success": False,
                    "error": "Multi generation requires slide indices"
                }
        elif request.generation_type != "full":
            return {
                "success": False,
                "error": "Invalid generation type"
            }

        # Start async generation task
        async def generate_async():
            try:
                logger.info(f"Starting async generation for task {task_id}")

                # Generate scripts based on type
                if request.generation_type == "single":
                    # Use multi_slide_scripts_with_retry for single slide to get progress tracking
                    result = await speech_service.generate_multi_slide_scripts_with_retry(
                        project, request.slide_indices, customization, task_id=task_id
                    )
                elif request.generation_type == "multi":
                    result = await speech_service.generate_multi_slide_scripts_with_retry(
                        project, request.slide_indices, customization, task_id=task_id
                    )
                elif request.generation_type == "full":
                    result = await speech_service.generate_full_presentation_scripts(
                        project, customization, progress_callback=None, task_id=task_id
                    )

                # Save scripts to database if successful
                if result.success:
                    logger.info(f"Generation successful for task {task_id}, saving to database")
                    from ..services.speech_script_repository import SpeechScriptRepository
                    repo = SpeechScriptRepository()

                    generation_params = {
                        'generation_type': request.generation_type,
                        'tone': customization.tone.value,
                        'target_audience': customization.target_audience.value,
                        'language_complexity': customization.language_complexity.value,
                        'custom_audience': request.customization.get('custom_audience'),
                        'custom_style_prompt': customization.custom_style_prompt,
                        'include_transitions': customization.include_transitions,
                        'include_timing_notes': customization.include_timing_notes,
                        'speaking_pace': customization.speaking_pace
                    }

                    saved_count = 0
                    for script in result.scripts:
                        await repo.save_speech_script(
                            project_id=project_id,
                            slide_index=script.slide_index,
                            slide_title=script.slide_title,
                            script_content=script.script_content,
                            generation_params=generation_params,
                            estimated_duration=script.estimated_duration
                        )
                        saved_count += 1
                        logger.debug(f"Saved script {saved_count}/{len(result.scripts)} for slide {script.slide_index}")

                    # Ensure all changes are committed before closing
                    repo.db.commit()
                    repo.close()
                    logger.info(f"All {saved_count} scripts saved and committed to database for task {task_id}")

                    # NOW mark the task as completed after database save
                    from ..services.progress_tracker import progress_tracker
                    progress_tracker.complete_task(
                        task_id,
                        f"ç”Ÿæˆå®Œæˆ!æˆåŠŸ {saved_count} é¡µ"
                    )
                    logger.info(f"Task {task_id} marked as completed")
                else:
                    logger.error(f"Generation failed for task {task_id}: {result.error_message}")

            except Exception as e:
                logger.error(f"Async speech script generation failed for task {task_id}: {e}")
                from ..services.progress_tracker import progress_tracker
                progress_tracker.fail_task(task_id, str(e))

        # Start the async task
        asyncio.create_task(generate_async())

        # Return immediately with task_id
        return {
            "success": True,
            "task_id": task_id,
            "message": "æ¼”è®²ç¨¿ç”Ÿæˆå·²å¼€å§‹,è¯·æŸ¥çœ‹è¿›åº¦"
        }

    except Exception as e:
        logger.error(f"Speech script generation failed: {e}")
        return {
            "success": False,
            "error": str(e)
        }

@router.post("/api/projects/{project_id}/speech-script/export")
async def export_speech_script(
    project_id: str,
    request: SpeechScriptExportRequest
):
    """Export speech scripts to document format"""
    try:
        # Get project for title
        project = await ppt_service.project_manager.get_project(project_id)
        if not project:
            return {
                "success": False,
                "error": "Project not found"
            }

        # Import exporter
        from ..services.speech_script_exporter import get_speech_script_exporter
        from ..services.speech_script_service import SlideScriptData

        exporter = get_speech_script_exporter()

        # Validate scripts data
        if not request.scripts_data or len(request.scripts_data) == 0:
            return {
                "success": False,
                "error": "No speech scripts data provided"
            }

        # Convert scripts data to SlideScriptData objects
        scripts = []
        for script_data in request.scripts_data:
            # Validate required fields
            if not script_data.get('script_content'):
                continue  # Skip empty scripts

            script = SlideScriptData(
                slide_index=script_data.get('slide_index', 0),
                slide_title=script_data.get('slide_title', ''),
                script_content=script_data.get('script_content', ''),
                estimated_duration=script_data.get('estimated_duration'),
                speaker_notes=script_data.get('speaker_notes')
            )
            scripts.append(script)

        # Check if we have any valid scripts after filtering
        if not scripts:
            return {
                "success": False,
                "error": "No valid speech scripts found"
            }

        # Prepare metadata
        metadata = {}
        if request.include_metadata:
            # Calculate total estimated duration from all scripts
            total_duration = None
            if scripts:
                total_minutes = 0
                for script in scripts:
                    if script.estimated_duration and 'åˆ†é’Ÿ' in script.estimated_duration:
                        try:
                            minutes = float(script.estimated_duration.replace('åˆ†é’Ÿ', ''))
                            total_minutes += minutes
                        except ValueError:
                            pass
                if total_minutes > 0:
                    total_duration = f"{total_minutes:.1f}åˆ†é’Ÿ"

            metadata = {
                'generation_time': time.time(),
                'total_estimated_duration': total_duration,
                'customization': {}
            }

        # Export based on format
        if request.export_format == "docx":
            if not exporter.is_docx_available():
                return {
                    "success": False,
                    "error": "DOCX export not available. Please install python-docx."
                }

            docx_content = await exporter.export_to_docx(
                scripts, project.topic, metadata
            )

            # Return file response
            import urllib.parse
            filename = f"{project.topic}_æ¼”è®²ç¨¿.docx"
            safe_filename = urllib.parse.quote(filename, safe='')

            from fastapi.responses import Response
            return Response(
                content=docx_content,
                media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                headers={
                    "Content-Disposition": f"attachment; filename*=UTF-8''{safe_filename}"
                }
            )

        elif request.export_format == "markdown":
            markdown_content = await exporter.export_to_markdown(
                scripts, project.topic, metadata
            )

            # Return file response
            import urllib.parse
            filename = f"{project.topic}_æ¼”è®²ç¨¿.md"
            safe_filename = urllib.parse.quote(filename, safe='')

            from fastapi.responses import Response
            return Response(
                content=markdown_content.encode('utf-8'),
                media_type="text/markdown",
                headers={
                    "Content-Disposition": f"attachment; filename*=UTF-8''{safe_filename}"
                }
            )

        else:
            return {
                "success": False,
                "error": "Unsupported export format"
            }

    except Exception as e:
        logger.error(f"Speech script export failed: {e}")
        return {
            "success": False,
            "error": str(e)
        }


@router.get("/api/projects/{project_id}/speech-scripts")
async def get_current_speech_scripts(
    project_id: str
):
    """è·å–é¡¹ç›®çš„å½“å‰æ¼”è®²ç¨¿"""
    try:
        from ..services.speech_script_repository import SpeechScriptRepository

        # æ£€æŸ¥é¡¹ç›®æ˜¯å¦å­˜åœ¨
        project = await ppt_service.project_manager.get_project(project_id)
        if not project:
            return {
                "success": False,
                "error": "Project not found"
            }

        repo = SpeechScriptRepository()

        # Expire all objects to ensure fresh data from database
        repo.db.expire_all()

        # è·å–é¡¹ç›®çš„å½“å‰æ¼”è®²ç¨¿
        scripts = await repo.get_current_speech_scripts_by_project(project_id)
        logger.info(f"Found {len(scripts)} speech scripts for project {project_id}")

        # è½¬æ¢ä¸ºJSONæ ¼å¼
        scripts_data = []
        for script in scripts:
            scripts_data.append({
                "id": script.id,
                "slide_index": script.slide_index,
                "slide_title": script.slide_title,
                "script_content": script.script_content,
                "estimated_duration": script.estimated_duration,
                "speaker_notes": script.speaker_notes,
                "generation_type": script.generation_type,
                "tone": script.tone,
                "target_audience": script.target_audience,
                "custom_audience": script.custom_audience,
                "language_complexity": script.language_complexity,
                "speaking_pace": script.speaking_pace,
                "custom_style_prompt": script.custom_style_prompt,
                "include_transitions": script.include_transitions,
                "include_timing_notes": script.include_timing_notes,
                "created_at": script.created_at,
                "updated_at": script.updated_at
            })

        repo.close()

        return {
            "success": True,
            "scripts": scripts_data
        }

    except Exception as e:
        logger.error(f"Get current speech scripts failed: {e}")
        return {
            "success": False,
            "error": str(e)
        }



@router.delete("/api/projects/{project_id}/speech-scripts/slide/{slide_index}")
async def delete_speech_script_by_slide(
    project_id: str,
    slide_index: int
):
    """åˆ é™¤æŒ‡å®šå¹»ç¯ç‰‡çš„æ¼”è®²ç¨¿"""
    try:
        from ..services.speech_script_repository import SpeechScriptRepository

        # æ£€æŸ¥é¡¹ç›®æ˜¯å¦å­˜åœ¨
        project = await ppt_service.project_manager.get_project(project_id)
        if not project:
            return {
                "success": False,
                "error": "Project not found"
            }

        repo = SpeechScriptRepository()

        # è·å–å¹¶åˆ é™¤æŒ‡å®šå¹»ç¯ç‰‡çš„æ¼”è®²ç¨¿
        script = await repo.get_speech_script_by_slide(project_id, slide_index)
        if not script:
            return {
                "success": False,
                "error": "Speech script not found"
            }

        success = await repo.delete_speech_script(script.id)

        return {
            "success": success,
            "message": f"ç¬¬{slide_index + 1}é¡µæ¼”è®²ç¨¿å·²åˆ é™¤" if success else "åˆ é™¤æ¼”è®²ç¨¿å¤±è´¥"
        }

    except Exception as e:
        logger.error(f"Delete speech script failed: {e}")
        return {
            "success": False,
            "error": str(e)
        }


@router.get("/api/projects/{project_id}/speech-scripts/result/{task_id}")
async def get_speech_script_result(
    project_id: str,
    task_id: str
):
    """è·å–æ¼”è®²ç¨¿ç”Ÿæˆç»“æœ"""
    try:
        from ..services.progress_tracker import progress_tracker
        from ..services.speech_script_repository import SpeechScriptRepository

        # æ£€æŸ¥é¡¹ç›®æ˜¯å¦å­˜åœ¨
        project = await ppt_service.project_manager.get_project(project_id)
        if not project:
            return {
                "success": False,
                "error": "Project not found"
            }

        # è·å–è¿›åº¦ä¿¡æ¯
        progress_info = progress_tracker.get_progress(task_id)

        if not progress_info:
            return {
                "success": False,
                "error": "Task not found"
            }

        # éªŒè¯ä»»åŠ¡æ˜¯å¦å±äºè¯¥é¡¹ç›®
        if progress_info.project_id != project_id:
            return {
                "success": False,
                "error": "Access denied"
            }

        # å¦‚æœä»»åŠ¡è¿˜æœªå®Œæˆ,è¿”å›è¿›åº¦ä¿¡æ¯
        if progress_info.status != "completed":
            return {
                "success": False,
                "error": "Task not completed yet",
                "status": progress_info.status,
                "progress": progress_info.to_dict()
            }

        # è·å–ç”Ÿæˆçš„æ¼”è®²ç¨¿
        repo = SpeechScriptRepository()
        scripts = await repo.get_current_speech_scripts_by_project(project_id)

        # è½¬æ¢ä¸ºAPIæ ¼å¼
        scripts_data = []
        total_duration_seconds = 0

        for script in scripts:
            script_data = {
                "slide_index": script.slide_index,
                "slide_title": script.slide_title,
                "script_content": script.script_content,
                "estimated_duration": script.estimated_duration,
                "speaker_notes": getattr(script, 'speaker_notes', None)
            }
            scripts_data.append(script_data)

            # è®¡ç®—æ€»æ—¶é•¿
            if script.estimated_duration:
                try:
                    if 'åˆ†é’Ÿ' in script.estimated_duration:
                        minutes = float(script.estimated_duration.replace('åˆ†é’Ÿ', ''))
                        total_duration_seconds += minutes * 60
                    elif 'ç§’' in script.estimated_duration:
                        seconds = float(script.estimated_duration.replace('ç§’', ''))
                        total_duration_seconds += seconds
                except:
                    pass

        # æ ¼å¼åŒ–æ€»æ—¶é•¿
        if total_duration_seconds < 60:
            total_duration = f"{int(total_duration_seconds)}ç§’"
        else:
            minutes = total_duration_seconds / 60
            total_duration = f"{minutes:.1f}åˆ†é’Ÿ"

        repo.close()

        return {
            "success": True,
            "scripts": scripts_data,
            "total_estimated_duration": total_duration,
            "generation_metadata": {
                "task_id": task_id,
                "completed_at": progress_info.last_update,
                "total_slides": progress_info.total_slides,
                "completed_slides": progress_info.completed_slides,
                "failed_slides": progress_info.failed_slides,
                "skipped_slides": progress_info.skipped_slides
            }
        }

    except Exception as e:
        logger.error(f"Get speech script result failed: {e}")
        return {
            "success": False,
            "error": str(e)
        }


@router.get("/api/projects/{project_id}/speech-scripts/progress/{task_id}")
async def get_speech_script_progress(
    project_id: str,
    task_id: str
):
    """è·å–æ¼”è®²ç¨¿ç”Ÿæˆè¿›åº¦"""
    try:
        from ..services.progress_tracker import progress_tracker

        # æ£€æŸ¥é¡¹ç›®æ˜¯å¦å­˜åœ¨
        project = await ppt_service.project_manager.get_project(project_id)
        if not project:
            return {
                "success": False,
                "error": "Project not found"
            }

        # è·å–è¿›åº¦ä¿¡æ¯
        progress_info = progress_tracker.get_progress(task_id)

        if not progress_info:
            return {
                "success": False,
                "error": "Task not found"
            }

        # éªŒè¯ä»»åŠ¡æ˜¯å¦å±äºè¯¥é¡¹ç›®
        if progress_info.project_id != project_id:
            return {
                "success": False,
                "error": "Access denied"
            }

        return {
            "success": True,
            "progress": progress_info.to_dict()
        }

    except Exception as e:
        logger.error(f"Get speech script progress failed: {e}")
        return {
            "success": False,
            "error": str(e)
        }


@router.put("/api/projects/{project_id}/speech-scripts/slide/{slide_index}")
async def update_speech_script_content(
    project_id: str,
    slide_index: int,
    request: dict
):
    """æ›´æ–°æ¼”è®²ç¨¿å†…å®¹"""
    try:
        from ..services.speech_script_repository import SpeechScriptRepository

        # æ£€æŸ¥é¡¹ç›®æ˜¯å¦å­˜åœ¨
        project = await ppt_service.project_manager.get_project(project_id)
        if not project:
            return {
                "success": False,
                "error": "Project not found"
            }

        # è·å–è¯·æ±‚æ•°æ®
        script_content = request.get('script_content', '').strip()
        slide_title = request.get('slide_title', f'ç¬¬{slide_index + 1}é¡µ')
        estimated_duration = request.get('estimated_duration')
        speaker_notes = request.get('speaker_notes')

        if not script_content:
            return {
                "success": False,
                "error": "æ¼”è®²ç¨¿å†…å®¹ä¸èƒ½ä¸ºç©º"
            }

        repo = SpeechScriptRepository()

        # è·å–ç°æœ‰æ¼”è®²ç¨¿
        existing_script = await repo.get_speech_script_by_slide(project_id, slide_index)
        if not existing_script:
            return {
                "success": False,
                "error": "æ¼”è®²ç¨¿ä¸å­˜åœ¨"
            }

        # æ›´æ–°å†…å®¹
        existing_script.script_content = script_content
        existing_script.slide_title = slide_title
        if estimated_duration:
            existing_script.estimated_duration = estimated_duration
        if speaker_notes is not None:
            existing_script.speaker_notes = speaker_notes
        existing_script.updated_at = time.time()

        repo.db.commit()
        repo.db.refresh(existing_script)
        repo.close()

        return {
            "success": True,
            "message": "æ¼”è®²ç¨¿å·²æ›´æ–°",
            "script": {
                "id": existing_script.id,
                "slide_index": existing_script.slide_index,
                "slide_title": existing_script.slide_title,
                "script_content": existing_script.script_content,
                "estimated_duration": existing_script.estimated_duration,
                "speaker_notes": existing_script.speaker_notes,
                "updated_at": existing_script.updated_at
            }
        }

    except Exception as e:
        logger.error(f"Update speech script content failed: {e}")
        return {
            "success": False,
            "error": str(e)
        }


@router.get("/api/projects/{project_id}/selected-global-template")
async def get_selected_global_template(
    project_id: str
):
    """è·å–é¡¹ç›®é€‰æ‹©çš„å…¨å±€æ¯ç‰ˆæ¨¡æ¿"""
    try:
        # æ£€æŸ¥é¡¹ç›®æ˜¯å¦çœŸæ­£é€‰æ‹©äº†æ¨¡æ¿
        selected_template = await ppt_service.get_selected_global_template(project_id)
        if selected_template:
            logger.info(f"Project {project_id} has selected template: {selected_template.get('template_name', 'Unknown')}")
            return {
                "status": "success",
                "template": selected_template,
                "is_user_selected": True
            }
        else:
            # å¦‚æœæ²¡æœ‰é€‰æ‹©çš„æ¨¡æ¿,å°è¯•è·å–é»˜è®¤æ¨¡æ¿
            default_template = await ppt_service.global_template_service.get_default_template()
            if default_template:
                logger.info(f"Project {project_id} using default template: {default_template.get('template_name', 'Unknown')}")
                return {
                    "status": "success",
                    "template": default_template,
                    "is_user_selected": False
                }
            else:
                logger.warning(f"No template available for project {project_id}")
                return {
                    "status": "success",
                    "template": None,
                    "is_user_selected": False
                }
    except Exception as e:
        logger.error(f"Error getting selected global template for project {project_id}: {e}")
        raise HTTPException(status_code=500, detail=str(e))

@router.post("/api/projects/{project_id}/slides/{slide_index}/save")
async def save_single_slide_content(
    project_id: str,
    slide_index: int,
    request: Request
):
    """ä¿å­˜å•ä¸ªå¹»ç¯ç‰‡å†…å®¹åˆ°æ•°æ®åº“"""
    try:
        logger.info(f"ğŸ”„ å¼€å§‹ä¿å­˜é¡¹ç›® {project_id} çš„ç¬¬ {slide_index + 1} é¡µ (ç´¢å¼•: {slide_index})")

        data = await request.json()
        html_content = data.get('html_content', '')

        logger.info(f"ğŸ“„ æ¥æ”¶åˆ°HTMLå†…å®¹,é•¿åº¦: {len(html_content)} å­—ç¬¦")

        if not html_content:
            logger.error("âŒ HTMLå†…å®¹ä¸ºç©º")
            raise HTTPException(status_code=400, detail="HTML content is required")

        project = await ppt_service.project_manager.get_project(project_id)
        if not project:
            logger.error(f"âŒ é¡¹ç›® {project_id} ä¸å­˜åœ¨")
            raise HTTPException(status_code=404, detail="Project not found")

        # è¯¦ç»†éªŒè¯å¹»ç¯ç‰‡ç´¢å¼•
        total_slides = len(project.slides_data) if project.slides_data else 0
        logger.debug(f"ğŸ“Š é¡¹ç›®å¹»ç¯ç‰‡ä¿¡æ¯: æ€»é¡µæ•°={total_slides}, è¯·æ±‚ç´¢å¼•={slide_index}")

        if slide_index < 0:
            logger.error(f"âŒ å¹»ç¯ç‰‡ç´¢å¼•ä¸èƒ½ä¸ºè´Ÿæ•°: {slide_index}")
            raise HTTPException(status_code=400, detail=f"Slide index cannot be negative: {slide_index}")

        if slide_index >= total_slides:
            logger.error(f"âŒ å¹»ç¯ç‰‡ç´¢å¼•è¶…å‡ºèŒƒå›´: {slide_index},é¡¹ç›®å…±æœ‰ {total_slides} é¡µ")
            raise HTTPException(status_code=400, detail=f"Slide index {slide_index} out of range (total: {total_slides})")

        logger.debug(f"ğŸ“ æ›´æ–°ç¬¬ {slide_index + 1} é¡µçš„å†…å®¹")

        # æ›´æ–°å¹»ç¯ç‰‡æ•°æ®
        project.slides_data[slide_index]['html_content'] = html_content
        project.slides_data[slide_index]['is_user_edited'] = True
        project.updated_at = time.time()

        # é‡æ–°ç”Ÿæˆç»„åˆHTML
        if project.slides_data:
            outline_title = project.title
            if isinstance(project.outline, dict):
                outline_title = project.outline.get('title', project.title)
            elif hasattr(project.outline, 'title'):
                outline_title = project.outline.title

            project.slides_html = ppt_service._combine_slides_to_full_html(
                project.slides_data, outline_title
            )

        # ä¿å­˜åˆ°æ•°æ®åº“
        try:
            logger.debug(f"ğŸ’¾ å¼€å§‹ä¿å­˜åˆ°æ•°æ®åº“... (ç¬¬{slide_index + 1}é¡µ)")

            from ..services.db_project_manager import DatabaseProjectManager
            db_manager = DatabaseProjectManager()

            # ä¿å­˜å•ä¸ªå¹»ç¯ç‰‡
            slide_data = project.slides_data[slide_index]
            slide_title = slide_data.get('title', 'æ— æ ‡é¢˜')
            is_user_edited = slide_data.get('is_user_edited', False)

            logger.debug(f"ğŸ“Š å¹»ç¯ç‰‡æ•°æ®: æ ‡é¢˜='{slide_title}', ç”¨æˆ·ç¼–è¾‘={is_user_edited}, ç´¢å¼•={slide_index}")
            logger.debug(f"ğŸ” ä¿å­˜å‰éªŒè¯: é¡¹ç›®ID={project_id}, å¹»ç¯ç‰‡ç´¢å¼•={slide_index}")

            save_success = await db_manager.save_single_slide(project_id, slide_index, slide_data)

            if save_success:
                logger.debug(f"âœ… ç¬¬ {slide_index + 1} é¡µå·²æˆåŠŸä¿å­˜åˆ°æ•°æ®åº“")

                return {
                    "success": True,
                    "message": f"Slide {slide_index + 1} saved successfully to database",
                    "slide_data": slide_data,
                    "database_saved": True
                }
            else:
                logger.error(f"âŒ ä¿å­˜ç¬¬ {slide_index + 1} é¡µåˆ°æ•°æ®åº“å¤±è´¥")
                return {
                    "success": False,
                    "error": "Failed to save slide to database",
                    "database_saved": False
                }

        except Exception as save_error:
            logger.error(f"âŒ ä¿å­˜ç¬¬ {slide_index + 1} é¡µæ—¶å‘ç”Ÿå¼‚å¸¸: {save_error}")
            import traceback
            traceback.print_exc()
            return {
                "success": False,
                "error": f"Database error: {str(save_error)}",
                "database_saved": False
            }

    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"âŒ ä¿å­˜å•ä¸ªå¹»ç¯ç‰‡æ—¶å‘ç”Ÿé”™è¯¯: {e}")
        import traceback
        traceback.print_exc()
        return {
            "success": False,
            "error": str(e),
            "database_saved": False
        }

@router.get("/api/projects/{project_id}/slides/stream")
async def stream_slides_generation(project_id: str):
    """Stream slides generation process"""
    try:
        async def generate_slides_stream():
            async for chunk in ppt_service.generate_slides_streaming(project_id):
                yield chunk

        return StreamingResponse(
            generate_slides_stream(),
            media_type="text/event-stream",
            headers={
                "Cache-Control": "no-cache",
                "Connection": "keep-alive",
                "Access-Control-Allow-Origin": "*",
                "Access-Control-Allow-Headers": "Cache-Control"
            }
        )

    except Exception as e:
        return {"error": str(e)}


@router.post("/api/projects/{project_id}/slides/cleanup")
async def cleanup_excess_slides(
    project_id: str,
    request: Request
):
    """æ¸…ç†é¡¹ç›®ä¸­å¤šä½™çš„å¹»ç¯ç‰‡"""
    try:
        logger.info(f"ğŸ§¹ å¼€å§‹æ¸…ç†é¡¹ç›® {project_id} çš„å¤šä½™å¹»ç¯ç‰‡")

        data = await request.json()
        current_slide_count = data.get('current_slide_count', 0)

        if current_slide_count <= 0:
            logger.error("âŒ æ— æ•ˆçš„å¹»ç¯ç‰‡æ•°é‡")
            raise HTTPException(status_code=400, detail="Invalid slide count")

        project = await ppt_service.project_manager.get_project(project_id)
        if not project:
            logger.error(f"âŒ é¡¹ç›® {project_id} ä¸å­˜åœ¨")
            raise HTTPException(status_code=404, detail="Project not found")

        # æ¸…ç†æ•°æ®åº“ä¸­å¤šä½™çš„å¹»ç¯ç‰‡
        from ..services.db_project_manager import DatabaseProjectManager
        db_manager = DatabaseProjectManager()
        deleted_count = await db_manager.cleanup_excess_slides(project_id, current_slide_count)

        logger.info(f"âœ… é¡¹ç›® {project_id} æ¸…ç†å®Œæˆ,åˆ é™¤äº† {deleted_count} å¼ å¤šä½™çš„å¹»ç¯ç‰‡")

        return {
            "success": True,
            "message": f"Successfully cleaned up {deleted_count} excess slides",
            "deleted_count": deleted_count
        }

    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"âŒ æ¸…ç†å¹»ç¯ç‰‡å¤±è´¥: {e}")
        return {"success": False, "error": str(e)}


@router.post("/api/projects/{project_id}/slides/batch-save")
async def batch_save_slides(
    project_id: str,
    request: Request
):
    """æ‰¹é‡ä¿å­˜æ‰€æœ‰å¹»ç¯ç‰‡ - é«˜æ•ˆç‰ˆæœ¬"""
    try:
        logger.debug(f"ğŸ”„ å¼€å§‹æ‰¹é‡ä¿å­˜é¡¹ç›® {project_id} çš„æ‰€æœ‰å¹»ç¯ç‰‡")

        data = await request.json()
        slides_data = data.get('slides_data', [])

        if not slides_data:
            logger.error("âŒ å¹»ç¯ç‰‡æ•°æ®ä¸ºç©º")
            raise HTTPException(status_code=400, detail="Slides data is required")

        project = await ppt_service.project_manager.get_project(project_id)
        if not project:
            logger.error(f"âŒ é¡¹ç›® {project_id} ä¸å­˜åœ¨")
            raise HTTPException(status_code=404, detail="Project not found")

        # æ›´æ–°é¡¹ç›®å†…å­˜ä¸­çš„æ•°æ®
        project.slides_data = slides_data
        project.updated_at = time.time()

        # é‡æ–°ç”Ÿæˆå®Œæ•´HTML
        outline_title = project.title
        if hasattr(project, 'outline') and project.outline:
            outline_title = project.outline.get('title', project.title)

        project.slides_html = ppt_service._combine_slides_to_full_html(
            project.slides_data, outline_title
        )

        # ä½¿ç”¨æ‰¹é‡ä¿å­˜åˆ°æ•°æ®åº“
        from ..services.db_project_manager import DatabaseProjectManager
        db_manager = DatabaseProjectManager()

        # æ‰¹é‡ä¿å­˜å¹»ç¯ç‰‡
        batch_success = await db_manager.batch_save_slides(project_id, slides_data)

        # æ›´æ–°é¡¹ç›®ä¿¡æ¯
        if batch_success:
            await db_manager.update_project_data(project_id, {
                "slides_html": project.slides_html,
                "slides_data": project.slides_data,
                "updated_at": project.updated_at
            })

        logger.debug(f"âœ… é¡¹ç›® {project_id} æ‰¹é‡ä¿å­˜å®Œæˆ,å…± {len(slides_data)} å¼ å¹»ç¯ç‰‡")

        return {
            "success": batch_success,
            "message": f"Successfully batch saved {len(slides_data)} slides" if batch_success else "Batch save failed",
            "slides_count": len(slides_data)
        }

    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"âŒ æ‰¹é‡ä¿å­˜å¹»ç¯ç‰‡å¤±è´¥: {e}")
        return {"success": False, "error": str(e)}


@router.get("/api/projects/{project_id}/export/pdf")
async def export_project_pdf(project_id: str, individual: bool = False):
    """Export project as PDF using Pyppeteer"""
    try:
        project = await ppt_service.project_manager.get_project(project_id)
        if not project:
            raise HTTPException(status_code=404, detail="Project not found")

        # Check if we have slides data
        if not project.slides_data or len(project.slides_data) == 0:
            raise HTTPException(status_code=400, detail="PPT not generated yet")

        # Check if Playwright is available
        pdf_converter = get_pdf_converter()
        logger.info(f"PDF converter (Playwright) available: {pdf_converter.is_available()}")
        if not pdf_converter.is_available():
            raise HTTPException(
                status_code=503,
                detail="PDF generation service unavailable. Please ensure Playwright is installed: pip install playwright && python -m playwright install chromium"
            )

        # Create temp file in thread pool to avoid blocking
        temp_pdf_path = await run_blocking_io(
            lambda: tempfile.NamedTemporaryFile(suffix='.pdf', delete=False).name
        )

        logging.info("Generating PDF with Pyppeteer")
        success = await _generate_pdf_with_pyppeteer(project, temp_pdf_path, individual)

        if not success:
            # Clean up temp file and raise error
            await run_blocking_io(lambda: os.unlink(temp_pdf_path) if os.path.exists(temp_pdf_path) else None)
            raise HTTPException(status_code=500, detail="PDF generation failed")

        # Return PDF file
        logging.info("PDF generated successfully using Pyppeteer")
        safe_filename = urllib.parse.quote(f"{project.topic}_PPT.pdf", safe='')

        # ä½¿ç”¨BackgroundTaskæ¥æ¸…ç†ä¸´æ—¶æ–‡ä»¶
        from starlette.background import BackgroundTask

        def cleanup_temp_file():
            try:
                os.unlink(temp_pdf_path)
            except:
                pass

        return FileResponse(
            temp_pdf_path,
            media_type="application/pdf",
            headers={
                "Content-Disposition": f"attachment; filename*=UTF-8''{safe_filename}",
                "X-PDF-Generator": "Pyppeteer"
            },
            background=BackgroundTask(cleanup_temp_file)
        )

    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@router.get("/api/projects/{project_id}/export/pdf/individual")
async def export_project_pdf_individual(project_id: str):
    """Export project as individual PDF files for each slide"""
    return await export_project_pdf(project_id, individual=True)

@router.get("/api/projects/{project_id}/export/pptx")
async def export_project_pptx(project_id: str, use_apryse: bool = False):
    """Export project as PPTX - supports both Apryse (PDF-based) and Playwright (screenshot-based) methods"""
    try:
        project = await ppt_service.project_manager.get_project(project_id)
        if not project:
            raise HTTPException(status_code=404, detail="Project not found")

        # Check if we have slides data
        if not project.slides_data or len(project.slides_data) == 0:
            raise HTTPException(status_code=400, detail="PPT not generated yet")

        # æ ¹æ®use_apryseå‚æ•°é€‰æ‹©å¯¼å‡ºæ–¹å¼
        if use_apryse:
            # ä½¿ç”¨Apryseæ–¹å¼ï¼šHTML -> PDF -> PPTX
            logger.info(f"Using Apryse method for PPTX export (project: {project_id})")
            
            # Get PDF to PPTX converter
            converter = get_pdf_to_pptx_converter()
            logger.info(f"Apryse converter available: {converter.is_available()}")
            if not converter.is_available():
                raise HTTPException(
                    status_code=503,
                    detail="PPTX conversion service unavailable. Please ensure Apryse SDK is installed and licensed."
                )

            # Check if Playwright is available for PDF generation
            pdf_converter = get_pdf_converter()
            logger.info(f"PDF converter (Playwright) available: {pdf_converter.is_available()}")
            if not pdf_converter.is_available():
                raise HTTPException(
                    status_code=503,
                    detail="PDF generation service unavailable. Please ensure Playwright is installed: pip install playwright && python -m playwright install chromium"
                )
        else:
            # ä½¿ç”¨Playwrightæ–¹å¼ï¼šHTML -> Screenshot -> PPTX
            logger.info(f"Using Playwright screenshot method for PPTX export (project: {project_id})")
            
            # Check if Playwright is available
            pdf_converter = get_pdf_converter()
            logger.info(f"PDF converter (Playwright) available: {pdf_converter.is_available()}")
            if not pdf_converter.is_available():
                raise HTTPException(
                    status_code=503,
                    detail="Screenshot service unavailable. Please ensure Playwright is installed: pip install playwright && python -m playwright install chromium"
                )

        from ..services.background_tasks import get_task_manager
        task_manager = get_task_manager()

        if use_apryse:
            # ===== Apryseæ–¹å¼ï¼šHTML -> PDF -> PPTX =====
            # Step 1: Generate PDF using existing PDF export functionality
            with tempfile.NamedTemporaryFile(suffix='.pdf', delete=False) as temp_pdf_file:
                temp_pdf_path = temp_pdf_file.name

            logging.info("Step 1: Generating PDF for PPTX conversion (Apryse method)")
            pdf_success = await _generate_pdf_with_pyppeteer(project, temp_pdf_path, individual=False)

            if not pdf_success:
                # Clean up temp file and raise error
                try:
                    os.unlink(temp_pdf_path)
                except:
                    pass
                raise HTTPException(status_code=500, detail="PDF generation failed")

            # Step 2: å¯åŠ¨PDFè½¬PPTXåå°ä»»åŠ¡
            logging.info("Step 2: Starting PDF to PPTX conversion task (Apryse)")

            # åˆ›å»ºä¸´æ—¶PPTXæ–‡ä»¶è·¯å¾„
            with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as temp_pptx_file:
                temp_pptx_path = temp_pptx_file.name

            # å®šä¹‰è½¬æ¢ä»»åŠ¡å‡½æ•°
            async def pdf_to_pptx_task():
                """PDF to PPTX conversion task (runs in subprocess)."""
                try:
                    success, result = await converter.convert_pdf_to_pptx_async(
                        temp_pdf_path,
                        temp_pptx_path
                    )
                    if success:
                        # è½¬æ¢æˆåŠŸåï¼Œæ·»åŠ æ¼”è®²ç¨¿åˆ°å¤‡æ³¨
                        try:
                            from pptx import Presentation
                            from ..services.speech_script_repository import SpeechScriptRepository

                            # è·å–æ¼”è®²ç¨¿æ•°æ®
                            repo = SpeechScriptRepository()
                            scripts_list = await repo.get_current_speech_scripts_by_project(project_id)
                            speech_scripts = {script.slide_index: script.script_content for script in scripts_list}
                            repo.close()

                            if len(speech_scripts) > 0:
                                # æ‰“å¼€ç”Ÿæˆçš„PPTXæ–‡ä»¶
                                prs = Presentation(temp_pptx_path)

                                # ä¸ºæ¯å¼ å¹»ç¯ç‰‡æ·»åŠ æ¼”è®²ç¨¿å¤‡æ³¨
                                for i, slide in enumerate(prs.slides):
                                    if i in speech_scripts:
                                        notes_slide = slide.notes_slide
                                        text_frame = notes_slide.notes_text_frame
                                        text_frame.text = speech_scripts[i]
                                        logging.info(f"Added speech script to slide {i+1} notes")

                                # ä¿å­˜ä¿®æ”¹åçš„PPTX
                                prs.save(temp_pptx_path)
                                logging.info(f"Added {len(speech_scripts)} speech scripts to PPTX notes")
                        except Exception as e:
                            logging.warning(f"Failed to add speech scripts to PPTX: {e}")
                            # ç»§ç»­æ‰§è¡Œï¼Œå³ä½¿æ·»åŠ æ¼”è®²ç¨¿å¤±è´¥ä¹Ÿè¿”å›PPTX

                        return {
                            "success": True,
                            "pptx_path": temp_pptx_path,
                            "pdf_path": temp_pdf_path
                        }
                    else:
                        return {
                            "success": False,
                            "error": result
                        }
                except Exception as e:
                    return {
                        "success": False,
                        "error": str(e)
                    }

            # æäº¤åå°ä»»åŠ¡
            task_id = task_manager.submit_task(
                task_type="pdf_to_pptx_conversion",
                func=pdf_to_pptx_task,
                metadata={
                    "project_id": project_id,
                    "project_topic": project.topic,
                    "pdf_path": temp_pdf_path,
                    "pptx_path": temp_pptx_path,
                    "method": "apryse"
                }
            )

            # ç«‹å³è¿”å›ä»»åŠ¡IDï¼Œä¸ç­‰å¾…ä»»åŠ¡å®Œæˆ
            return JSONResponse({
                "status": "processing",
                "task_id": task_id,
                "message": "PPTX conversion started in background (Apryse method)",
                "polling_endpoint": f"/api/ai_slides/tasks/{task_id}"
            })
        
        else:
            # ===== Playwrightæ–¹å¼ï¼šHTML -> Screenshot -> PPTX =====
            logging.info("Starting Playwright screenshot-based PPTX export")
            
            # åˆ›å»ºä¸´æ—¶ç›®å½•å’ŒPPTXæ–‡ä»¶è·¯å¾„
            temp_dir = tempfile.mkdtemp()
            with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as temp_pptx_file:
                temp_pptx_path = temp_pptx_file.name

            # å®šä¹‰HTMLåˆ°å›¾ç‰‡åˆ°PPTXçš„ä»»åŠ¡å‡½æ•°
            async def html_to_pptx_task():
                """ä½¿ç”¨Playwrightæˆªå›¾å¹¶ç”ŸæˆPPTX"""
                screenshot_paths = []
                try:
                    slides_data = project.slides_data
                    logging.info(f"Starting screenshot-based PPTX export for {len(slides_data)} slides")

                    # ç¬¬1æ­¥ï¼šè·å–æ¼”è®²ç¨¿æ•°æ®
                    speech_scripts = {}
                    try:
                        from ..services.speech_script_repository import SpeechScriptRepository
                        repo = SpeechScriptRepository()
                        scripts_list = await repo.get_current_speech_scripts_by_project(project_id)
                        # æ„å»ºå¹»ç¯ç‰‡ç´¢å¼•åˆ°æ¼”è®²ç¨¿çš„æ˜ å°„
                        for script in scripts_list:
                            speech_scripts[script.slide_index] = script.script_content
                        repo.close()
                        logging.info(f"Loaded {len(speech_scripts)} speech scripts for slides")
                    except Exception as e:
                        logging.warning(f"Failed to load speech scripts: {e}")
                        # ç»§ç»­æ‰§è¡Œï¼Œå³ä½¿æ²¡æœ‰æ¼”è®²ç¨¿ä¹Ÿå¯ä»¥ç”ŸæˆPPTX

                    # ç¬¬2æ­¥ï¼šä¸ºæ¯å¼ å¹»ç¯ç‰‡åˆ›å»ºä¸´æ—¶HTMLæ–‡ä»¶
                    html_files = []
                    for i, slide in enumerate(slides_data):
                        html_file = os.path.join(temp_dir, f"slide_{i}.html")
                        html_content = slide.get('html_content', '')
                        with open(html_file, 'w', encoding='utf-8') as f:
                            f.write(html_content)
                        html_files.append(html_file)

                    # ç¬¬3æ­¥ï¼šä½¿ç”¨Playwrightå¯¹æ¯å¼ å¹»ç¯ç‰‡è¿›è¡Œæˆªå›¾
                    for i, html_file in enumerate(html_files):
                        screenshot_path = os.path.join(temp_dir, f"slide_{i}.png")

                        # ä½¿ç”¨PDF converterçš„æˆªå›¾åŠŸèƒ½
                        success = await pdf_converter.screenshot_html(
                            html_file,
                            screenshot_path,
                            width=1280,
                            height=720
                        )

                        if success:
                            screenshot_paths.append(screenshot_path)
                            logging.info(f"Screenshot {i+1}/{len(html_files)} completed")
                        else:
                            logging.warning(f"Screenshot {i+1} failed, skipping")

                    if len(screenshot_paths) == 0:
                        raise Exception("No screenshots were generated")

                    # ç¬¬4æ­¥ï¼šå°†æˆªå›¾è½¬æ¢ä¸ºPPTX
                    from pptx import Presentation
                    from pptx.util import Inches
                    
                    logging.info("Creating PPTX from screenshots...")
                    prs = Presentation()

                    # è®¾ç½®å¹»ç¯ç‰‡å°ºå¯¸ä¸º16:9
                    prs.slide_width = Inches(10)
                    prs.slide_height = Inches(5.625)

                    for i, screenshot_path in enumerate(screenshot_paths):
                        # æ·»åŠ ç©ºç™½å¹»ç¯ç‰‡
                        blank_slide_layout = prs.slide_layouts[6]
                        slide = prs.slides.add_slide(blank_slide_layout)

                        # æ·»åŠ æˆªå›¾ï¼Œå¡«å……æ•´ä¸ªå¹»ç¯ç‰‡
                        left = Inches(0)
                        top = Inches(0)
                        width = prs.slide_width
                        height = prs.slide_height

                        slide.shapes.add_picture(screenshot_path, left, top, width=width, height=height)

                        # å¦‚æœè¯¥å¹»ç¯ç‰‡æœ‰æ¼”è®²ç¨¿ï¼Œæ·»åŠ åˆ°å¤‡æ³¨ä¸­
                        if i in speech_scripts:
                            notes_slide = slide.notes_slide
                            text_frame = notes_slide.notes_text_frame
                            text_frame.text = speech_scripts[i]
                            logging.info(f"Added speech script to slide {i+1} notes")

                    # ä¿å­˜PPTXæ–‡ä»¶
                    prs.save(temp_pptx_path)
                    logging.info(f"PPTX saved to {temp_pptx_path}")

                    return {
                        "success": True,
                        "pptx_path": temp_pptx_path
                    }

                except Exception as e:
                    logging.error(f"HTML to PPTX conversion failed: {e}")
                    import traceback
                    traceback.print_exc()
                    return {
                        "success": False,
                        "error": str(e)
                    }
                finally:
                    # æ¸…ç†ä¸´æ—¶HTMLå’Œæˆªå›¾æ–‡ä»¶
                    try:
                        import shutil
                        if os.path.exists(temp_dir):
                            shutil.rmtree(temp_dir)
                            logging.info(f"Cleaned up temp directory: {temp_dir}")
                    except Exception as cleanup_error:
                        logging.warning(f"Failed to cleanup temp directory: {cleanup_error}")

            # æäº¤åå°ä»»åŠ¡
            task_id = task_manager.submit_task(
                task_type="html_to_pptx_screenshot",
                func=html_to_pptx_task,
                metadata={
                    "project_id": project_id,
                    "project_topic": project.topic,
                    "slide_count": len(project.slides_data),
                    "pptx_path": temp_pptx_path,
                    "method": "playwright"
                }
            )

            # ç«‹å³è¿”å›ä»»åŠ¡ID
            return JSONResponse({
                "status": "processing",
                "task_id": task_id,
                "message": "PPTX generation with screenshots started in background (Playwright method)",
                "polling_endpoint": f"/api/ai_slides/tasks/{task_id}"
            })

    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@router.post("/api/projects/{project_id}/export/pptx-images")
async def export_project_pptx_from_images(project_id: str, request: ImagePPTXExportRequest):
    """Export project as PPTX using high-quality Playwright screenshots"""
    try:
        from io import BytesIO
        from pptx import Presentation
        from pptx.util import Inches

        project = await ppt_service.project_manager.get_project(project_id)
        if not project:
            raise HTTPException(status_code=404, detail="Project not found")

        # éªŒè¯æ˜¯å¦æœ‰å¹»ç¯ç‰‡æ•°æ®
        slides = getattr(request, 'slides', None)
        if not slides or len(slides) == 0:
            raise HTTPException(status_code=400, detail="No slides provided")

        # æ£€æŸ¥Playwrightæ˜¯å¦å¯ç”¨
        pdf_converter = get_pdf_converter()
        if not pdf_converter.is_available():
            raise HTTPException(
                status_code=503,
                detail="Screenshot service unavailable. Please ensure Playwright is installed."
            )

        # åˆ›å»ºåå°ä»»åŠ¡
        from ..services.background_tasks import get_task_manager
        task_manager = get_task_manager()

        # åˆ›å»ºä¸´æ—¶ç›®å½•å’ŒPPTXæ–‡ä»¶è·¯å¾„
        temp_dir = tempfile.mkdtemp()
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as temp_pptx_file:
            temp_pptx_path = temp_pptx_file.name

        # å®šä¹‰HTMLåˆ°å›¾ç‰‡åˆ°PPTXçš„ä»»åŠ¡å‡½æ•°
        async def html_to_pptx_task():
            """ä½¿ç”¨Playwrightæˆªå›¾å¹¶ç”ŸæˆPPTX"""
            screenshot_paths = []
            try:
                logging.info(f"Starting screenshot-based PPTX export for {len(slides)} slides")

                # ç¬¬1æ­¥ï¼šè·å–æ¼”è®²ç¨¿æ•°æ®
                speech_scripts = {}
                try:
                    from ..services.speech_script_repository import SpeechScriptRepository
                    repo = SpeechScriptRepository()
                    scripts_list = await repo.get_current_speech_scripts_by_project(project_id)
                    # æ„å»ºå¹»ç¯ç‰‡ç´¢å¼•åˆ°æ¼”è®²ç¨¿çš„æ˜ å°„
                    for script in scripts_list:
                        speech_scripts[script.slide_index] = script.script_content
                    repo.close()
                    logging.info(f"Loaded {len(speech_scripts)} speech scripts for slides")
                except Exception as e:
                    logging.warning(f"Failed to load speech scripts: {e}")
                    # ç»§ç»­æ‰§è¡Œï¼Œå³ä½¿æ²¡æœ‰æ¼”è®²ç¨¿ä¹Ÿå¯ä»¥ç”ŸæˆPPTX

                # ç¬¬2æ­¥ï¼šä¸ºæ¯å¼ å¹»ç¯ç‰‡åˆ›å»ºä¸´æ—¶HTMLæ–‡ä»¶
                html_files = []
                for i, slide in enumerate(slides):
                    html_file = os.path.join(temp_dir, f"slide_{i}.html")
                    with open(html_file, 'w', encoding='utf-8') as f:
                        f.write(slide['html_content'])
                    html_files.append(html_file)

                # ç¬¬3æ­¥ï¼šä½¿ç”¨Playwrightå¯¹æ¯å¼ å¹»ç¯ç‰‡è¿›è¡Œæˆªå›¾
                for i, html_file in enumerate(html_files):
                    screenshot_path = os.path.join(temp_dir, f"slide_{i}.png")

                    # ä½¿ç”¨PDF converterçš„æˆªå›¾åŠŸèƒ½
                    success = await pdf_converter.screenshot_html(
                        html_file,
                        screenshot_path,
                        width=1280,
                        height=720
                    )

                    if success:
                        screenshot_paths.append(screenshot_path)
                        logging.info(f"Screenshot {i+1}/{len(html_files)} completed")
                    else:
                        logging.warning(f"Screenshot {i+1} failed, skipping")

                if len(screenshot_paths) == 0:
                    raise Exception("No screenshots were generated")

                # ç¬¬4æ­¥ï¼šå°†æˆªå›¾è½¬æ¢ä¸ºPPTX
                logging.info("Creating PPTX from screenshots...")
                prs = Presentation()

                # è®¾ç½®å¹»ç¯ç‰‡å°ºå¯¸ä¸º16:9
                prs.slide_width = Inches(10)
                prs.slide_height = Inches(5.625)

                for i, screenshot_path in enumerate(screenshot_paths):
                    # æ·»åŠ ç©ºç™½å¹»ç¯ç‰‡
                    blank_slide_layout = prs.slide_layouts[6]
                    slide = prs.slides.add_slide(blank_slide_layout)

                    # æ·»åŠ æˆªå›¾ï¼Œå¡«å……æ•´ä¸ªå¹»ç¯ç‰‡
                    left = Inches(0)
                    top = Inches(0)
                    width = prs.slide_width
                    height = prs.slide_height

                    slide.shapes.add_picture(screenshot_path, left, top, width=width, height=height)

                    # å¦‚æœè¯¥å¹»ç¯ç‰‡æœ‰æ¼”è®²ç¨¿ï¼Œæ·»åŠ åˆ°å¤‡æ³¨ä¸­
                    if i in speech_scripts:
                        notes_slide = slide.notes_slide
                        text_frame = notes_slide.notes_text_frame
                        text_frame.text = speech_scripts[i]
                        logging.info(f"Added speech script to slide {i+1} notes")

                # ä¿å­˜PPTXæ–‡ä»¶
                prs.save(temp_pptx_path)
                logging.info(f"PPTX saved to {temp_pptx_path}")

                return {
                    "success": True,
                    "pptx_path": temp_pptx_path
                }

            except Exception as e:
                logging.error(f"HTML to PPTX conversion failed: {e}")
                import traceback
                traceback.print_exc()
                return {
                    "success": False,
                    "error": str(e)
                }
            finally:
                # æ¸…ç†ä¸´æ—¶HTMLå’Œæˆªå›¾æ–‡ä»¶
                try:
                    import shutil
                    if os.path.exists(temp_dir):
                        shutil.rmtree(temp_dir)
                        logging.info(f"Cleaned up temp directory: {temp_dir}")
                except Exception as cleanup_error:
                    logging.warning(f"Failed to cleanup temp directory: {cleanup_error}")

        # æäº¤åå°ä»»åŠ¡
        task_id = task_manager.submit_task(
            task_type="html_to_pptx_screenshot",
            func=html_to_pptx_task,
            metadata={
                "project_id": project_id,
                "project_topic": project.topic,
                "slide_count": len(slides),
                "pptx_path": temp_pptx_path
            }
        )

        # ç«‹å³è¿”å›ä»»åŠ¡ID
        return JSONResponse({
            "status": "processing",
            "task_id": task_id,
            "message": "PPTX generation with screenshots started in background",
            "polling_endpoint": f"/api/ai_slides/tasks/{task_id}"
        })

    except HTTPException:
        raise
    except Exception as e:
        logging.error(f"PPTX screenshot export error: {e}")
        import traceback
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=str(e))


# åå°ä»»åŠ¡æŸ¥è¯¢ç«¯ç‚¹
@router.get("/api/ai_slides/tasks/{task_id}")
async def get_task_status(task_id: str):
    """æŸ¥è¯¢åå°ä»»åŠ¡çŠ¶æ€"""
    from ..services.background_tasks import get_task_manager

    task_manager = get_task_manager()
    task = task_manager.get_task(task_id)

    if not task:
        raise HTTPException(status_code=404, detail="Task not found")

    response = {
        "task_id": task.task_id,
        "task_type": task.task_type,
        "status": task.status.value,
        "progress": task.progress,
        "created_at": task.created_at.isoformat(),
        "updated_at": task.updated_at.isoformat(),
        "metadata": task.metadata
    }

    # å¦‚æœä»»åŠ¡å®Œæˆï¼Œæ·»åŠ ç»“æœä¿¡æ¯
    if task.status.value == "completed" and task.result:
        response["result"] = task.result
        # å¦‚æœæ˜¯PDFè½¬PPTXä»»åŠ¡ï¼Œæä¾›ä¸‹è½½é“¾æ¥
        if task.task_type == "pdf_to_pptx_conversion" and task.result.get("success"):
            response["download_url"] = f"/api/ai_slides/tasks/{task_id}/download"

    # å¦‚æœä»»åŠ¡å¤±è´¥ï¼Œæ·»åŠ é”™è¯¯ä¿¡æ¯
    if task.status.value == "failed":
        response["error"] = task.error

    return JSONResponse(response)


@router.get("/api/ai_slides/tasks/{task_id}/download")
async def download_task_result(task_id: str):
    """ä¸‹è½½ä»»åŠ¡ç»“æœæ–‡ä»¶"""
    from ..services.background_tasks import get_task_manager, TaskStatus
    from starlette.background import BackgroundTask

    task_manager = get_task_manager()
    task = task_manager.get_task(task_id)

    if not task:
        raise HTTPException(status_code=404, detail="Task not found")

    if task.status != TaskStatus.COMPLETED:
        raise HTTPException(status_code=400, detail=f"Task not completed yet (status: {task.status.value})")

    if not task.result or not task.result.get("success"):
        raise HTTPException(status_code=400, detail="Task failed or no result available")

    pptx_path = task.result.get("pptx_path")
    pdf_path = task.result.get("pdf_path")

    if not pptx_path or not os.path.exists(pptx_path):
        raise HTTPException(status_code=404, detail="Result file not found")

    # è·å–é¡¹ç›®ä¸»é¢˜ä½œä¸ºæ–‡ä»¶å
    project_topic = task.metadata.get("project_topic", "PPT")
    safe_filename = urllib.parse.quote(f"{project_topic}_PPT.pptx", safe='')

    # æ¸…ç†ä¸´æ—¶æ–‡ä»¶çš„åå°ä»»åŠ¡
    def cleanup_temp_files():
        try:
            if pdf_path and os.path.exists(pdf_path):
                os.unlink(pdf_path)
        except:
            pass
        try:
            if pptx_path and os.path.exists(pptx_path):
                os.unlink(pptx_path)
        except:
            pass

    return FileResponse(
        pptx_path,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={
            "Content-Disposition": f"attachment; filename*=UTF-8''{safe_filename}",
            "X-Conversion-Method": "PDF-to-PPTX-Background"
        },
        background=BackgroundTask(cleanup_temp_files)
    )


@router.get("/api/projects/{project_id}/export/html")
async def export_project_html(project_id: str):
    """Export project as HTML ZIP package with slideshow index"""
    try:
        project = await ppt_service.project_manager.get_project(project_id)
        if not project:
            raise HTTPException(status_code=404, detail="Project not found")

        # Check if we have slides data
        if not project.slides_data or len(project.slides_data) == 0:
            raise HTTPException(status_code=400, detail="PPT not generated yet")

        # Create temporary directory and generate files in thread pool
        zip_content = await run_blocking_io(_generate_html_export_sync, project)

        # URL encode the filename to handle Chinese characters
        zip_filename = f"{project.topic}_PPT.zip"
        safe_filename = urllib.parse.quote(zip_filename, safe='')

        from fastapi.responses import Response
        return Response(
            content=zip_content,
            media_type="application/zip",
            headers={
                "Content-Disposition": f"attachment; filename*=UTF-8''{safe_filename}"
            }
        )

    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


def _generate_html_export_sync(project) -> bytes:
    """åŒæ­¥ç”ŸæˆHTMLå¯¼å‡ºæ–‡ä»¶ï¼ˆåœ¨çº¿ç¨‹æ± ä¸­è¿è¡Œï¼‰"""
    with tempfile.TemporaryDirectory() as temp_dir:
        temp_path = Path(temp_dir)

        # Generate individual HTML files for each slide
        slide_files = []
        for i, slide in enumerate(project.slides_data):
            slide_filename = f"slide_{i+1}.html"
            slide_files.append(slide_filename)

            # Create complete HTML document for each slide
            slide_html = _generate_individual_slide_html_sync(slide, i+1, len(project.slides_data), project.topic)

            slide_path = temp_path / slide_filename
            with open(slide_path, 'w', encoding='utf-8') as f:
                f.write(slide_html)

        # Generate index.html slideshow page
        index_html = _generate_slideshow_index_sync(project, slide_files)
        index_path = temp_path / "index.html"
        with open(index_path, 'w', encoding='utf-8') as f:
            f.write(index_html)

        # Create ZIP file
        zip_filename = f"{project.topic}_PPT.zip"
        zip_path = temp_path / zip_filename

        with zipfile.ZipFile(zip_path, 'w', zipfile.ZIP_DEFLATED) as zipf:
            # Add index.html
            zipf.write(index_path, "index.html")

            # Add all slide files
            for slide_file in slide_files:
                slide_path = temp_path / slide_file
                zipf.write(slide_path, slide_file)

        # Read ZIP file content
        with open(zip_path, 'rb') as f:
            return f.read()


def _generate_individual_slide_html_sync(slide, slide_number: int, total_slides: int, topic: str) -> str:
    """åŒæ­¥ç”Ÿæˆå•ä¸ªå¹»ç¯ç‰‡HTMLï¼ˆåœ¨çº¿ç¨‹æ± ä¸­è¿è¡Œï¼‰"""
    slide_html = slide.get('html_content', '')
    slide_title = slide.get('title', f'ç¬¬{slide_number}é¡µ')

    # Check if it's already a complete HTML document
    import re
    if slide_html.strip().lower().startswith('<!doctype') or slide_html.strip().lower().startswith('<html'):
        # It's a complete HTML document, enhance it with navigation
        return _enhance_complete_html_with_navigation(slide_html, slide_number, total_slides, topic, slide_title)
    else:
        # It's just content, wrap it in a complete structure
        slide_content = slide_html

    return f"""<!DOCTYPE html>
<html lang="zh-CN">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>{topic} - {slide_title}</title>
    <style>
        body {{
            margin: 0;
            padding: 0;
            font-family: 'Microsoft YaHei', 'PingFang SC', sans-serif;
            background: #f5f5f5;
            display: flex;
            align-items: center;
            justify-content: center;
            min-height: 100vh;
        }}
        .slide-container {{
            width: 90vw;
            height: 90vh;
            background: white;
            border-radius: 10px;
            box-shadow: 0 4px 20px rgba(0,0,0,0.1);
            overflow: hidden;
            position: relative;
        }}
        .slide-content {{
            width: 100%;
            height: 100%;
            padding: 20px;
            box-sizing: border-box;
        }}
        .slide-number {{
            position: absolute;
            bottom: 20px;
            right: 20px;
            background: rgba(0,0,0,0.7);
            color: white;
            padding: 5px 10px;
            border-radius: 5px;
            font-size: 14px;
        }}
    </style>
</head>
<body>
    <div class="slide-container">
        <div class="slide-content">
            {slide_content}
        </div>
        <div class="slide-number">{slide_number} / {total_slides}</div>
    </div>
</body>
</html>"""


def _enhance_complete_html_with_navigation(html_content: str, slide_number: int, total_slides: int, topic: str, slide_title: str) -> str:
    """ä¸ºå®Œæ•´çš„ HTML æ–‡æ¡£æ·»åŠ å¯¼èˆªä¿¡æ¯"""
    # ç®€å•åœ°åœ¨ body ç»“æŸæ ‡ç­¾å‰æ·»åŠ å¹»ç¯ç‰‡ç¼–å·
    import re
    
    # å°è¯•åœ¨ </body> å‰æ’å…¥å¹»ç¯ç‰‡ç¼–å·
    if '</body>' in html_content.lower():
        slide_number_html = f'''
    <div style="position: fixed; bottom: 20px; right: 20px; background: rgba(0,0,0,0.7); color: white; padding: 5px 10px; border-radius: 5px; font-size: 14px; z-index: 9999;">
        {slide_number} / {total_slides}
    </div>
'''
        html_content = re.sub(r'</body>', slide_number_html + '</body>', html_content, flags=re.IGNORECASE)
    
    return html_content


def _generate_slideshow_index_sync(project, slide_files: list) -> str:
    """åŒæ­¥ç”Ÿæˆå¹»ç¯ç‰‡ç´¢å¼•é¡µé¢ï¼ˆåœ¨çº¿ç¨‹æ± ä¸­è¿è¡Œï¼‰"""
    slides_list = ""
    for i, slide_file in enumerate(slide_files):
        slide = project.slides_data[i]
        slide_title = slide.get('title', f'ç¬¬{i+1}é¡µ')
        slides_list += f"""
        <div class="slide-item" onclick="openSlide('{slide_file}')">
            <div class="slide-preview">
                <div class="slide-number">{i+1}</div>
                <div class="slide-title">{slide_title}</div>
            </div>
        </div>"""

    return f"""<!DOCTYPE html>
<html lang="zh-CN">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>{project.topic} - PPTæ”¾æ˜ </title>
    <style>
        body {{
            margin: 0;
            padding: 0;
            font-family: 'Microsoft YaHei', 'PingFang SC', sans-serif;
            background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
            min-height: 100vh;
        }}
        .header {{
            text-align: center;
            padding: 40px 20px;
            color: white;
        }}
        .header h1 {{
            margin: 0;
            font-size: 2.5em;
            font-weight: 300;
        }}
        .slides-grid {{
            display: grid;
            grid-template-columns: repeat(auto-fill, minmax(300px, 1fr));
            gap: 20px;
            padding: 20px;
        }}
        .slide-item {{
            background: white;
            border-radius: 10px;
            padding: 20px;
            cursor: pointer;
            transition: all 0.3s ease;
            box-shadow: 0 4px 15px rgba(0,0,0,0.1);
        }}
        .slide-item:hover {{
            transform: translateY(-5px);
            box-shadow: 0 8px 25px rgba(0,0,0,0.2);
        }}
        .slide-number {{
            background: #007bff;
            color: white;
            width: 40px;
            height: 40px;
            border-radius: 50%;
            display: flex;
            align-items: center;
            justify-content: center;
            margin: 0 auto 15px auto;
            font-weight: bold;
        }}
        .slide-title {{
            font-size: 1.1em;
            color: #333;
            margin: 0;
            text-align: center;
        }}
    </style>
</head>
<body>
    <div class="header">
        <h1>{project.topic}</h1>
        <p>PPTæ¼”ç¤ºæ–‡ç¨¿ - å…±{len(slide_files)}é¡µ</p>
    </div>
    <div class="slides-grid">
        {slides_list}
    </div>
    <script>
        function openSlide(slideFile) {{
            window.open(slideFile, '_blank');
        }}
    </script>
</body>
</html>"""



@router.get("/upload", response_class=HTMLResponse)
async def web_upload_page(
    request: Request
):
    """File upload page"""
    return templates.TemplateResponse("upload.html", {
        "request": request
    })
async def _process_uploaded_files_for_outline(
    file_uploads: List[UploadFile],
    topic: str,
    target_audience: str,
    page_count_mode: str,
    min_pages: int,
    max_pages: int,
    fixed_pages: int,
    ppt_style: str,
    custom_style_prompt: str,
    file_processing_mode: str,
    content_analysis_depth: str,
    requirements: str = None,
    enable_web_search: bool = False,
    scenario: str = "general",
    language: str = "zh"
) -> Optional[Dict[str, Any]]:
    """å¤„ç†ä¸Šä¼ çš„å¤šä¸ªæ–‡ä»¶å¹¶ç”ŸæˆPPTå¤§çº²,æ”¯æŒè”ç½‘æœç´¢é›†æˆ"""
    try:
        from ..services.file_processor import FileProcessor
        file_processor = FileProcessor()

        # è¿‡æ»¤æ‰Noneå€¼(å¦‚æœæ²¡æœ‰æ–‡ä»¶ä¸Šä¼ )
        files = [f for f in file_uploads if f is not None]
        if not files:
            logger.error("No files provided")
            return None

        saved_file_paths = []
        all_processed_content = []

        # Process each uploaded file

        try:
            # å¤„ç†æ¯ä¸ªæ–‡ä»¶
            for file_upload in files:
                # éªŒè¯æ–‡ä»¶
                is_valid, message = file_processor.validate_file(file_upload.filename, file_upload.size)
                if not is_valid:
                    logger.error(f"File validation failed for {file_upload.filename}: {message}")
                    continue

                # è¯»å–æ–‡ä»¶å†…å®¹å¹¶ä¿å­˜åˆ°é¡¹ç›®æ–‡ä»¶ç›®å½•
                content = await file_upload.read()
                # logger.info(f"æ–‡ä»¶å†…å®¹: {content}")
                project_file_path = await run_blocking_io(
                    _save_project_file_sync, content, file_upload.filename
                )
                saved_file_paths.append(project_file_path)

                # å¤„ç†å•ä¸ªæ–‡ä»¶å†…å®¹
                file_result = await file_processor.process_file(project_file_path, file_upload.filename)
                all_processed_content.append({
                    "filename": file_upload.filename,
                    "content": file_result.processed_content
                })
                logger.debug(f"æ–‡ä»¶å¤„ç†å†…å®¹: {file_result.processed_content}")
            if not all_processed_content:
                logger.error("No files were successfully processed")
                return None

            # åˆå¹¶æ‰€æœ‰æ–‡ä»¶å†…å®¹
            merged_content = file_processor.merge_multiple_files_to_markdown(all_processed_content)

            # åˆ›å»ºä¸´æ—¶åˆå¹¶æ–‡ä»¶
            import tempfile
            import os
            with tempfile.NamedTemporaryFile(mode='w', delete=False, suffix='.md', encoding='utf-8') as merged_file:
                merged_file.write(merged_content)
                merged_file_path = merged_file.name

            merged_filename = f"merged_content_{len(files)}_files.md"

            saved_file_paths.append(merged_file_path)

            # åˆ›å»ºæ–‡ä»¶å¤§çº²ç”Ÿæˆè¯·æ±‚
            from ..api.models import FileOutlineGenerationRequest
            filenames_str = ", ".join([f.filename for f in files])
            merged_filename = f"merged_content_{len(files)}_files.md"
            outline_request = FileOutlineGenerationRequest(
                file_path=merged_file_path,
                filename=merged_filename,
                topic=topic if topic.strip() else None,
                scenario="general",
                requirements=requirements,
                target_audience=target_audience,
                page_count_mode=page_count_mode,
                min_pages=min_pages,
                max_pages=max_pages,
                fixed_pages=fixed_pages,
                ppt_style=ppt_style,
                custom_style_prompt=custom_style_prompt,
                file_processing_mode=file_processing_mode,
                content_analysis_depth=content_analysis_depth
            )

            # ä½¿ç”¨enhanced_ppt_serviceç”Ÿæˆå¤§çº²
            result = await ppt_service.generate_outline_from_file(outline_request)

            if result.success:
                logger.info(f"Successfully generated outline from {len(files)} files: {filenames_str}")
                # åœ¨å¤§çº²ä¸­æ·»åŠ æ–‡ä»¶ä¿¡æ¯,ç”¨äºé‡æ–°ç”Ÿæˆ
                outline_with_file_info = result.outline.copy()
                original_filenames = [f.filename for f in files]
                file_paths_without_merge = saved_file_paths[:-1]  # æ’é™¤ä¸´æ—¶åˆå¹¶æ–‡ä»¶
                uploaded_files_info = [
                    {'filename': name, 'file_path': path}
                    for name, path in zip(original_filenames, file_paths_without_merge)
                ]
                outline_with_file_info['file_info'] = {
                    'file_paths': file_paths_without_merge,
                    'merged_file_path': merged_file_path,
                    'merged_filename': merged_filename,
                    'filenames': original_filenames,
                    'files_count': len(files),
                    'processing_mode': file_processing_mode,
                    'analysis_depth': content_analysis_depth,
                    'file_path': merged_file_path,
                    'filename': merged_filename,
                    'uploaded_files': uploaded_files_info
                }
                return outline_with_file_info
            else:
                logger.error(f"Failed to generate outline from files: {result.error}")
                # å¦‚æœç”Ÿæˆå¤±è´¥,æ¸…ç†æ–‡ä»¶
                for file_path in saved_file_paths:
                    await run_blocking_io(_cleanup_project_file_sync, file_path)
                return None

        except Exception as e:
            # æ¸…ç†æ‰€æœ‰å·²ä¿å­˜çš„æ–‡ä»¶
            for file_path in saved_file_paths:
                try:
                    await run_blocking_io(_cleanup_project_file_sync, file_path)
                except:
                    pass
            raise e

    except Exception as e:
        logger.error(f"Error processing uploaded files for outline: {e}")
        return None


async def _process_uploaded_file_for_outline(
    file_upload: UploadFile,
    topic: str,
    target_audience: str,
    page_count_mode: str,
    min_pages: int,
    max_pages: int,
    fixed_pages: int,
    ppt_style: str,
    custom_style_prompt: str,
    file_processing_mode: str,
    content_analysis_depth: str,
    requirements: str = None
) -> Optional[Dict[str, Any]]:
    """å¤„ç†ä¸Šä¼ çš„å•ä¸ªæ–‡ä»¶å¹¶ç”ŸæˆPPTå¤§çº²(å‘åå…¼å®¹)"""
    return await _process_uploaded_files_for_outline(
        [file_upload], topic, target_audience, page_count_mode, min_pages, max_pages,
        fixed_pages, ppt_style, custom_style_prompt, file_processing_mode,
        content_analysis_depth, requirements
    )


def _save_temp_file_sync(content: bytes, filename: str) -> str:
    """åŒæ­¥ä¿å­˜ä¸´æ—¶æ–‡ä»¶(åœ¨çº¿ç¨‹æ± ä¸­è¿è¡Œ)"""
    import tempfile
    import os

    with tempfile.NamedTemporaryFile(
        delete=False,
        suffix=os.path.splitext(filename)[1]
    ) as temp_file:
        temp_file.write(content)
        return temp_file.name


def _save_project_file_sync(content: bytes, filename: str) -> str:
    """åŒæ­¥ä¿å­˜é¡¹ç›®æ–‡ä»¶åˆ°æ°¸ä¹…ä½ç½®(åœ¨çº¿ç¨‹æ± ä¸­è¿è¡Œ)"""
    import os
    import time
    from pathlib import Path

    # åˆ›å»ºé¡¹ç›®æ–‡ä»¶ç›®å½•
    project_files_dir = Path("temp/project_files")
    project_files_dir.mkdir(parents=True, exist_ok=True)

    # ç”Ÿæˆå”¯ä¸€æ–‡ä»¶å
    timestamp = int(time.time())
    file_ext = os.path.splitext(filename)[1]
    safe_filename = f"{timestamp}_{filename}"
    file_path = project_files_dir / safe_filename

    # ä¿å­˜æ–‡ä»¶
    with open(file_path, 'wb') as f:
        f.write(content)

    return str(file_path)


def _cleanup_temp_file_sync(temp_file_path: str):
    """åŒæ­¥æ¸…ç†ä¸´æ—¶æ–‡ä»¶(åœ¨çº¿ç¨‹æ± ä¸­è¿è¡Œ)"""
    import os
    if os.path.exists(temp_file_path):
        os.unlink(temp_file_path)


def _cleanup_project_file_sync(project_file_path: str):
    """åŒæ­¥æ¸…ç†é¡¹ç›®æ–‡ä»¶(åœ¨çº¿ç¨‹æ± ä¸­è¿è¡Œ)"""
    import os
    if os.path.exists(project_file_path):
        os.unlink(project_file_path)


@router.get("/global-master-templates", response_class=HTMLResponse)
async def global_master_templates_page(
    request: Request
):
    """Global master templates management page"""
    try:
        return templates.TemplateResponse("global_master_templates.html", {
            "request": request
        })
    except Exception as e:
        logger.error(f"Error loading global master templates page: {e}")
        return templates.TemplateResponse("error.html", {
            "request": request,
            "error": str(e)
        })


@router.get("/image-gallery", response_class=HTMLResponse)
async def image_gallery_page(
    request: Request
):
    """æœ¬åœ°å›¾åºŠç®¡ç†é¡µé¢"""
    try:
        return templates.TemplateResponse("image_gallery.html", {
            "request": request,
            "user": user
        })
    except Exception as e:
        logger.error(f"Error rendering image gallery page: {e}")
        return templates.TemplateResponse("error.html", {
            "request": request,
            "error": str(e)
        })


@router.get("/image-generation-test", response_class=HTMLResponse)
async def image_generation_test_page(
    request: Request
):
    """AIå›¾ç‰‡ç”Ÿæˆæµ‹è¯•é¡µé¢"""
    try:
        return templates.TemplateResponse("image_generation_test.html", {
            "request": request,
            "user": user
        })
    except Exception as e:
        logger.error(f"Error rendering image generation test page: {e}")
        return templates.TemplateResponse("error.html", {
            "request": request,
            "error": str(e)
        })


@router.get("/projects/{project_id}/template-selection", response_class=HTMLResponse)
async def template_selection_page(
    request: Request,
    project_id: str
):
    """Template selection page for PPT generation"""
    try:
        # Get project info
        project = await ppt_service.project_manager.get_project(project_id)
        if not project:
            raise HTTPException(status_code=404, detail="Project not found")

        return templates.TemplateResponse("template_selection.html", {
            "request": request,
            "project_id": project_id,
            "project_topic": project.topic
        })
    except Exception as e:
        logger.error(f"Error loading template selection page: {e}")
        return templates.TemplateResponse("error.html", {
            "request": request,
            "error": str(e)
        })


# å›¾åƒé‡æ–°ç”Ÿæˆç›¸å…³è¾…åŠ©å‡½æ•°
async def analyze_image_context(image_info: Dict[str, Any], slide_content: Dict[str, Any],
                               project_topic: str, project_scenario: str) -> Dict[str, Any]:
    """åˆ†æå›¾åƒåœ¨å¹»ç¯ç‰‡ä¸­çš„ä¸Šä¸‹æ–‡"""
    return {
        "slide_title": slide_content.get("title", ""),
        "slide_content": slide_content.get("html_content", ""),
        "image_alt": image_info.get("alt", ""),
        "image_title": image_info.get("title", ""),
        "image_size": f"{image_info.get('width', 0)}x{image_info.get('height', 0)}",
        "image_position": image_info.get("position", {}),
        "project_topic": project_topic,
        "project_scenario": project_scenario,
        "image_purpose": determine_image_purpose(image_info, slide_content)
    }

def determine_image_purpose(image_info: Dict[str, Any], slide_content: Dict[str, Any]) -> str:
    """ç¡®å®šå›¾åƒåœ¨å¹»ç¯ç‰‡ä¸­çš„ç”¨é€”"""
    # ç®€å•çš„å¯å‘å¼è§„åˆ™æ¥ç¡®å®šå›¾åƒç”¨é€”
    width = image_info.get('width', 0)
    height = image_info.get('height', 0)
    alt_text = image_info.get('alt', '').lower()

    if width > 800 or height > 600:
        return "background"  # å¤§å›¾åƒå¯èƒ½æ˜¯èƒŒæ™¯
    elif 'icon' in alt_text or 'logo' in alt_text:
        return "icon"
    elif 'chart' in alt_text or 'graph' in alt_text:
        return "chart_support"
    elif width < 200 and height < 200:
        return "decoration"
    else:
        return "illustration"

# å›¾åƒé‡æ–°ç”Ÿæˆç›¸å…³è¾…åŠ©å‡½æ•°

def select_best_image_source(enabled_sources: List, image_config: Dict[str, Any], image_context: Dict[str, Any]):
    """æ™ºèƒ½é€‰æ‹©æœ€ä½³çš„å›¾ç‰‡æ¥æº"""
    from ..services.models.slide_image_info import ImageSource

    # å¦‚æœåªæœ‰ä¸€ä¸ªå¯ç”¨çš„æ¥æº,ç›´æ¥ä½¿ç”¨
    if len(enabled_sources) == 1:
        return enabled_sources[0]

    # æ ¹æ®å›¾åƒç”¨é€”å’Œé…ç½®æ™ºèƒ½é€‰æ‹©
    image_purpose = image_context.get('image_purpose', 'illustration')

    # ä¼˜å…ˆçº§è§„åˆ™
    if image_purpose == 'background':
        # èƒŒæ™¯å›¾ä¼˜å…ˆä½¿ç”¨AIç”Ÿæˆ,å…¶æ¬¡ç½‘ç»œæœç´¢
        if ImageSource.AI_GENERATED in enabled_sources:
            return ImageSource.AI_GENERATED
        elif ImageSource.NETWORK in enabled_sources:
            return ImageSource.NETWORK
        elif ImageSource.LOCAL in enabled_sources:
            return ImageSource.LOCAL

    elif image_purpose == 'icon':
        # å›¾æ ‡ä¼˜å…ˆä½¿ç”¨æœ¬åœ°,å…¶æ¬¡AIç”Ÿæˆ
        if ImageSource.LOCAL in enabled_sources:
            return ImageSource.LOCAL
        elif ImageSource.AI_GENERATED in enabled_sources:
            return ImageSource.AI_GENERATED
        elif ImageSource.NETWORK in enabled_sources:
            return ImageSource.NETWORK

    elif image_purpose in ['illustration', 'chart_support', 'decoration']:
        # è¯´æ˜æ€§å›¾ç‰‡ä¼˜å…ˆä½¿ç”¨ç½‘ç»œæœç´¢,å…¶æ¬¡AIç”Ÿæˆ
        if ImageSource.NETWORK in enabled_sources:
            return ImageSource.NETWORK
        elif ImageSource.AI_GENERATED in enabled_sources:
            return ImageSource.AI_GENERATED
        elif ImageSource.LOCAL in enabled_sources:
            return ImageSource.LOCAL

    # é»˜è®¤ä¼˜å…ˆçº§:AIç”Ÿæˆ > ç½‘ç»œæœç´¢ > æœ¬åœ°
    for source in [ImageSource.AI_GENERATED, ImageSource.NETWORK, ImageSource.LOCAL]:
        if source in enabled_sources:
            return source

    # å¦‚æœéƒ½æ²¡æœ‰,è¿”å›ç¬¬ä¸€ä¸ªå¯ç”¨çš„
    return enabled_sources[0] if enabled_sources else ImageSource.AI_GENERATED

# æ³¨æ„:generate_image_prompt_for_replacement å‡½æ•°å·²è¢«PPTImageProcessorçš„æ ‡å‡†æµç¨‹æ›¿ä»£
# ç°åœ¨ä½¿ç”¨ PPTImageProcessor._ai_generate_image_prompt æ–¹æ³•æ¥ç”Ÿæˆæç¤ºè¯

def replace_image_in_html(html_content: str, image_info: Dict[str, Any], new_image_url: str) -> str:
    """åœ¨HTMLå†…å®¹ä¸­æ›¿æ¢æŒ‡å®šçš„å›¾åƒ,æ”¯æŒimgæ ‡ç­¾,èƒŒæ™¯å›¾åƒå’ŒSVG,ä¿æŒå¸ƒå±€å’Œæ ·å¼"""
    try:
        from bs4 import BeautifulSoup
        import re

        soup = BeautifulSoup(html_content, 'html.parser')

        old_src = image_info.get('src', '')
        image_type = image_info.get('type', 'img')

        if not old_src:
            logger.warning("å›¾åƒä¿¡æ¯ä¸­æ²¡æœ‰srcå±æ€§,æ— æ³•æ›¿æ¢")
            return html_content

        replacement_success = False

        if image_type == 'img':
            # å¤„ç† <img> æ ‡ç­¾
            replacement_success = replace_img_tag(soup, image_info, new_image_url, old_src)

        elif image_type == 'background':
            # å¤„ç†èƒŒæ™¯å›¾åƒ
            replacement_success = replace_background_image(soup, image_info, new_image_url, old_src)

        elif image_type == 'svg':
            # å¤„ç†SVGå›¾åƒ
            replacement_success = replace_svg_image(soup, image_info, new_image_url, old_src)

        if replacement_success:
            logger.info(f"æˆåŠŸæ›¿æ¢{image_type}å›¾åƒ: {old_src} -> {new_image_url}")
            return str(soup)
        else:
            logger.warning(f"æœªæ‰¾åˆ°åŒ¹é…çš„{image_type}å›¾åƒè¿›è¡Œæ›¿æ¢")
            return fallback_string_replacement(html_content, old_src, new_image_url)

    except Exception as e:
        logger.error(f"æ›¿æ¢HTMLä¸­çš„å›¾åƒå¤±è´¥: {e}")
        return fallback_string_replacement(html_content, image_info.get('src', ''), new_image_url)

def replace_img_tag(soup, image_info: Dict[str, Any], new_image_url: str, old_src: str) -> bool:
    """æ›¿æ¢imgæ ‡ç­¾"""
    img_elements = soup.find_all('img')

    for img in img_elements:
        img_src = img.get('src', '')

        # æ¯”è¾ƒå›¾åƒæºURL(å¤„ç†ç›¸å¯¹è·¯å¾„å’Œç»å¯¹è·¯å¾„)
        if (img_src == old_src or
            img_src.endswith(old_src.split('/')[-1]) or
            old_src.endswith(img_src.split('/')[-1])):

            # æ›¿æ¢å›¾åƒURL
            img['src'] = new_image_url

            # ä¿æŒåŸæœ‰çš„é‡è¦å±æ€§
            preserved_attributes = ['class', 'style', 'width', 'height', 'id']
            for attr in preserved_attributes:
                if attr in image_info and image_info[attr]:
                    img[attr] = image_info[attr]

            # æ›´æ–°æˆ–ä¿æŒaltå’Œtitle
            if image_info.get('alt'):
                img['alt'] = image_info['alt']
            if image_info.get('title'):
                img['title'] = image_info['title']

            # ç¡®ä¿å›¾åƒåŠ è½½é”™è¯¯æ—¶æœ‰åå¤‡å¤„ç†
            if not img.get('onerror'):
                img['onerror'] = "this.style.display='none'"

            return True

    return False

def replace_background_image(soup, image_info: Dict[str, Any], new_image_url: str, old_src: str) -> bool:
    """æ›¿æ¢CSSèƒŒæ™¯å›¾åƒ"""
    # æŸ¥æ‰¾æ‰€æœ‰å…ƒç´ 
    all_elements = soup.find_all()

    for element in all_elements:
        # æ£€æŸ¥å†…è”æ ·å¼ä¸­çš„èƒŒæ™¯å›¾åƒ
        style = element.get('style', '')
        if 'background-image' in style and old_src in style:
            # æ›¿æ¢å†…è”æ ·å¼ä¸­çš„èƒŒæ™¯å›¾åƒURL
            new_style = style.replace(old_src, new_image_url)
            element['style'] = new_style
            return True

        # æ£€æŸ¥classå±æ€§,å¯èƒ½å¯¹åº”CSSè§„åˆ™ä¸­çš„èƒŒæ™¯å›¾åƒ
        class_names = element.get('class', [])
        if class_names and image_info.get('className'):
            # å¦‚æœclassåŒ¹é…,æˆ‘ä»¬å‡è®¾è¿™æ˜¯ç›®æ ‡å…ƒç´ 
            if any(cls in image_info.get('className', '') for cls in class_names):
                # ä¸ºå…ƒç´ æ·»åŠ å†…è”èƒŒæ™¯å›¾åƒæ ·å¼
                current_style = element.get('style', '')
                if current_style and not current_style.endswith(';'):
                    current_style += ';'
                new_style = f"{current_style}background-image: url('{new_image_url}');"
                element['style'] = new_style
                return True

    return False

def replace_svg_image(soup, image_info: Dict[str, Any], new_image_url: str, old_src: str) -> bool:
    """æ›¿æ¢SVGå›¾åƒ"""
    # æŸ¥æ‰¾SVGå…ƒç´ 
    svg_elements = soup.find_all('svg')

    for svg in svg_elements:
        # å¦‚æœSVGæœ‰srcå±æ€§(è™½ç„¶ä¸å¸¸è§)
        if svg.get('src') == old_src:
            svg['src'] = new_image_url
            return True

        # æ£€æŸ¥SVGçš„å†…å®¹æˆ–å…¶ä»–æ ‡è¯†
        if image_info.get('outerHTML') and svg.get_text() in image_info.get('outerHTML', ''):
            # å¯¹äºå†…è”SVG,æˆ‘ä»¬å¯èƒ½éœ€è¦æ›¿æ¢æ•´ä¸ªå…ƒç´ 
            # è¿™é‡Œç®€åŒ–å¤„ç†,æ·»åŠ ä¸€ä¸ªdataå±æ€§æ¥æ ‡è®°å·²æ›¿æ¢
            svg['data-replaced-image'] = new_image_url
            return True

    return False

def fallback_string_replacement(html_content: str, old_src: str, new_image_url: str) -> str:
    """åå¤‡çš„å­—ç¬¦ä¸²æ›¿æ¢æ–¹æ¡ˆ"""
    try:
        import re

        if old_src and old_src in html_content:
            # å°è¯•å¤šç§æ›¿æ¢æ¨¡å¼
            patterns = [
                # imgæ ‡ç­¾çš„srcå±æ€§
                (rf'(<img[^>]*src=")[^"]*({re.escape(old_src)}[^"]*")([^>]*>)', rf'\1{new_image_url}\3'),
                # CSSèƒŒæ™¯å›¾åƒ
                (rf'(background-image:\s*url\([\'"]?)[^\'")]*({re.escape(old_src)}[^\'")]*)', rf'\1{new_image_url}'),
                # ç›´æ¥å­—ç¬¦ä¸²æ›¿æ¢
                (re.escape(old_src), new_image_url)
            ]

            for pattern, replacement in patterns:
                updated_html = re.sub(pattern, replacement, html_content, flags=re.IGNORECASE)
                if updated_html != html_content:
                    logger.info(f"ä½¿ç”¨åå¤‡æ–¹æ¡ˆæˆåŠŸæ›¿æ¢å›¾åƒ: {old_src} -> {new_image_url}")
                    return updated_html

        return html_content

    except Exception as e:
        logger.error(f"åå¤‡æ›¿æ¢æ–¹æ¡ˆä¹Ÿå¤±è´¥: {e}")
        return html_content
